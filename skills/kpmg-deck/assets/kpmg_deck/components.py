"""
components.py -- the design vocabulary. A closed set, deliberately.

THIS IS WHERE TASTE IS ENCODED AS CONSTRAINT.

The point of this package is not that it can draw anything. PowerPoint can already draw
anything, and that is precisely the problem: a tool with no constraints hands every decision to
whoever is driving, and most of those decisions are ones a non-designer should not have to make
correctly at 2am. So this module offers a SMALL, OPINIONATED set of moves, each of which is
already correct, rather than a general drawing API plus advice.

Three mechanisms do the work:

1. GROUND-AWARE COLOUR. A component is told what it is sitting on (`Ground.LIGHT`,
   `Ground.DARK`, `Ground.BRAND`) and picks its own foreground from the legal set for that
   ground. A caller cannot put Cobalt on Dark blue -- a pairing that measures 2.35:1 and is
   invisible in a lit room -- because a caller never names a foreground colour at all.

2. NO CALLER-SUPPLIED GEOMETRY. Components take a `Box` from the grid, never raw coordinates.
   Drift is structurally impossible rather than merely discouraged.

3. NO BULLET CHARACTERS ANYWHERE. There is no `bullets()` function in this module and that is
   not an oversight. A stack of bullets is a page of notes; the equivalent component here
   (`rows`) separates items with hairline rules and space, which is what a designed deck does
   and what forces the writer to keep each item short.

The negative space matters as much as the list: no clip art, no decorative icons, no drop
shadows outside `card()`, no gradients except the one signature treatment, no rounded corners
on structural elements, no centre-aligned body text, no 3D anything.
"""

from __future__ import annotations

from dataclasses import dataclass
from enum import Enum
from typing import Sequence

from pptx.enum.shapes import MSO_SHAPE

from . import oxml as X
from . import text as T
from .canvas import GRID, Box, Grid, inches, points
from .tokens import (
    LINE_SPACING,
    SCALE,
    SPACE,
    Brand,
    contrast_ratio,
    contrast_verdict,
    tracking_pt,
)


@dataclass(frozen=True)
class Placed:
    """
    What a text component actually occupied, as opposed to the region it was offered.

    This exists because of a defect class that only a render catches. An archetype that places
    a rule at `top + allocated_height` is positioning against the BOX, and the box is always
    at least as tall as the text -- so on any slide where the text came up short, the rule
    floats; and where the allocation was optimistic, the rule lands ON the text. In the first
    render of this package the accent rule struck through the headline on four slides out of
    ten, and every one of those slides was structurally valid, passed every schema check, and
    would have gone on a screen.

    `bottom` is the measured bottom edge of the TEXT. Advance from that.
    """

    shape: object
    box: Box
    text_height: int  # EMU actually occupied by the laid-out text
    size_pt: float  # the size finally used (may differ if shrinking was permitted)

    @property
    def bottom(self) -> int:
        return self.box.top + self.text_height


class Ground(Enum):
    """
    What a component is sitting on. Determines the entire legal foreground set.

    This is the mechanism that makes contrast failures unreachable rather than merely
    detectable. LIGHT and ALT are white and off-white; DARK is the brand's dark ground;
    BRAND is the primary brand colour used as a full field.
    """

    LIGHT = "light"
    ALT = "alt"
    DARK = "dark"
    BRAND = "brand"


@dataclass(frozen=True)
class Ctx:
    """
    Everything a component needs: the brand, the grid, and what it is standing on.

    Passed down rather than looked up globally, so two slides with different grounds can be
    built in the same run without a mode flag leaking between them.
    """

    brand: Brand
    ground: Ground = Ground.LIGHT
    grid: Grid = GRID

    def on(self, ground: Ground) -> "Ctx":
        return Ctx(self.brand, ground, self.grid)

    # -- the legal foreground set for this ground ---------------------------

    @property
    def canvas(self) -> str:
        return {
            Ground.LIGHT: self.brand.color("canvas"),
            Ground.ALT: self.brand.color("canvas_alt"),
            Ground.DARK: self.brand.color("canvas_dark"),
            Ground.BRAND: self.brand.color("canvas_brand"),
        }[self.ground]

    @property
    def dark_ground(self) -> bool:
        return self.ground in (Ground.DARK, Ground.BRAND)

    @property
    def ink(self) -> str:
        """Body text."""
        return (
            self.brand.color("on_dark") if self.dark_ground else self.brand.color("ink")
        )

    @property
    def headline(self) -> str:
        """
        Headline colour.

        On a light ground this is the brand blue, per KPMG's own published rule that headlines
        are set in KPMG Blue on light and white on dark. On a dark or brand ground it inverts
        to white -- the brand blue on the dark blue ground measures 1.41:1 and is unreadable.
        """
        return (
            self.brand.color("on_dark")
            if self.dark_ground
            else self.brand.color("headline")
        )

    @property
    def muted(self) -> str:
        return (
            self.brand.color("on_dark_muted")
            if self.dark_ground
            else self.brand.color("muted")
        )

    @property
    def faint(self) -> str:
        return (
            self.brand.color("on_dark_muted")
            if self.dark_ground
            else self.brand.color("faint")
        )

    @property
    def accent(self) -> str:
        """
        The accent that is legal on THIS ground.

        On light this is Cobalt. On dark it inverts to the light tint, because the mid-blue
        accent on the dark ground is illegal at every size. This inversion is the single
        easiest mistake to make with a blue-heavy palette and the reason accent is a property
        rather than a constant.
        """
        if self.ground is Ground.DARK:
            return self.brand.color("on_dark_accent")
        if self.ground is Ground.BRAND:
            return self.brand.color("on_brand_accent")
        return self.brand.color("accent")

    @property
    def rule(self) -> str:
        return (
            self.brand.color("on_dark_muted")
            if self.dark_ground
            else self.brand.color("rule")
        )

    @property
    def rule_strong(self) -> str:
        """
        The heavier of the two greys: column dividers, connectors, unemphasised markers.

        #E5E5E5 and #B2B2B2 are close enough to look interchangeable in code and are not --
        `rule` is the chart track and only that. This exists because three components were
        each re-deriving the dark-ground branch of it by hand, and a rule re-derived in three
        places is a rule that will diverge in two of them.
        """
        return (
            self.brand.color("on_dark_muted")
            if self.dark_ground
            else self.brand.color("rule_strong")
        )

    def assert_legible(
        self, fg: str, size_pt: float, bold: bool = False, *, what: str = "text"
    ) -> None:
        """
        Fail loudly on an illegal pairing.

        Called by the text components. A component that quietly rendered unreadable text would
        pass every structural check and fail only in the room.
        """
        verdict, ratio = contrast_verdict(fg, self.canvas, size_pt, bold)
        if verdict == "FAIL":
            raise ValueError(
                f"{what} at {size_pt}pt would be illegible: {fg} on {self.canvas} is "
                f"{ratio:.2f}:1, below the WCAG AA floor. Change the ground or the role."
            )


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------


def _rect(slide, box: Box, color: str | object, *, alpha: float | None = None):
    """A flat rectangle: no outline, no shadow, no rounded corners. The structural atom."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *box.as_tuple())
    X.set_fill(shape, color, alpha=alpha)
    X.set_line(shape, None)
    X.clear_effects(shape)
    shape.shadow.inherit = False
    return shape


def background(slide, ctx: Ctx) -> None:
    """Fill the slide with the ground colour. Call first on every slide."""
    X.set_slide_background(slide, color=ctx.canvas)


# The house panel widths, as fractions of a 960pt page. These are FREQUENCIES in a measured
# corpus, not proportions anyone derived: 283.3pt occurs 6 times across 3 decks and 324.2pt
# 5 times, while the quote panel is 558.55pt on both TMT p6 and D2 p6. Nobody recovered the
# rule that generated them -- 283.3 is within 0.5pt of 100mm and 324.2 is 114.36mm, which is
# nothing in particular -- so they are used as measurements rather than reasoned about.
PANEL_DEFAULT = 0.29511  # 283.3pt -- the house default
PANEL_WIDE = 0.33770  # 324.2pt -- the wide variant
PANEL_QUOTE = 0.58182  # 558.55pt -- the attributed-quote panel


def full_bleed(
    slide,
    ctx: Ctx,
    color: str | object,
    *,
    side: str = "left",
    fraction: float = PANEL_DEFAULT,
):
    """
    A full-bleed colour block running off the slide edge.

    The most effective single move available in PowerPoint for making a slide look composed
    rather than filled in, and it costs nothing: an asymmetric field of flat colour anchored to
    an edge creates a strong reading order and a place to put a headline that is not "the top".

    `fraction` is deliberately NOT 0.5, and it is no longer a reasoned-about number either.
    The previous default 0.38 was chosen as "close to the golden section's minor part" and has
    no instance behind it anywhere in 95 measured pages. It is now the same 0.29511 the panel
    system uses, which does.
    """
    g = ctx.grid
    if side == "left":
        box = Box(0, 0, int(g.slide_w * fraction), g.slide_h)
    elif side == "right":
        box = Box(
            int(g.slide_w * (1 - fraction)), 0, int(g.slide_w * fraction), g.slide_h
        )
    elif side == "top":
        box = Box(0, 0, g.slide_w, int(g.slide_h * fraction))
    elif side == "bottom":
        box = Box(
            0, int(g.slide_h * (1 - fraction)), g.slide_w, int(g.slide_h * fraction)
        )
    else:
        raise ValueError(f"side must be left/right/top/bottom, got {side!r}")
    return _rect(slide, box, color)


def accent_rule(
    slide,
    ctx: Ctx,
    box: Box,
    *,
    width_pt: float = 30.91,
    thickness_pt: float = 1.465,
    color: str | None = None,
):
    """
    The short rule that sits UNDER THE SECTION NUMBER. Not under every headline.

    THIS FUNCTION WAS A TEMPLATE FIELD PRETENDING TO BE A MOTIF, and the correction is the
    point of its current signature. It used to be 1.6in x 4pt and every archetype's `_head()`
    drew one, so it appeared on 28 of 29 light pages, all starting at exactly x = 8.1%. A mark
    that never varies has stopped being a mark and become a form field -- and in a side-by-side
    against real KPMG pages it was the third-fastest tell a designer would spot.

    KPMG's only equivalent is far smaller and far rarer: a 1.465pt Cobalt rule, 30.91pt long,
    at (56.69, 110.63), sitting under the SECTION NUMBER on section-opener pages, and nowhere
    else. Those are the defaults now. Use it only in the section-number lockup.

    For the vertical column divider -- a genuinely different mark, 0.5pt #B2B2B2, 32 measured
    instances -- see `column_divider()`.
    """
    rule_box = Box(box.left, box.top, points(width_pt), points(thickness_pt))
    return _rect(slide, rule_box, color or ctx.accent)


def column_divider(
    slide,
    ctx: Ctx,
    x: int,
    top: int,
    bottom: int,
    *,
    color: str | None = None,
    thickness_pt: float = 0.5,
):
    """
    The thin vertical rule between the headline column and the body columns.

    Real, measured 32 times, seven of them sharing the exact extent y 127.38 -> 489.99 at
    x = 490.22 -- which is `grid.col(6, 6).left` to within 0.3pt. We did not have this mark at
    all, and it is a large part of why a KPMG page reads as a DOCUMENT with columns rather than
    as a slide with text on both sides.

    Set in `rule_strong` #B2B2B2, not `rule` #E5E5E5. The two are close enough to look
    interchangeable in code and are not: #E5E5E5 is the chart track and only that.
    """
    return _rect(
        slide, Box(x, top, points(thickness_pt), bottom - top), color or ctx.rule_strong
    )


def hairline(
    slide, ctx: Ctx, box: Box, *, color: str | None = None, thickness_pt: float = 0.75
):
    """
    A full-width hairline, for separating rows.

    0.75pt, in the palette's lightest grey. This is what replaces both bullet characters and
    table borders: it separates without enclosing. A 1pt+ rule in a mid grey reads as a border
    and starts to box things in, which adds visual weight the content has not earned.
    """
    return _rect(
        slide,
        Box(box.left, box.top, box.width, points(thickness_pt)),
        color or ctx.rule,
    )


# ---------------------------------------------------------------------------
# Text components
# ---------------------------------------------------------------------------


def eyebrow(slide, ctx: Ctx, box: Box, text_str: str, *, color: str | None = None):
    """
    The small tracked label above a headline: section name, or the question being answered.

    Set in caps at 14pt with +0.1em tracking. Caps ONLY here and only at this size -- caps at
    a readable size destroys word-shape and slows reading measurably, but for a two-or-three
    word label that is scanned rather than read, caps plus tracking reads as a considered
    typographic mark rather than as shouting.
    """
    size = SCALE["footnote"]
    fg = color or ctx.muted
    ctx.assert_legible(fg, size, bold=True, what="eyebrow")
    return T.add_textbox(
        slide,
        box,
        text_str,
        family=ctx.brand.font_minor,
        size_pt=size,
        bold=True,
        color=fg,
        caps=True,
        tracking_pt=tracking_pt("eyebrow", size),
        line_spacing=LINE_SPACING["footnote"],
        anchor="top",
    )


def headline(
    slide,
    ctx: Ctx,
    box: Box,
    text_str: str,
    *,
    role: str = "h2",
    color: str | None = None,
    align: str = "left",
    anchor: str = "top",
    allow_shrink_to: str | None = None,
):
    """
    The slide's claim. Measured before it is written.

    `allow_shrink_to` names the SMALLEST role in the scale this headline may drop to if it does
    not fit. Passing None means the size is fixed and an overflow raises -- which is the right
    default, because a headline that does not fit is nearly always a headline carrying two
    ideas, and the correct fix is to cut it rather than to shrink it.

    When shrinking is permitted the size still comes from the scale, never from a continuous
    search, so sibling slides stay siblings.
    """
    fg = color or ctx.headline
    size = SCALE[role]
    ctx.assert_legible(fg, size, bold=True, what="headline")

    measured = T.fit(
        text_str,
        box,
        ctx.brand.font_major,
        size,
        bold=True,
        tracking_pt=tracking_pt(role, size),
        line_spacing=LINE_SPACING.get(role, 1.1),
    )

    if not measured.fits:
        if allow_shrink_to is None:
            raise ValueError(
                f"headline does not fit at {role} ({size}pt): {measured.describe()}\n"
                f"  text: {text_str!r}\n"
                f"  CUT WORDS rather than shrinking -- a headline that overflows is usually "
                f"carrying two ideas. If shrinking is genuinely right, pass allow_shrink_to."
            )
        floor = SCALE[allow_shrink_to]
        allowed = [s for s in sorted(set(SCALE.values())) if floor <= s <= size]
        chosen = T.largest_size_that_fits(
            text_str,
            box,
            ctx.brand.font_major,
            allowed,
            bold=True,
            line_spacing=LINE_SPACING.get(role, 1.1),
        )
        if chosen is None:
            raise ValueError(
                f"headline does not fit even at the {allow_shrink_to} floor ({floor}pt): "
                f"{measured.describe()}\n  text: {text_str!r}\n  The text must be cut."
            )
        size = chosen

    final = T.fit(
        text_str,
        box,
        ctx.brand.font_major,
        size,
        bold=True,
        tracking_pt=tracking_pt(role, size),
        line_spacing=LINE_SPACING.get(role, 1.1),
    )
    shape = T.add_textbox(
        slide,
        box,
        text_str,
        family=ctx.brand.font_major,
        size_pt=size,
        bold=True,
        color=fg,
        align=align,
        anchor=anchor,
        tracking_pt=tracking_pt(role, size),
        line_spacing=LINE_SPACING.get(role, 1.1),
    )
    return Placed(shape, box, points(final.height_pt), size)


def body(
    slide,
    ctx: Ctx,
    box: Box,
    text_str: str,
    *,
    role: str = "body",
    color: str | None = None,
    align: str = "left",
    anchor: str = "top",
    strict: bool = True,
):
    """
    Supporting prose. Always left-aligned by default.

    Centre-aligned body text is one of the most reliable amateur tells: a ragged left edge
    gives the eye no consistent return point, so multi-line centred text is measurably slower
    to read. Centring is legitimate for a single short line and essentially nothing else.

    `strict=False` NOW STEPS THE SIZE DOWN. IT USED TO MEAN "OVERFLOW SILENTLY", and that was
    the single largest source of defects in this package's second rebuild: a text-overflow gate
    added afterwards found 23 of them across 37 pages, every one a run of prose drawn straight
    through whatever sat beneath it. Nothing else could see them -- an overflowing text frame is
    valid XML, on the slide, in the palette, above the type floor. Only a render or a
    re-measurement finds it, and by then it has been shipped.

    So the three states are now:
      strict=True    (default) does not fit at `role` -> RAISE. The text is wrong.
      strict=False   step down through body / small / caption until it fits.
      neither        if it does not fit at `caption` either, it overflows and is reported,
                     because going below 16pt would silently surrender the documented
                     room-legibility deviation and nobody would ever see that happen.
    """
    fg = color or ctx.ink

    size = SCALE[role]
    measured = T.fit(
        text_str, box, ctx.brand.font_minor, size, line_spacing=LINE_SPACING["body"]
    )

    if not measured.fits:
        if strict:
            raise ValueError(
                f"body text does not fit at {role} ({size}pt): {measured.describe()}\n"
                f"  text: {text_str!r}\n"
                f"  Cut the text or give it a larger region. Do not reduce below "
                f"{SCALE['footnote']}pt -- that is the legibility floor for a large room."
            )
        # The ladder is ordered, and it never goes below `caption`. Steps smaller than the
        # role are the only candidates: a `lead` that does not fit must not become a `body`
        # that is bigger than the caller asked for.
        for candidate in ("body", "small", "caption"):
            if SCALE[candidate] >= size:
                continue
            trial = T.fit(
                text_str,
                box,
                ctx.brand.font_minor,
                SCALE[candidate],
                line_spacing=LINE_SPACING["body"],
            )
            if trial.fits:
                size, measured = SCALE[candidate], trial
                break

    ctx.assert_legible(fg, size, what="body")
    shape = T.add_textbox(
        slide,
        box,
        text_str,
        family=ctx.brand.font_minor,
        size_pt=size,
        color=fg,
        align=align,
        anchor=anchor,
        line_spacing=LINE_SPACING["body"],
    )
    return Placed(shape, box, points(measured.height_pt), size)


def source_line(slide, ctx: Ctx, text_str: str):
    """
    The source note, bottom-left, at the legibility floor.

    Present on every exhibit without exception. A figure with no source is an assertion; a
    figure with one is evidence, and in a professional-services context the difference is the
    whole point. It is also the cheapest credibility signal available.
    """
    g = ctx.grid
    box = Box(
        g.margin_x,
        g.slide_h - g.margin_bottom + points(6),
        int(g.slide_w * 0.62),
        points(28),
    )
    return T.add_textbox(
        slide,
        box,
        text_str,
        family=ctx.brand.font_minor,
        size_pt=SCALE["footnote"],
        color=ctx.faint,
        line_spacing=LINE_SPACING["footnote"],
    )


def page_number(slide, ctx: Ctx, number: int, *, color: str | None = None):
    """Page number, bottom-right, aligned to the same baseline as the source line."""
    g = ctx.grid
    box = Box(
        g.slide_w - g.margin_x - inches(1.0),
        g.slide_h - g.margin_bottom + points(6),
        inches(1.0),
        points(28),
    )
    return T.add_textbox(
        slide,
        box,
        str(number),
        family=ctx.brand.font_minor,
        size_pt=SCALE["footnote"],
        color=color or ctx.faint,
        align="right",
    )


# ---------------------------------------------------------------------------
# Running chrome -- what makes a page read as a DOCUMENT rather than a slide
# ---------------------------------------------------------------------------
#
# 100% of non-cover, non-back-matter KPMG pages carry both a top nav band (y 12.4 -> 36.0) and
# a footer band (y 513.7 -> 529.1). Our footer was 6.6pt wide -- a single digit -- against
# KPMG's 845pt of document furniture. That gap is the single cheapest signal that an artifact
# is a slide template rather than a published document, and it costs two components to close.
#
# TWO DELIBERATE DEVIATIONS, both stated rather than silently taken:
#
# (a) NAV LABELS GO AT 10PT, NOT KPMG'S 8. Arial is ~57% wider per character than Univers, so
#     8pt Arial is not the same mark at all -- and 8pt from the back of a 140-person room is
#     decorative rather than legible. 10pt still sits well below the 20pt body and reads as
#     furniture. This is the one item in the chrome that could be a net negative: it adds a
#     band of small type nobody reads, purely for the document texture it signals. It has NOT
#     been tested in the room. `Deck(nav=False)` drops it and keeps the footer, which is the
#     fallback if the room says otherwise.
#
# (b) THE THREE NAV BUTTONS ARE NOT DRAWN. KPMG's measured chrome ends with three circles
#     (back / home / forward, diameter 15.61pt) at the right of the band. Those are the
#     navigation affordances of an INTERACTIVE PDF. In a .pptx they would be three circles that
#     do nothing -- a control that cannot be operated is not furniture, it is clip art, and the
#     module docstring's refusal of clip art does not have an exception for clip art that was
#     measured. The band ends at the last tab.


def nav_strip(
    slide,
    ctx: Ctx,
    sections: Sequence[str],
    current: int | None = None,
    *,
    right_limit_pt: float = 870.0,
    dark_range: tuple[int, int] | None = None,
) -> list[int]:
    """
    The top section navigation. Returns the divider x-positions in EMU.

    The return value is not incidental. In D1/D2, 11 of 12 panel edges land on a nav-divider
    position, ten of them exact to <=0.11pt -- so the dividers are the only positional rule the
    panel system demonstrably has, and a caller can pass one to `panel(edge_x=...)`.

    Measured geometry (D1/D2):

        band        y 12.4 -> 36.0pt
        labels      2 lines, tops at y 15.84 / 25.44, Univers Light 8pt (we set 10 -- above)
        inactive    #B2B2B2
        active      Bold #1E49E2 on a white rect 92.60 x 21.06pt
        dividers    0.5pt #B2B2B2 verticals, 13.80pt tall, y 18.87 -> 32.67
        D1/D2 divider x: 144.68 - 245.12 - 324.17 - 440.30 - 558.55 - 676.80 - 792.41

    THOSE SEVEN DIVIDER POSITIONS ARE NOT REPRODUCED HERE and the reason is a genuine unknown.
    All four measured CEO decks happen to have exactly EIGHT sections; we have six, and nobody
    could establish how the template degrades -- whether the tabs keep their widths and leave a
    gap, or redistribute. This lays them out evenly across the same band, which is the
    assumption that keeps the strip looking deliberate at any count. If a real 6-tab KPMG page
    ever turns up, that is a measurement worth taking.

    `current` is a zero-based index, or None on a page belonging to no section.

    `dark_range` is the (x0, x1) EMU span of any colour panel bleeding under the band. Tabs
    whose text falls inside it invert to the on-dark tints. The band is drawn AFTER the panel
    and runs straight over it -- which is what D1 p6 does, where the quote panel bleeds to
    y=0 and the first three tabs sit reversed on top of it. Splitting the band around the panel
    would be the intuitive move and is not what the system does.
    """
    g = ctx.grid
    n = len(sections)
    if n < 2:
        raise ValueError("nav_strip needs at least two sections")

    left = g.margin_x
    right = points(right_limit_pt)
    cell_w = (right - left) // n

    def over_panel(x0: int, x1: int) -> bool:
        # THE MIDPOINT OF THE INK, NOT ANY OVERLAP. A tab straddling the panel edge overlaps
        # the panel by one pixel and was taking the reversed colour for its whole width --
        # pale blue type on white, at 1.32:1. What decides a run's colour is what MOST OF IT
        # is standing on, which is the same rule verify._effective_bg applies from the other
        # side. Two places, one rule; they disagreed and the render was the tie-breaker.
        if dark_range is None:
            return False
        mid = (x0 + x1) // 2
        return dark_range[0] <= mid <= dark_range[1]

    dividers: list[int] = []
    size = 10.0
    inactive = ctx.brand.color("nav_inactive")
    active = ctx.brand.color("accent")

    for i, label in enumerate(sections):
        cell_left = left + i * cell_w
        if i > 0:
            # The divider sits at the cell boundary, 13.80pt tall, vertically centred in the
            # band -- shorter than the band so it separates without ruling the page.
            x = cell_left - points(6)
            dividers.append(x)
            _rect(
                slide,
                Box(x, points(18.87), points(0.5), points(13.80)),
                ctx.brand.color("on_dark_muted")
                if over_panel(x, x + points(0.5))
                else ctx.brand.color("rule_strong"),
            )

        # A DARK GROUND IS DARK EVERYWHERE, NOT ONLY UNDER A PANEL. `over_panel` answers "is
        # this tab standing on the bleeding colour field", which is the right question on a
        # white page carrying a panel and the wrong one on a slide whose whole ground is dark:
        # `dark_range` is None there, so every tab took the grey it uses on white and the strip
        # went out at #666666 on #0C233C -- 2.77:1, caught by verify only once an exhibit was
        # first built on Ground.DARK. The ground is the outer condition; the panel is a local
        # exception to a light one.
        dark = ctx.ground in (Ground.DARK, Ground.BRAND) or over_panel(
            cell_left, cell_left + cell_w
        )
        is_current = current is not None and i == current
        if is_current and not dark:
            # The active tab sits on a white patch. On a panel it does not: a white rectangle
            # punched into a bleeding colour field reads as a hole, so there the active tab is
            # marked by weight and the pale tint instead.
            _rect(
                slide,
                Box(cell_left - points(6), points(12.4), cell_w, points(21.06)),
                ctx.brand.color("canvas"),
            )
        if dark:
            # #B2B2B2 on Cobalt is 3.19:1 -- legal for large text and this is 10pt furniture,
            # so it inverts to the pale tint (5.66:1) rather than staying grey. Same rule as
            # everywhere else in this package: the accent inverts on a colour ground.
            inactive_here = ctx.brand.color("on_dark_accent")
            active_here = ctx.brand.color("on_dark")
        else:
            inactive_here, active_here = inactive, active

        T.add_textbox(
            slide,
            Box(cell_left, points(15.84), cell_w - points(12), points(20.0)),
            label,
            family=ctx.brand.font_minor,
            size_pt=size,
            bold=is_current,
            color=active_here if is_current else inactive_here,
            line_spacing=1.15,
            anchor="top",
        )

    return dividers


def footer(
    slide,
    ctx: Ctx,
    doc_title: str,
    page: int,
    *,
    legal: str | None = None,
    dark_range: tuple[int, int] | None = None,
):
    """
    The document footer: title and page number ranged right, standing text ranged left.

    Measured: 10pt #00338D title + " | " + a BOLD page number, ranged right to x = 903.40, at
    y ~519.4; a 6pt legal line at x = 56.69, y = 520.90 in licensed mode.

    The bold page number is not a flourish. It is what makes the right end of the footer read
    as a page reference rather than as a run-on of the title, and it is present on every one of
    the 95 measured pages.

    `legal` is emitted at 6pt in the faint grey. In UNLICENSED mode that slot carries the
    deck's own standing line, never KPMG's copyright -- a third party asserting KPMG's
    copyright over its own material is worse than omitting the line.

    On `#989898`: KPMG's own legal grey measures 2.88:1 on white and fails AA at every size, so
    the brand file maps this to #666666. That deviation is documented in brands/kpmg.json and
    is deliberate; do not "fix" it back.
    """
    g = ctx.grid
    y = points(513.7)
    right = points(903.40)

    def ink_at(x0: int, x1: int) -> str:
        over = dark_range is not None and x0 < dark_range[1] and x1 > dark_range[0]
        if over or ctx.dark_ground:
            return ctx.brand.color("on_dark")
        return ctx.brand.color("headline")

    if legal:
        w = int(g.slide_w * 0.62)
        # The legal line is LEFT-ALIGNED in a wide box, so what it sits on is decided at its
        # left edge, not across the box. Testing the whole box put it on the panel whenever a
        # panel existed anywhere to its right.
        over = dark_range is not None and dark_range[0] <= g.margin_x <= dark_range[1]
        T.add_textbox(
            slide,
            Box(g.margin_x, points(516.5), w, points(16)),
            legal,
            family=ctx.brand.font_minor,
            size_pt=8.0,
            # #B2B2B2 on Cobalt is 3.19:1, which is legal for large text and this is 8pt
            # furniture. Same inversion as the nav strip: on a colour ground the quiet tone is
            # the pale tint (5.66:1), never a grey.
            color=ctx.brand.color("on_dark_accent")
            if (over or ctx.dark_ground)
            else ctx.faint,
            line_spacing=1.2,
            anchor="top",
        )

    # Title and page number are ONE text frame with two runs, so the " | " separator sits at a
    # real typographic gap rather than at whatever two independently-placed boxes happen to
    # leave. Ranged right, so the page number is the last ink on the line at every page count.
    box_w = int(g.slide_w * 0.45)
    box = Box(right - box_w, y, box_w, points(18))
    # The colour is decided by what the RIGHT END sits on, since the text is ranged right and
    # a short title leaves the left of the box empty. This is the same defect the page-number
    # fix caught: a run coloured for the slide default rather than for its actual ground
    # measured 1.18:1 on five slides and was invisible.
    fg = ink_at(right - points(240), right)
    shape = T.add_textbox(
        slide,
        box,
        f"{doc_title}  |  ",
        family=ctx.brand.font_minor,
        size_pt=10.0,
        color=fg,
        align="right",
        anchor="top",
        line_spacing=1.2,
    )
    para = shape.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = str(page)
    run.font.size = para.runs[0].font.size
    run.font.name = ctx.brand.font_minor
    run.font.bold = True
    T._apply_color(run, fg, None)
    return shape


def classification(slide, ctx: Ctx, text_str: str):
    """
    The mandated document-classification marker.

    Corporate furniture, not content. It goes bottom-left at the smallest size in the scale,
    quiet enough to ignore and present enough to satisfy the requirement. Never style it to be
    noticed, and never omit it because it spoils a composition -- on a Big-Four deck a missing
    classification marker is a compliance finding, and a designer's objection does not survive
    contact with risk review.
    """
    g = ctx.grid
    box = Box(g.margin_x, g.slide_h - points(26), int(g.slide_w * 0.5), points(20))
    return T.add_textbox(
        slide,
        box,
        text_str,
        family=ctx.brand.font_minor,
        size_pt=SCALE["micro"],
        color=ctx.faint,
    )


# ---------------------------------------------------------------------------
# Content components
# ---------------------------------------------------------------------------


def stat(
    slide,
    ctx: Ctx,
    box: Box,
    value: str,
    label: str,
    *,
    caption: str | None = None,
    color: str | None = None,
    role: str = "display",
    label_height: int | None = None,
):
    """
    An oversized figure with its label beneath.

    The rule that makes this work, and that most attempts get backwards: the NUMBER is the
    headline and the label is subordinate. A 60pt figure over a 20pt label reads instantly. A
    32pt figure over an 18pt label reads as a form field, because the ratio is too close to
    establish which one is the idea.

    The number is set in the accent, the label in body ink. Colouring both is a common error
    that removes the distinction the size was establishing.
    """
    # THE FIGURE IS SIZED TO ITS COLUMN, AND ITS HEIGHT IS MEASURED.
    #
    # `SCALE[role]` with a `size * 1.15` allocation assumed the value fits on one line at the
    # role's size. In a 145pt aside column "57 min" wraps to two lines at display size, and the
    # allocation then reserved one -- so the label was drawn through the word "min". Stepping
    # down the scale until the figure fits on one line keeps it a figure; letting it wrap makes
    # it a paragraph that happens to be large.
    fg = color or ctx.accent
    size = SCALE[role]
    for candidate in sorted(
        (x for x in set(SCALE.values()) if x <= SCALE[role]), reverse=True
    ):
        if (
            T.fit(
                value,
                Box(0, 0, box.width, points(candidate * 4)),
                ctx.brand.font_major,
                candidate,
                bold=True,
                tracking_pt=tracking_pt(role, candidate),
                line_spacing=LINE_SPACING.get(role, 1.0),
            ).line_count
            == 1
        ):
            size = candidate
            break
    ctx.assert_legible(fg, size, bold=True, what="stat value")

    value_h = points(
        T.fit(
            value,
            Box(0, 0, box.width, points(size * 4)),
            ctx.brand.font_major,
            size,
            bold=True,
            tracking_pt=tracking_pt(role, size),
            line_spacing=LINE_SPACING.get(role, 1.0),
        ).height_pt
    )
    T.add_textbox(
        slide,
        box.top_slice(value_h),
        value,
        family=ctx.brand.font_major,
        size_pt=size,
        bold=True,
        color=fg,
        tracking_pt=tracking_pt(role, size),
        line_spacing=LINE_SPACING.get(role, 1.0),
        anchor="top",
    )

    label_top = box.top + value_h + points(SPACE["xs"])
    label_box = Box(box.left, label_top, box.width, box.bottom - label_top)
    label_fit = T.fit(
        label,
        label_box,
        ctx.brand.font_minor,
        SCALE["body"],
        line_spacing=LINE_SPACING["body"],
    )
    T.add_textbox(
        slide,
        label_box,
        label,
        family=ctx.brand.font_minor,
        size_pt=SCALE["body"],
        color=ctx.ink,
        line_spacing=LINE_SPACING["body"],
    )

    if caption:
        # Placed below the label's MEASURED bottom. A fixed label height collided with the
        # caption on every three-line label in the first render.
        # `label_height`, when supplied, overrides this stat's own measurement with the
        # tallest label in the row, so every caption in a metrics row sits on ONE baseline.
        # Without it each caption floats to its own label's depth and the row reads ragged --
        # the sibling-consistency rule applied within a slide rather than across slides.
        used_h = (
            label_height if label_height is not None else points(label_fit.height_pt)
        )
        cap_top = label_top + used_h + points(SPACE["xs"])
        # `max(points(20), ...)` was here and it is why a caption ran into the footer band: a
        # floor on the BOX does nothing once the box has already passed the bottom of the
        # region -- it just guarantees 20pt of ink somewhere off the page. Draw the caption
        # only where there is real room. A missing caption is a page that says slightly less;
        # a caption over the footer is a broken page.
        room = box.bottom - cap_top
        if room < points(18):
            return
        cap_box = Box(box.left, cap_top, box.width, room)
        T.add_textbox(
            slide,
            cap_box,
            caption,
            family=ctx.brand.font_minor,
            size_pt=SCALE["footnote"],
            color=ctx.faint,
            line_spacing=LINE_SPACING["footnote"],
        )


def callout(
    slide,
    ctx: Ctx,
    box: Box,
    text_str: str,
    *,
    role: str = "lead",
    color: str | None = None,
):
    """
    The "so what" -- the conclusion drawn from the exhibit beside it.

    Rendered as a LEFT ACCENT BAR plus text, not as a bordered or filled box. This is a
    deliberate choice worth stating: a filled callout box competes with the exhibit it is
    commenting on, and a bordered one adds a rectangle to a slide that already has enough
    edges. A 4pt bar at the left margin of the text does the same job -- marks this text as
    different in kind -- while adding almost no visual weight.
    """
    bar_w = points(4)
    gap = points(SPACE["sm"])
    accent_color = color or ctx.accent

    _rect(slide, Box(box.left, box.top, bar_w, box.height), accent_color)

    text_box = Box(box.left + bar_w + gap, box.top, box.width - bar_w - gap, box.height)
    size = SCALE[role]
    ctx.assert_legible(ctx.ink, size, what="callout")
    return T.add_textbox(
        slide,
        text_box,
        text_str,
        family=ctx.brand.font_minor,
        size_pt=size,
        color=ctx.ink,
        line_spacing=LINE_SPACING["h3"],
        anchor="middle",
    )


def rows(
    slide,
    ctx: Ctx,
    box: Box,
    items: Sequence[tuple[str, str]],
    *,
    rule_between: bool = True,
    label_role: str = "h3",
    body_role: str = "small",
):
    """
    A stack of label/description pairs separated by hairlines. THE REPLACEMENT FOR BULLETS.

    Each item is (label, description). The label carries the idea and the description supports
    it -- which is the same answer-first discipline a good slide headline follows, applied one
    level down.

    Capped at five items, and the cap is the point rather than a limitation. Working memory
    holds about four items; a list of seven is not a list anyone retains, it is a document
    someone read aloud. If there are more than five, the grouping is wrong and wants a level.
    """
    if not items:
        raise ValueError("rows() needs at least one item")
    if len(items) > 5:
        raise ValueError(
            f"rows() was given {len(items)} items. The cap is five, and it is deliberate: "
            f"beyond about four the list stops being retained. Regroup into fewer, higher-level "
            f"items, or split across two slides."
        )

    n = len(items)
    gap = points(SPACE["md"])
    row_h = (box.height - gap * (n - 1)) // n

    # THE GAPS ARE FIXED AND THE ROWS TAKE WHAT IS LEFT, so a box that is short for its item
    # count leaves `row_h` at or below zero and every row after the first is placed ABOVE the
    # one before it. What that produced was four text boxes with negative height stacked on one
    # another -- valid XML, invisible to every check except the geometry one, and drawn by
    # PowerPoint straight through the row beneath. The package's rule is that text is measured
    # before it is written and never silently made to fit; a box too small for its rows is the
    # same finding as a headline that will not fit, so it is raised rather than drawn.
    min_row = points(SCALE[label_role] * 1.35) + points(SCALE["small"] * 1.2)
    if row_h < min_row:
        raise ValueError(
            f"rows(): {n} rows need at least {(min_row * n + gap * (n - 1)) / 12700:.0f}pt "
            f"and the region is {box.height / 12700:.0f}pt -- each row gets "
            f"{row_h / 12700:.1f}pt against a {min_row / 12700:.0f}pt minimum. Use fewer rows, "
            f"or give the component a taller region. Do not shrink the type: 14pt is the "
            f"legibility floor for a large room."
        )

    for i, (label, description) in enumerate(items):
        top = box.top + i * (row_h + gap)
        if rule_between and i > 0:
            hairline(slide, ctx, Box(box.left, top - gap // 2, box.width, 0))

        label_h = points(SCALE[label_role] * 1.35)
        T.add_textbox(
            slide,
            Box(box.left, top, box.width, label_h),
            label,
            family=ctx.brand.font_major,
            size_pt=SCALE[label_role],
            bold=True,
            color=ctx.headline,
            line_spacing=LINE_SPACING["h3"],
            tracking_pt=tracking_pt(label_role, SCALE[label_role]),
        )
        if description:
            T.add_textbox(
                slide,
                Box(
                    box.left,
                    top + label_h + points(SPACE["hair"]),
                    box.width,
                    row_h - label_h - points(SPACE["hair"]),
                ),
                description,
                family=ctx.brand.font_minor,
                size_pt=SCALE[body_role],
                color=ctx.ink,
                line_spacing=LINE_SPACING["body"],
            )


def steps(
    slide,
    ctx: Ctx,
    box: Box,
    items: Sequence[tuple[str, str]],
    *,
    connector: bool = True,
):
    """
    A numbered horizontal sequence: process, timeline, or staged argument.

    Numerals are set in the accent at display weight and the connector is a hairline running
    BEHIND them, which is what makes it read as one continuous process rather than as separate
    cards. Drawing a box around each step is the obvious move and the wrong one -- it turns a
    sequence into a set, losing the ordering that was the whole reason to use this component.
    """
    if not 2 <= len(items) <= 5:
        raise ValueError(f"steps() takes 2 to 5 items, got {len(items)}")

    n = len(items)
    gap = points(SPACE["lg"])
    num_size = SCALE["h1"]
    num_h = points(num_size * 1.2)

    if connector:
        y = box.top + num_h // 2
        _rect(slide, Box(box.left, y, box.width, points(0.75)), ctx.rule)

    for i, (label, description) in enumerate(items):
        cell = box.row(n, i, gap=gap)

        # A patch of ground behind each numeral, so the connector appears to pass behind it.
        _rect(
            slide, Box(cell.left, box.top, points(num_size * 0.85), num_h), ctx.canvas
        )

        T.add_textbox(
            slide,
            Box(cell.left, box.top, cell.width, num_h),
            str(i + 1),
            family=ctx.brand.font_major,
            size_pt=num_size,
            bold=True,
            color=ctx.accent,
            tracking_pt=tracking_pt("h1", num_size),
            line_spacing=1.0,
        )

        label_top = box.top + num_h + points(SPACE["sm"])
        label_h = points(SCALE["h3"] * 1.4)
        T.add_textbox(
            slide,
            Box(cell.left, label_top, cell.width, label_h),
            label,
            family=ctx.brand.font_major,
            size_pt=SCALE["h3"],
            bold=True,
            color=ctx.headline,
            line_spacing=LINE_SPACING["h3"],
        )
        if description:
            T.add_textbox(
                slide,
                Box(
                    cell.left,
                    label_top + label_h + points(SPACE["hair"]),
                    cell.width,
                    box.bottom - label_top - label_h,
                ),
                description,
                family=ctx.brand.font_minor,
                size_pt=SCALE["small"],
                color=ctx.ink,
                line_spacing=LINE_SPACING["body"],
            )


def quote(
    slide, ctx: Ctx, box: Box, text_str: str, attribution: str, *, role: str = "h1"
):
    """
    A pull quote, WITHOUT the oversized opening glyph.

    THIS DOCSTRING USED TO STATE A RULE THAT MEASUREMENT REFUTED, and the refutation is worth
    recording rather than quietly deleting. It read: "No quotation-mark graphic. An oversized
    decorative quote glyph is clip art with a typographic alibi." That is a reasonable general
    prior and it is wrong about this system specifically -- KPMG's oversized glyph is a fixed-
    size recurring device, #00B8F5 Pacific at a bbox of 60.02 x 49.41pt, appearing at that
    exact size on TMT pages 2, 6, 8, 11, 13 and on D1 p6. It is not decoration that crept in;
    it is a component with a spec.

    The glyph belongs to the ATTRIBUTED PANEL quote, so it lives in `quote_glyph()` and is
    drawn by `deck.quote_panel()`. This function stays glyph-free because it is the quieter
    treatment for a quote on white, where the prior above does hold.

    The detail that separates a real KPMG quote page from a blue one: the glyph is PACIFIC on
    the Cobalt panel, not white. A white glyph is the obvious choice and the wrong one.
    """
    size = SCALE[role]
    quote_h = box.height - points(SCALE["body"] * 2.2)

    T.add_textbox(
        slide,
        Box(box.left, box.top, box.width, quote_h),
        f"“{text_str}”",
        family=ctx.brand.font_major,
        size_pt=size,
        color=ctx.headline,
        tracking_pt=tracking_pt(role, size),
        line_spacing=LINE_SPACING.get(role, 1.1),
        anchor="bottom",
    )
    T.add_textbox(
        slide,
        Box(
            box.left,
            box.bottom - points(SCALE["body"] * 2.0),
            box.width,
            points(SCALE["body"] * 2.0),
        ),
        attribution,
        family=ctx.brand.font_minor,
        size_pt=SCALE["body"],
        color=ctx.muted,
        line_spacing=LINE_SPACING["body"],
    )


# `section_numeral()` WAS HERE AND HAS BEEN DELETED. It drew a 400pt numeral at 12% alpha
# bleeding off the right edge of a divider -- a composition that appears NOWHERE in 95 measured
# KPMG pages. It had already been removed from `deck.section()`, but leaving the function in
# the module is how a removed device comes back: the next caller looking for "a divider device"
# finds it, and it is the only thing in here that looks like one. A vocabulary is defined as
# much by what it refuses as by what it offers, so the refusal is the deletion, not a comment
# on the call site.
#
# What replaces it is `section_number()` below: the real KPMG lockup, which is small.


def quote_glyph(slide, ctx: Ctx, left: int, top: int, *, color: str | None = None):
    """
    The oversized opening quotation mark. A fixed-size KPMG device, not decoration.

    Measured identically on TMT pages 2, 6, 8, 11 and 13 and on D1 p6: an ink bbox of
    60.02 x 49.41pt. Because it recurs at that EXACT size across pages and decks, it is drawn
    at a fixed size rather than scaled to its container -- scaling it would be the thing that
    turned it back into decoration.

    PACIFIC #00B8F5 on the Cobalt panel. Not white. This is the detail that separates a real
    KPMG quote page from a merely blue one, and white is the choice everyone makes by default.

    The CLOSING mark is typed as a plain " at the end of the quotation text, never drawn as a
    second glyph. KPMG does it that way and the asymmetry is deliberate: the opening mark is a
    graphic element announcing the register, the closing one is punctuation.
    """
    # 60.02 x 49.41pt is the INK bbox. A text frame's box must be larger than the ink it holds,
    # so the point size is derived from the ink height and the frame given generous slack --
    # measure the mark, not the box.
    size_pt = 118.0  # Arial Bold '"' ink cap-height lands ~49.4pt at this size
    box = Box(
        left, top - points(size_pt * 0.16), points(size_pt * 1.2), points(size_pt * 1.1)
    )
    return T.add_textbox(
        slide,
        box,
        "“",
        family=ctx.brand.font_major,
        size_pt=size_pt,
        bold=True,
        color=color or ctx.brand.color("accent_bright"),
        align="left",
        anchor="top",
        line_spacing=1.0,
    )


def section_number(
    slide,
    ctx: Ctx,
    number: int | str,
    *,
    left: int | None = None,
    top: int | None = None,
) -> int:
    """
    The section-number lockup: a two-digit numeral with a short rule under it.

    Measured span-by-span off D1 p5 and TMT p5, which are identical to 0.1pt:

        "01"   Univers Bold 26.36pt  #1E49E2  at (56.69, 80.70)
        rule   1.465pt #1E49E2, x 56.69 -> 87.61 (30.91pt long), at y 110.63

    Returns the y below the rule, so the headline can advance from a measurement.

    Set at 28pt rather than 26.36: 26.36 is Univers, and the scale's `h3` is the nearest real
    step for Arial. Sizes come from the scale even when a measurement is available in between,
    because a scale with one off-step value in it has stopped being a scale.
    """
    g = ctx.grid
    x = g.margin_x if left is None else left
    y = points(80.70) if top is None else top
    size = SCALE["h3"]
    fg = ctx.brand.color("on_dark") if ctx.dark_ground else ctx.brand.color("accent")

    T.add_textbox(
        slide,
        Box(x, y, points(120), points(size * 1.4)),
        f"{int(number):02d}" if str(number).isdigit() else str(number),
        family=ctx.brand.font_major,
        size_pt=size,
        bold=True,
        color=fg,
        line_spacing=1.0,
        anchor="top",
    )
    rule_y = points(110.63) if top is None else y + points(size * 1.06)
    accent_rule(slide, ctx, Box(x, rule_y, 0, 0), color=fg)
    return rule_y + points(1.465)


# ---------------------------------------------------------------------------
# The panel system -- the single most load-bearing correction to this package
# ---------------------------------------------------------------------------
#
# MEASURED ACROSS 95 PAGES of five current KPMG thought-leadership decks:
#   - 34 occurrences of a full-height flat colour panel bleeding off one vertical edge
#   - ALL 24 measured full-height panels are Cobalt #1E49E2. NONE is #00338D.
#   - two house widths carry most of them: 283.3pt (n=6, 3 decks) and 324.2pt (n=5)
#   - of 24 non-flood panels: 17 bleed RIGHT, 7 bleed LEFT -- right-dominant, alternating
#   - content inset inside the panel is 56.69pt, the same as the page margin
#   - corner radius 0, no stroke, no shadow, always bleeding, never a floating card
#
# THE PORTRAIT PHOTO WINDOW THIS REPLACES OCCURS 5 TIMES IN 95 PAGES AND EVERY ONE IS A
# COVER. Zero of 90 interior pages carry it. The panel is not a stand-in for a photograph we
# do not have -- it is what KPMG's interiors actually use, and the 2015 brand book's
# "what not to do" page explicitly rejects "a no-photo layout without using a vertical
# element". Of nine sanctioned layouts, every one contains a vertical element.
#
# THERE IS NO GRID BEHIND THE PANELS, and this was tested rather than assumed. Fitting 23
# panel edges to a k-column model leaves a 36.93pt mean residual at k=3 and explains 5 of 23;
# k=15 "wins" only by overfitting. The one real positional rule is local to D1/D2, where 11 of
# 12 panel edges land on a top-nav divider position, ten of them exact to <=0.11pt -- hence
# `edge_x`, so a caller can snap to a divider `nav_strip()` returned. In D3/D4, 0 of 11 land
# on a divider. We have the numbers; nobody recovered the rule that generated them.
#
# WHAT THE PANEL IS FOR, and this is the correction that survived the first rebuild:
# A COLOUR FIELD MARKS A CHANGE OF VOICE, NEVER A CHANGE OF SECTION. KPMG's full-colour
# fields are reserved for the cover, attributed pull-quotes, the methodology page and the
# closing. Their section openers are WHITE WORKING PAGES. We used colour for section
# boundaries and nothing else; the first rebuild changed the hue and left the category error
# in place.


def panel(
    slide,
    ctx: Ctx,
    *,
    side: str = "right",
    fraction: float = PANEL_DEFAULT,
    color: str | None = None,
    edge_x: int | None = None,
) -> tuple[Box, Box]:
    """
    A full-height colour panel bleeding off one vertical edge.

    Returns (panel_box, remainder_box) -- the remainder already inset to the grid margin, so
    a caller can lay content into it directly.

    `color` defaults to the brand accent (Cobalt), NOT the master blue. Measurement confirmed
    this exactly: 24 of 24 full-height panels are #1E49E2 and none is #00338D. Filling a panel
    with the master blue is a characteristic near-miss -- correct palette, wrong role.

    `fraction` defaults to the house width, 283.3pt on a 960pt page. `PANEL_WIDE` (324.2pt)
    and `PANEL_QUOTE` (558.55pt) are the other two measured sizes. A width outside that set is
    not forbidden but has no instance behind it.

    `edge_x`, when given, sets the panel's inner edge in EMU directly and overrides `fraction`.
    Pass a divider position from `nav_strip()` to reproduce the only positional rule the system
    demonstrably has.

    `side` defaults to "right" because that is the corpus majority (17 of 24), reversing this
    function's previous "left" default which was read off a smaller and differently-filtered
    sample.

    The panel ALWAYS bleeds -- full slide height, hard to one edge, never floated with margin
    on all four sides. Observed without exception.
    """
    g = ctx.grid
    fill = color or ctx.brand.color("accent")
    width = (
        int(g.slide_w * fraction)
        if edge_x is None
        else (g.slide_w - edge_x if side == "right" else edge_x)
    )

    if side == "left":
        pbox = Box(0, 0, width, g.slide_h)
        rest = Box(
            width + g.margin_x,
            g.margin_top,
            g.slide_w - width - 2 * g.margin_x,
            g.slide_h - g.margin_top - g.margin_bottom,
        )
    elif side == "right":
        pbox = Box(g.slide_w - width, 0, width, g.slide_h)
        rest = Box(
            g.margin_x,
            g.margin_top,
            g.slide_w - width - 2 * g.margin_x,
            g.slide_h - g.margin_top - g.margin_bottom,
        )
    else:
        raise ValueError(f"panel side must be left or right, got {side!r}")

    _rect(slide, pbox, fill)
    return pbox, rest


def panel_inner(
    ctx: Ctx,
    pbox: Box,
    *,
    top: int | None = None,
    bottom: int | None = None,
    inset_pt: float = 34.3,
) -> Box:
    """
    The content region inside a panel.

    TWO INSETS WERE MEASURED AND THEY DISAGREE: 34.3pt from the panel's outer edge on D1 p12's
    quote, and 57.46pt on D1 p2, which S2 reports as 56.69 -- the page margin. One of those two
    pages is doing something different and nobody established which. Both are real numbers off
    real pages.

    THIS DEFAULTS TO 34.3, AND THE REASON IS OUR TYPE SIZE RATHER THAN A VERDICT ON THE
    EVIDENCE. KPMG sets panel copy at 10.5pt; we set it at 18-20 for a 140-person room, which
    is the stated deviation in the brand file. On the house 283.3pt panel a 56.69pt inset
    leaves 170pt of measure -- 30 characters a line at KPMG's size and SEVENTEEN at ours, which
    is not a column, it is a stack of fragments. The 34.3pt inset leaves 215pt, giving ~24
    characters at 18pt: the same READING measure, reached with a different number because the
    type is twice the size.

    Pass `inset_pt=56.693` where the content is display-scale and the measure does not bind --
    a big statistic, a ring, a two-word label.
    """
    g = ctx.grid
    m = points(inset_pt)
    t = g.margin_top if top is None else top
    b = (g.slide_h - g.margin_bottom) if bottom is None else bottom
    return Box(pbox.left + m, t, pbox.width - 2 * m, b - t)


def stepped_headline(
    slide,
    ctx: Ctx,
    box: Box,
    lines: Sequence[str],
    *,
    role: str = "display",
    step_in: float = 0.055,
    colors: Sequence[str] | None = None,
    on_panel: bool = True,
    allow_shrink_to: str | None = "h3",
):
    """
    The COVER staircase: a multi-line title whose lines step rightward, alternating colour.

    THIS IS A COVER MOVE AND ONLY A COVER MOVE. That is the correction. Measured cover
    x-origins are 72.0 -> 119.3 -> 300.8, constant across D3/D4 -- while every interior
    two-tone headline in the corpus is FLUSH LEFT at 55.4-56.7pt, on every instance. We had
    this exactly inverted: the step applied everywhere and the two-tone nowhere. For the
    interior form see `two_tone_headline()`, which is what content and section pages use.

        "KPMG Global"       pale blue, flush left      x = 72.0
        "private company"   white,     indented        x = 119.3
        "CEO Outlook"       pale blue, indented more   x = 300.8

    THE COLOUR RULE IS POSITIONAL, NOT SEMANTIC, and the previous implementation got this
    wrong in a way that produced a different answer on every deck. It assigned the strong
    colour to the LONGEST line (`max(..., key=len)`), reasoning that the longest line is the
    subject. KPMG assigns by INDEX: tint / white / tint, measured on both the private (75/80/60
    pt) and tmt (50/67/40 pt) covers. A rule that reads the content produces a staircase whose
    emphasis moves when the words change, which is precisely what a signature move must not do.

    Each line is a separate text box so the indent is real geometry rather than leading spaces,
    which would be measured wrong and would break if the line wrapped.

    Returns the y coordinate below the last line.
    """
    if not lines:
        raise ValueError("stepped_headline needs at least one line")

    # Pick the largest size from the SCALE at which EVERY line fits on one line.
    #
    # A stepped headline is only legible as one gesture if each line is a single line; the
    # moment one wraps, the diagonal breaks and the alternating colours stop mapping to
    # phrases. But a panel is narrow (a 42% panel leaves roughly 285pt of measure), so the
    # display size that works on a full-width slide routinely will not fit here. Stepping
    # down through the scale -- never to a continuous size -- keeps sibling dividers
    # consistent while letting a long section title survive.
    def _all_fit(size_pt: float) -> bool:
        lh = points(
            T.line_height_pt(
                ctx.brand.font_major, size_pt, LINE_SPACING.get(role, 1.0), bold=True
            )
        )
        ind = int(box.width * step_in)
        return all(
            T.fit(
                line,
                Box(box.left + i * ind, box.top, box.width - i * ind, int(lh * 1.4)),
                ctx.brand.font_major,
                size_pt,
                bold=True,
                tracking_pt=tracking_pt(role, size_pt),
                line_spacing=LINE_SPACING.get(role, 1.0),
            ).line_count
            == 1
            for i, line in enumerate(lines)
        )

    size = SCALE[role]
    if not _all_fit(size):
        if allow_shrink_to is None:
            raise ValueError(
                f"stepped_headline does not fit at {role} ({size}pt) in "
                f"{box.width / 12700:.0f}pt of measure: {list(lines)}"
            )
        floor = SCALE[allow_shrink_to]
        candidates = [
            x for x in sorted(set(SCALE.values()), reverse=True) if floor <= x <= size
        ]
        chosen = next((x for x in candidates if _all_fit(x)), None)
        if chosen is None:
            raise ValueError(
                f"stepped_headline does not fit even at the {allow_shrink_to} floor "
                f"({floor}pt) in {box.width / 12700:.0f}pt of measure: {list(lines)}\n"
                f"  Shorten the headline, or widen the panel."
            )
        size = chosen

    if colors is None:
        if on_panel:
            tint, strong = ctx.brand.color("on_dark_accent"), ctx.brand.color("on_dark")
        else:
            tint, strong = ctx.brand.color("accent"), ctx.brand.color("headline")
        # POSITIONAL: tint, strong, tint, strong... by line INDEX. Never by line length.
        colors = [strong if i % 2 else tint for i in range(len(lines))]

    line_h = points(
        T.line_height_pt(
            ctx.brand.font_major, size, LINE_SPACING.get(role, 1.0), bold=True
        )
    )
    indent = int(box.width * step_in)
    y = box.top

    for i, line in enumerate(lines):
        left = box.left + i * indent
        # The measuring box gets slack; the ADVANCE stays exactly one line.
        # Sizing it to exactly line_h makes every line fail the fit check by the safety
        # factor alone. Because the text is top-anchored, a taller box does not move the
        # glyphs -- so slack here is free, and y still advances by the true line height.
        lbox = Box(left, y, box.width - i * indent, int(line_h * 1.4))
        fit = T.fit(
            line,
            lbox,
            ctx.brand.font_major,
            size,
            bold=True,
            tracking_pt=tracking_pt(role, size),
            line_spacing=LINE_SPACING.get(role, 1.0),
        )
        if not fit.fits:
            raise ValueError(
                f"stepped_headline line {i} ({line!r}) does not fit at {role} ({size}pt) in "
                f"{fit.box_width_pt:.0f}pt. Shorten the line or split it differently -- each "
                f"line of a stepped headline must occupy exactly one line."
            )
        T.add_textbox(
            slide,
            lbox,
            line,
            family=ctx.brand.font_major,
            size_pt=size,
            bold=True,
            color=colors[i],
            tracking_pt=tracking_pt(role, size),
            line_spacing=LINE_SPACING.get(role, 1.0),
            anchor="top",
        )
        y += line_h

    return y


def two_tone_headline(
    slide,
    ctx: Ctx,
    box: Box,
    lines: Sequence[str],
    *,
    role: str = "h1",
    on_panel: bool = False,
    allow_shrink_to: str | None = "h3",
) -> Placed:
    """
    THE INTERIOR HEADLINE. Opens in KPMG Blue, closes in Cobalt, flush left.

    This is the signature typographic move of the whole system and we did not have it. All five
    section openers in D1 split mid-sentence: the opening line in #00338D, every remaining line
    in #1E49E2, same size, same left edge, NO INDENT. 13 of 14 sampled two-colour headline
    blocks across four other decks pair exactly those two values.

        Economic outlook       #00338D   x = 55.4
        and business           #1E49E2   x = 55.4
        confidence             #1E49E2   x = 55.4

    THE RULE IS POSITIONAL AND INVARIANT: line 1 takes the master blue, every remaining line
    takes Cobalt, and the headline ALWAYS ENDS IN COBALT. The split point moves with the
    sentence; the colours do not move with it. On a panel the same rule runs white -> pale
    blue.

    A headline set entirely in one blue is the single clearest tell that a deck is off this
    system -- which is what ours were, on every page.

    `lines` may be a plain string, and passing one is the better call. THE SPLIT IS THEN
    MEASURED RATHER THAN GUESSED, which is the difference between a component that works at any
    column width and one that works at the width its caller happened to test.

    Guessing looked like `split_headline(title, 3, target_chars=22)` at every call site -- a
    character count standing in for a measurement. It survives exactly until a page has a
    narrower column: the deck's own page 2 has a 293pt headline measure next to a 283pt panel,
    where 22 characters is 26pt of overflow at every size down to the floor, and the component
    could only raise. Given the string, it searches instead: largest size first, fewest lines
    first, and it accepts the first combination whose every line renders on one line.

    Returns a `Placed` measured from the LAST line, so callers advance from real geometry.
    """
    if isinstance(lines, str):
        text_str, lines = lines, None
    else:
        text_str, lines = None, list(lines)
        if not lines:
            raise ValueError("two_tone_headline needs at least one line")

    if on_panel:
        first, rest_color = (
            ctx.brand.color("on_dark"),
            ctx.brand.color("on_dark_accent"),
        )
    else:
        first, rest_color = ctx.brand.color("headline"), ctx.brand.color("accent")

    def _all_fit(candidate: Sequence[str], size_pt: float) -> bool:
        lh = points(
            T.line_height_pt(
                ctx.brand.font_major, size_pt, LINE_SPACING.get(role, 1.0), bold=True
            )
        )
        return all(
            T.fit(
                line,
                Box(box.left, box.top, box.width, int(lh * 1.4)),
                ctx.brand.font_major,
                size_pt,
                bold=True,
                tracking_pt=tracking_pt(role, size_pt),
                line_spacing=LINE_SPACING.get(role, 1.0),
            ).line_count
            == 1
            for line in candidate
        )

    top = SCALE[role]
    floor = SCALE[allow_shrink_to] if allow_shrink_to else top
    sizes = [x for x in sorted(set(SCALE.values()), reverse=True) if floor <= x <= top]

    size = None
    if text_str is not None:
        # Size before line count: a headline one size larger on four lines still reads as the
        # page's claim, while the same headline shrunk to fit on two has stopped being one.
        for candidate_size in sizes:
            per_line = max(6, int((box.width / 12700) / (candidate_size * 0.52)))
            for n in (2, 3, 4, 5):
                candidate = split_headline(text_str, n, target_chars=per_line)
                if _all_fit(candidate, candidate_size):
                    lines, size = candidate, candidate_size
                    break
            if size is not None:
                break
        if size is None:
            raise ValueError(
                f"two_tone_headline cannot set {text_str!r} in "
                f"{box.width / 12700:.0f}pt of measure, at any size down to "
                f"{allow_shrink_to} ({floor}pt) on up to five lines.\n"
                f"  CUT WORDS. A headline this long in a column this narrow is carrying more "
                f"than one idea."
            )
    else:
        size = next((x for x in sizes if _all_fit(lines, x)), None)
        if size is None:
            raise ValueError(
                f"two_tone_headline does not fit even at the {allow_shrink_to} floor "
                f"({floor}pt) in {box.width / 12700:.0f}pt of measure: {list(lines)}\n"
                f"  Pass the headline as a STRING instead of pre-split lines and it will "
                f"find a split that fits."
            )

    colors = [first] + [rest_color] * (len(lines) - 1)

    line_h = points(
        T.line_height_pt(
            ctx.brand.font_major, size, LINE_SPACING.get(role, 1.0), bold=True
        )
    )
    y = box.top
    last = None
    for i, line in enumerate(lines):
        # Flush left, every line. `step_in = 0` is not a parameter here because a stepped
        # interior headline is not a variant of this component -- it is the cover's component.
        lbox = Box(box.left, y, box.width, int(line_h * 1.4))
        last = T.add_textbox(
            slide,
            lbox,
            line,
            family=ctx.brand.font_major,
            size_pt=size,
            bold=True,
            color=colors[i],
            tracking_pt=tracking_pt(role, size),
            line_spacing=LINE_SPACING.get(role, 1.0),
            anchor="top",
        )
        y += line_h

    return Placed(last, box, y - box.top, size)


def split_headline(
    text: str, max_lines: int = 3, target_chars: int | None = None
) -> list[str]:
    """
    Break a headline into stepped lines at natural syntactic joints.

    Priority of break points, highest first:
      1. A SENTENCE boundary. A two-sentence claim steps at the full stop or not at all;
         stepping a mid-sentence fragment reads as a typesetting accident.
      2. A colon, comma or semicolon -- the author's own stated break.
      3. A preposition or conjunction, which is where the qualifier/subject split falls:
         "Tuning the workforce" / "into an AI world".
      4. Failing all of those, the midpoint.

    `target_chars` caps the length of any one line. When a line exceeds it and lines remain
    in the budget, that line is split again at its own best joint. This is what lets a long
    title survive a narrow panel: a 45-character headline that will not fit on two lines at a
    readable size fits comfortably on three.
    """

    def best_break(ws: list[str]) -> int | None:
        if len(ws) < 3:
            return None
        JOINTS = {
            "into",
            "through",
            "with",
            "for",
            "and",
            "to",
            "in",
            "on",
            "at",
            "by",
            "from",
            "then",
            "before",
            "after",
            "across",
            "without",
            "under",
            "over",
        }
        sentence = [
            i + 1 for i, w in enumerate(ws) if w.endswith(".") and 0 < i < len(ws) - 1
        ]
        punct = [
            i + 1
            for i, w in enumerate(ws)
            if w.endswith((":", ",", ";")) and 0 < i < len(ws) - 1
        ]
        joints = [
            i
            for i, w in enumerate(ws)
            if w.lower().strip(",:;.") in JOINTS and 0 < i < len(ws) - 1
        ]
        # The fallback balances by CHARACTERS, not by word index, and refuses to orphan a line.
        #
        # `len(ws) // 2` was the previous fallback and it produced "Set up / a / clean
        # workspace" -- a one-character second line. A word-index midpoint assumes every word is
        # the same width, which is exactly the assumption the rest of this package exists to
        # stop anyone making. The orphan guard is separate and stricter: a display line of
        # under MIN_LINE_CHARS is a rag defect no reader forgives, and it is worth an unbalanced
        # split to avoid.
        MIN_LINE_CHARS = 5
        total = sum(len(w) + 1 for w in ws) - 1
        candidates = []
        for i in range(1, len(ws)):
            left = sum(len(w) + 1 for w in ws[:i]) - 1
            right = total - left - 1
            if left < MIN_LINE_CHARS or right < MIN_LINE_CHARS:
                continue
            candidates.append((abs(left - right), i))
        if not candidates:
            return None
        for pool in (sentence, punct, joints):
            legal = [i for i in pool if any(i == c[1] for c in candidates)]
            if legal:
                return min(legal, key=lambda i: abs(i - len(ws) / 2))
        return min(candidates)[1]

    words = text.split()
    if len(words) <= 2 or max_lines < 2:
        return [text]

    b = best_break(words)
    lines = [" ".join(words[:b]), " ".join(words[b:])] if b else [text]

    # Split the longest line again while there is budget and it is over target.
    while len(lines) < max_lines:
        idx = max(range(len(lines)), key=lambda i: len(lines[i]))
        if target_chars is not None and len(lines[idx]) <= target_chars:
            break
        if target_chars is None and len(lines) >= 2:
            break
        ws = lines[idx].split()
        nb = best_break(ws)
        if not nb:
            break
        lines[idx : idx + 1] = [" ".join(ws[:nb]), " ".join(ws[nb:])]

    return lines


def fits_stepped(
    ctx: Ctx,
    box: Box,
    lines: Sequence[str],
    *,
    role: str = "h1",
    step_in: float = 0.055,
    floor: str = "h3",
) -> bool:
    """
    Would `stepped_headline` succeed here?

    Exists so a caller can CHOOSE the panel treatment rather than attempt it and catch a
    failure. Deciding by try/except would work, but it hides the editorial fact that matters:
    the panel treatment is for SHORT claims. A two-sentence statement belongs on a white
    slide, and that is a content judgment the caller should make explicitly.
    """

    def _all_fit(size_pt: float) -> bool:
        lh = points(
            T.line_height_pt(
                ctx.brand.font_major, size_pt, LINE_SPACING.get(role, 1.0), bold=True
            )
        )
        ind = int(box.width * step_in)
        return all(
            T.fit(
                line,
                Box(box.left + i * ind, box.top, box.width - i * ind, int(lh * 1.4)),
                ctx.brand.font_major,
                size_pt,
                bold=True,
                tracking_pt=tracking_pt(role, size_pt),
                line_spacing=LINE_SPACING.get(role, 1.0),
            ).line_count
            == 1
            for i, line in enumerate(lines)
        )

    lo, hi = SCALE[floor], SCALE[role]
    return any(
        _all_fit(x) for x in sorted(set(SCALE.values()), reverse=True) if lo <= x <= hi
    )


def stat_card(
    slide,
    ctx: Ctx,
    box: Box,
    value: str,
    label: str,
    *,
    caption: str | None = None,
    color: str | None = None,
) -> Box:
    """
    A filled Cobalt card carrying one large figure. Measured on TMT p9.

        rect      #1E49E2, 269.0 x 165.1pt -- EXACTLY one 3-column module wide
        label     #FFFFFF 10.5pt, inset 17.0 left / 14.5 top
        numeral   very large, white
        caption   #FFFFFF 10.5pt
        radius 0, shadow 0, stroke 0 -- verified at 600dpi

    The card being exactly one column module wide is the detail that matters. A filled card at
    an arbitrary width reads as a callout box someone drew; at a module width it reads as part
    of the page's structure, and it lines up with everything else on the page for free.

    Label and caption go at the deck's own 18/16pt rather than KPMG's 10.5, per the stated
    room-legibility deviation. The numeral carries the ratio that makes the card work, so it
    stays large: a figure only fractionally bigger than its label reads as a form field.
    """
    fill = color or ctx.brand.color("accent")
    _rect(slide, box, fill)

    card = Ctx(ctx.brand, Ground.BRAND, ctx.grid)
    inset_x = points(17.0)
    inset_y = points(14.5)
    x = box.left + inset_x
    w = box.width - 2 * inset_x
    y = box.top + inset_y

    T.add_textbox(
        slide,
        Box(x, y, w, points(SCALE["small"] * 1.5)),
        label,
        family=ctx.brand.font_minor,
        size_pt=SCALE["small"],
        color=ctx.brand.color("on_dark"),
        line_spacing=LINE_SPACING["body"],
        anchor="top",
    )
    y += points(SCALE["small"] * 1.5)

    value_size = SCALE["display"]
    T.add_textbox(
        slide,
        Box(x, y, w, points(value_size * 1.25)),
        value,
        family=ctx.brand.font_major,
        size_pt=value_size,
        bold=True,
        color=ctx.brand.color("on_dark"),
        tracking_pt=tracking_pt("display", value_size),
        line_spacing=1.0,
        anchor="top",
    )
    y += points(value_size * 1.18)

    if caption:
        T.add_textbox(
            slide,
            Box(x, y, w, box.bottom - y - inset_y),
            caption,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.brand.color("on_dark"),
            line_spacing=LINE_SPACING["body"],
            anchor="top",
        )
    del card
    return box


def contrast_pair(
    slide,
    ctx: Ctx,
    box: Box,
    left_label: str,
    left_body: str,
    right_label: str,
    right_body: str,
    *,
    emphasis: str = "right",
) -> Box:
    """
    Two things set against each other, separated by a vertical rule. The two-column contrast.

    This is the cheapest of the three-element page's second elements, and the one the course
    content reaches for most: "you type it / the agent types it", "the repository with a trail /
    the one without", "grown, never pruned / written when it broke".

    `emphasis` names which side carries the accent. Exactly one side should -- a contrast where
    both sides are equally coloured has presented two options rather than made an argument,
    which is the same defect `Datum.emphasis` exists to prevent in a chart.

    The separator is a vertical `rule_strong` hairline, not a box around each side. Boxing them
    turns a comparison into two cards, and two cards read as a set rather than as an opposition.
    """
    gap = points(SPACE["lg"])
    half = (box.width - gap) // 2
    left_box = Box(box.left, box.top, half, box.height)
    right_box = Box(box.left + half + gap, box.top, half, box.height)

    column_divider(slide, ctx, box.left + half + gap // 2, box.top, box.bottom)

    strong = (
        ctx.brand.color("on_dark") if ctx.dark_ground else ctx.brand.color("accent")
    )
    quiet = (
        ctx.brand.color("on_dark_muted")
        if ctx.dark_ground
        else ctx.brand.color("headline")
    )

    # BOTH LABELS ARE MEASURED, AND BOTH BODIES START BELOW THE TALLER OF THE TWO.
    #
    # Two separate defects hid in the one line this replaces (`label_h = SCALE["h3"] * 1.5`).
    # A fixed 1.5-line allocation is too short for a two-line label, so the body was drawn
    # straight through it -- and even measured PER SIDE, a one-line label opposite a two-line
    # one puts the two bodies on different baselines, which reads as a rendering fault rather
    # than as a comparison. One measurement, shared: the sibling-consistency rule applied
    # inside a slide.
    # AND THE PAIR FITS ITS BOX, or it steps down until it does.
    #
    # A supporting element is given whatever the page has left, which is not a fixed amount --
    # so a contrast set at one size works on the page with a short body and runs into the
    # footer band on the page with a long one. The structural gate cannot see that: an
    # overflowing text frame is a valid text frame. Stepping down through the scale is the same
    # move `_body_columns` makes, and for the same reason.
    def _measure(label_size: float, body_size: float) -> tuple[int, int]:
        lh = max(
            points(
                T.fit(
                    lbl,
                    Box(0, 0, half, points(500)),
                    ctx.brand.font_major,
                    label_size,
                    bold=True,
                    tracking_pt=tracking_pt("h3", label_size),
                    line_spacing=LINE_SPACING["h3"],
                ).height_pt
            )
            for lbl in (left_label, right_label)
        )
        bh = max(
            points(
                T.fit(
                    bdy,
                    Box(0, 0, half, points(500)),
                    ctx.brand.font_minor,
                    body_size,
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            for bdy in (left_body or " ", right_body or " ")
        )
        return lh, bh

    pairs = [
        (SCALE["h3"], SCALE["small"]),
        (SCALE["lead"], SCALE["small"]),
        (SCALE["lead"], SCALE["caption"]),
        (SCALE["small"], SCALE["caption"]),
    ]
    label_size, body_size = pairs[-1]
    label_h, body_h = _measure(*pairs[-1])
    for cand_label, cand_body in pairs:
        lh, bh = _measure(cand_label, cand_body)
        if lh + points(SPACE["xs"]) + bh <= box.height:
            label_size, body_size, label_h = cand_label, cand_body, lh
            break

    for side, cell, lbl, bdy in (
        ("left", left_box, left_label, left_body),
        ("right", right_box, right_label, right_body),
    ):
        fg = strong if side == emphasis else quiet
        T.add_textbox(
            slide,
            Box(cell.left, cell.top, cell.width, label_h),
            lbl,
            family=ctx.brand.font_major,
            size_pt=label_size,
            bold=True,
            color=fg,
            tracking_pt=tracking_pt("h3", label_size),
            line_spacing=LINE_SPACING["h3"],
            anchor="top",
        )
        if bdy:
            T.add_textbox(
                slide,
                Box(
                    cell.left,
                    cell.top + label_h + points(SPACE["xs"]),
                    cell.width,
                    cell.height - label_h - points(SPACE["xs"]),
                ),
                bdy,
                family=ctx.brand.font_minor,
                size_pt=body_size,
                color=ctx.ink,
                line_spacing=LINE_SPACING["body"],
                anchor="top",
            )
    return box


def code_block(
    slide, ctx: Ctx, box: Box, lines: Sequence[str], *, caption: str | None = None
) -> Box:
    """
    A worked terminal or file-tree example, set monospaced on a light neutral field.

    The one place a non-brand typeface is correct: a folder tree or a shell line set in Arial
    has stopped being an example of a thing you type and become a description of one. The
    ground is #E5E5E5, the real light neutral in this system -- NOT #F5F6FA, which does not
    appear in the top-28 flat colours on any of 95 measured pages.

    COURIER NEW, NOT MENLO, AND THIS IS THE PACKAGE'S OWN RULE APPLIED TO ITSELF. The first
    version named Menlo, which is a macOS face and is not on Windows -- and this deck opens on
    140 attendees' laptops. Naming a font the delivery machine lacks is worse than not using
    one: PowerPoint substitutes silently, the substitute has different advance widths, and
    every box measured against the declared font re-flows on the presenting machine after every
    check has passed on the building one. That is stated in `Deck._resolve_fonts`'s docstring
    and in the brand file, and it was broken here within an hour of being written. Courier New
    ships with both Windows and macOS.

    No frame, no rounded corners, no traffic-light dots. It is a field of colour with type on
    it, which is what a KPMG page does with everything.
    """
    _rect(slide, box, ctx.brand.color("rule"))
    inner = box.inset(points(SPACE["md"]), points(SPACE["sm"]))
    # The caption's height is MEASURED and subtracted. A flat 2.2-line reservation is short by
    # a line whenever the caption wraps, and the code is then drawn straight through it.
    cap_h = (
        0
        if not caption
        else points(
            T.fit(
                caption,
                Box(0, 0, inner.width, points(200)),
                ctx.brand.font_minor,
                SCALE["caption"],
                line_spacing=LINE_SPACING["body"],
            ).height_pt
        )
    )
    body_h = inner.height - cap_h - (points(SPACE["sm"]) if caption else 0)
    # A CODE BLOCK IS THE ONE PLACE A LINE MAY NOT WRAP. A wrapped shell line or file-tree row
    # has stopped being an example of a thing you type. So the size steps DOWN until every
    # line fits the measure on one line -- and Courier is wide, so this bites often. The floor
    # is the footnote size; below that it is not a worked example anyone can read from a room,
    # and the right fix is a shorter line, not smaller type.
    code_size = SCALE["caption"]
    text_block = "\n".join(lines)
    # The floor is `footnote` (14pt), NOT `micro`. Letting it reach 11pt to make a long line
    # fit produced a folder tree nobody past row three could read -- which is the type floor
    # this package sets, given away to save an editor from shortening a path. A code line that
    # will not fit at 14pt is too long, and that is a content finding.
    for candidate in (SCALE["caption"], SCALE["footnote"]):
        f = T.fit(
            text_block,
            Box(0, 0, inner.width, body_h),
            "Courier New",
            candidate,
            line_spacing=1.35,
        )
        if f.fits and f.line_count == len(lines):
            code_size = candidate
            break
        code_size = candidate
    T.add_textbox(
        slide,
        Box(inner.left, inner.top, inner.width, body_h),
        text_block,
        family="Courier New",
        size_pt=code_size,
        color=ctx.brand.color("ink_strong"),
        line_spacing=1.35,
        anchor="top",
    )
    if caption:
        T.add_textbox(
            slide,
            Box(inner.left, inner.bottom - cap_h, inner.width, cap_h),
            caption,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.brand.color("headline"),
            line_spacing=LINE_SPACING["body"],
            anchor="bottom",
        )
    return box


# ---------------------------------------------------------------------------
# Time, flow, and the annotated figure
# ---------------------------------------------------------------------------
#
# The three forms below exist because of a measured gap rather than a wish list. Rendered and
# read page by page, roughly thirty of this package's own 43 reference pages were the same
# shape: eyebrow, two-tone headline, two or three columns of body text, closing line. The
# archetype names differed and the pages did not, which is exactly what `check_siblings` cannot
# see -- it counts names.
#
# What was missing was not decoration. It was every argument that is NOT "here are some points":
# an argument about WHEN (timeline), about ROUTE (flow_diagram), and about ONE FIGURE and what
# hangs off it (annotated_figure). A deck with no time form puts its chronology in prose.


@dataclass(frozen=True)
class Moment:
    """
    One dated point on a `timeline`.

    `when` is the date exactly as it should read on the page -- "Sept 2026", "Day 1", "Q3" --
    and is never parsed. `emphasis` marks the single moment the headline is about; exactly one
    should carry it, for the same reason `Datum.emphasis` exists in charts.py.
    """

    when: str
    what: str
    detail: str | None = None
    emphasis: bool = False


def timeline(
    slide,
    ctx: Ctx,
    box: Box,
    moments: Sequence[Moment],
    *,
    axis: bool = True,
) -> Placed:
    """
    An ordered run of dated moments along one horizontal axis, with one moment marked.

    THE POSITIONS ARE ORDINAL, NOT PROPORTIONAL. Moments are spaced evenly whatever their dates
    say, so an argument that depends on the GAPS -- "nothing moved for three years and then
    everything did" -- is misrepresented by this component and wants a real time axis instead.
    What this draws is: these things happened, in this order, and that one is the point. Say so
    in the headline and the ordinal spacing is honest; imply elapsed time and it is not.

    Date above the axis, marker on it, moment below. The emphasised marker is a filled accent
    disc at twice the diameter and its label goes to headline weight; the rest sit in
    rule_strong. That is the same ghosting the charts use, and it is what makes a timeline
    argue rather than list.

    The axis runs the full width rather than stopping at the last marker, because a line that
    stops dead on the final date says the sequence ends there, which is almost never the claim.

    Returns the measured bottom of the deepest column, so a caller advances from real geometry.
    """
    if not 2 <= len(moments) <= 6:
        raise ValueError(
            f"timeline() takes 2 to 6 moments, got {len(moments)}. Beyond six the labels are "
            f"narrower than the words in them and the form has become a table."
        )

    n = len(moments)
    gap = points(SPACE["md"])
    marker_d = points(14.0)
    cells = [box.row(n, i, gap=gap) for i in range(n)]

    date_h = max(
        points(
            T.fit(
                m.when,
                Box(0, 0, cells[0].width, points(200)),
                ctx.brand.font_minor,
                SCALE["caption"],
                bold=True,
                tracking_pt=SCALE["caption"] * 0.08,
                line_spacing=LINE_SPACING["footnote"],
            ).height_pt
        )
        for m in moments
    )
    axis_y = box.top + date_h + points(SPACE["xs"]) + marker_d // 2
    label_top = axis_y + marker_d // 2 + points(SPACE["sm"])
    room = box.bottom - label_top

    # ONE MEASUREMENT, SHARED BY EVERY COLUMN. Measured per column, a one-line label beside a
    # two-line one puts the two details on different baselines, which reads as a rendering
    # fault rather than as a sequence -- the same defect `contrast_pair` records.
    def _measure(label_size: float, detail_size: float) -> tuple[int, int]:
        lh = max(
            points(
                T.fit(
                    m.what,
                    Box(0, 0, cells[0].width, points(500)),
                    ctx.brand.font_major,
                    label_size,
                    bold=True,
                    tracking_pt=tracking_pt("h3", label_size),
                    line_spacing=LINE_SPACING["h3"],
                ).height_pt
            )
            for m in moments
        )
        dh = max(
            points(
                T.fit(
                    m.detail or " ",
                    Box(0, 0, cells[0].width, points(500)),
                    ctx.brand.font_minor,
                    detail_size,
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            for m in moments
        )
        return lh, dh

    pairs = (
        (SCALE["h3"], SCALE["small"]),
        (SCALE["lead"], SCALE["small"]),
        (SCALE["lead"], SCALE["caption"]),
        (SCALE["small"], SCALE["caption"]),
    )
    label_size, detail_size = pairs[-1]
    label_h, detail_h = _measure(*pairs[-1])
    for cand_label, cand_detail in pairs:
        lh, dh = _measure(cand_label, cand_detail)
        if lh + points(SPACE["xs"]) + dh <= room:
            label_size, detail_size, label_h, detail_h = cand_label, cand_detail, lh, dh
            break

    if axis:
        _rect(
            slide,
            Box(box.left, axis_y - points(0.375), box.width, points(0.75)),
            ctx.rule,
        )

    last_shape = None
    for cell, m in zip(cells, moments):
        # THE LABEL DE-EMPHASISES BY HUE, NOT BY GOING GREY. Set in `muted`, the moment names
        # rendered #666666 while their own supporting detail below them was brand navy -- the
        # subordinate line louder than the idea it supports, on every column. Ghosting works in
        # a chart because a grey bar is still a bar; a grey headline has just been demoted.
        fg = ctx.accent if m.emphasis else ctx.headline
        T.add_textbox(
            slide,
            Box(cell.left, box.top, cell.width, date_h),
            m.when,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            bold=True,
            caps=True,
            color=ctx.accent if m.emphasis else ctx.muted,
            tracking_pt=SCALE["caption"] * 0.08,
            line_spacing=LINE_SPACING["footnote"],
        )

        # The marker is LEFT-ALIGNED on the cell, not centred in it, so its edge and the left
        # edge of the label beneath it are the same vertical -- which is what makes the column
        # read as one object rather than as a dot with some text near it.
        d = marker_d if m.emphasis else (marker_d * 2) // 3
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            *Box(cell.left, axis_y - d // 2, d, d).as_tuple(),
        )
        X.set_fill(dot, ctx.accent if m.emphasis else ctx.rule_strong)
        X.set_line(dot, None)
        X.clear_effects(dot)
        dot.shadow.inherit = False

        last_shape = T.add_textbox(
            slide,
            Box(cell.left, label_top, cell.width, label_h),
            m.what,
            family=ctx.brand.font_major,
            size_pt=label_size,
            bold=True,
            color=fg,
            tracking_pt=tracking_pt("h3", label_size),
            line_spacing=LINE_SPACING["h3"],
            anchor="top",
        )
        if m.detail:
            T.add_textbox(
                slide,
                Box(
                    cell.left,
                    label_top + label_h + points(SPACE["xs"]),
                    cell.width,
                    detail_h,
                ),
                m.detail,
                family=ctx.brand.font_minor,
                size_pt=detail_size,
                color=ctx.ink,
                line_spacing=LINE_SPACING["body"],
                anchor="top",
            )

    used = label_top + label_h + points(SPACE["xs"]) + detail_h - box.top
    return Placed(last_shape, box, used, label_size)


def flow_diagram(
    slide,
    ctx: Ctx,
    box: Box,
    nodes: Sequence[str],
    *,
    emphasis: int | None = None,
    height_pt: float = 96.0,
) -> Box:
    """
    Three to five named stages joined left to right, with the route into one of them marked.

    THE DIFFERENCE FROM `steps()` IS THE ARGUMENT, NOT THE PICTURE. `steps()` says "do these in
    order"; this says "the work travels this way, and it arrives HERE". So the emphasised node
    takes the accent field AND every connector leading into it does too -- the route lights up,
    the tail beyond it stays quiet. Emphasise nothing and it is a plain process diagram, which
    is a weaker exhibit and usually means `steps()` was the right component.

    A node is a flat field of colour with type on it, which is what this system does with
    everything: no outline, no radius, no shadow. The connector is a hairline with one small
    solid chevron at the end it points to -- an arrowHEAD, not an arrow, because the line is
    already doing the pointing and a two-ended arrow on every join is chart-library furniture.

    Do NOT use it for anything that branches. Five nodes is the cap and a single row is the only
    layout, so a diagram with a fork has to be drawn as two flows or rethought as a `matrix_2x2`.
    """
    if not 3 <= len(nodes) <= 5:
        raise ValueError(
            f"flow_diagram() takes 3 to 5 nodes, got {len(nodes)}. Two nodes is a "
            f"`contrast_pair`; six is a diagram nobody reads from a room."
        )
    if emphasis is not None and not 0 <= emphasis < len(nodes):
        raise ValueError(f"emphasis {emphasis} is not one of {len(nodes)} nodes")

    n = len(nodes)
    join = points(SPACE["lg"])
    node_w = (box.width - join * (n - 1)) // n
    node_h = points(height_pt)
    mid_y = box.top + node_h // 2
    # 13pt, not 11. At 11 the chevron rendered as a 17px triangle at projector scale, which is
    # a speck rather than a direction -- and the connector it terminates is only 1.25pt itself.
    chevron = points(13.0)

    quiet_fill = (
        ctx.brand.color("rule")
        if not ctx.dark_ground
        else ctx.brand.color("on_dark_muted")
    )
    live_fill = ctx.accent

    # THE LABEL IS MEASURED AGAINST THE NODE'S OWN FILL, NOT AGAINST THE SLIDE'S GROUND, and
    # `ctx.assert_legible` is deliberately NOT used here for exactly that reason -- it compares
    # against `ctx.canvas`, so on a white slide it reads white-on-white and rejects the one
    # correct answer. A node is a field, so the ground rule applies one level down.
    def _label_color(fill: str, size_pt: float) -> str:
        dark_ink = ctx.brand.color("ink_strong")
        light_ink = ctx.brand.color("on_dark")
        fg = (
            light_ink
            if contrast_ratio(light_ink, fill) >= contrast_ratio(dark_ink, fill)
            else dark_ink
        )
        verdict, ratio = contrast_verdict(fg, fill, size_pt, True)
        if verdict == "FAIL":
            raise ValueError(
                f"a flow node filled {fill} cannot carry a legible label: the better of "
                f"{light_ink} and {dark_ink} is only {ratio:.2f}:1 on it."
            )
        return fg

    # Measured once across every node, so five labels of different lengths still set at one
    # size. Sized per node they would step down independently and the row would read ragged.
    inner_w = node_w - 2 * points(SPACE["sm"])
    label_size = SCALE["small"]
    for candidate in (SCALE["h3"], SCALE["lead"], SCALE["body"], SCALE["small"]):
        if all(
            T.fit(
                label,
                Box(0, 0, inner_w, node_h - 2 * points(SPACE["xs"])),
                ctx.brand.font_major,
                candidate,
                bold=True,
                tracking_pt=tracking_pt("h3", candidate),
                line_spacing=LINE_SPACING["h3"],
            ).fits
            for label in nodes
        ):
            label_size = candidate
            break

    for i, label in enumerate(nodes):
        left = box.left + i * (node_w + join)
        live = emphasis is not None and i == emphasis
        fill = live_fill if live else quiet_fill
        _rect(slide, Box(left, box.top, node_w, node_h), fill)
        fg = _label_color(fill, label_size)
        T.add_textbox(
            slide,
            Box(
                left + points(SPACE["sm"]),
                box.top + points(SPACE["xs"]),
                inner_w,
                node_h - 2 * points(SPACE["xs"]),
            ),
            label,
            family=ctx.brand.font_major,
            size_pt=label_size,
            bold=True,
            color=fg,
            align="center",
            anchor="middle",
            tracking_pt=tracking_pt("h3", label_size),
            line_spacing=LINE_SPACING["h3"],
        )

        if i == n - 1:
            continue
        # Connector i joins node i to node i+1, so it is part of the route into the emphasised
        # node exactly when i + 1 <= emphasis.
        on_route = emphasis is not None and i + 1 <= emphasis
        stroke = ctx.accent if on_route else ctx.rule_strong
        _rect(
            slide,
            Box(
                left + node_w,
                mid_y - points(0.625),
                join - chevron,
                points(1.25),
            ),
            stroke,
        )
        head = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            *Box(
                left + node_w + join - chevron, mid_y - chevron // 2, chevron, chevron
            ).as_tuple(),
        )
        head.rotation = 90.0
        X.set_fill(head, stroke)
        X.set_line(head, None)
        X.clear_effects(head)
        head.shadow.inherit = False

    return Box(box.left, box.top, box.width, node_h)


def annotated_figure(
    slide,
    ctx: Ctx,
    box: Box,
    figure: str,
    annotations: Sequence[tuple[float, str]],
    *,
    label: str | None = None,
    color: str | None = None,
) -> Box:
    """
    One figure set large, with two or three elbow leaders taking annotations off it.

    The form a page wants when the number is not the whole story -- "57 minutes" is a fact, and
    "57 minutes, of which 40 are the agent working unattended" is an argument. A `stat` states
    the fact; this one takes it apart.

    Each annotation is `(anchor, text)`, where anchor is 0.0 at the TOP of the figure and 1.0 at
    its bottom: the height on the figure the leader departs from. ANCHORS MUST BE GIVEN IN
    ASCENDING ORDER and the function raises if they are not, because that is what makes crossing
    leaders structurally impossible rather than something to notice in a render. A crossed leader
    is the defect this component is most likely to produce and no checker sees it.

    The leader is an elbow -- out, across, in -- in the hairline grey, with a small solid dot
    where it meets the figure. Not a diagonal: a diagonal rule is the only slanted mark on an
    otherwise orthogonal page and it reads as a callout someone drew in.

    LEADERS DEPART FROM THE FIGURE'S MEASURED INK, NOT FROM ITS COLUMN. That distinction is the
    whole difference between this reading as annotation and reading as two unrelated halves: set
    to leave from the column edge, three dots floated in white space a full inch clear of the
    word they claimed to point at, and the elbows collapsed into a little ladder because they
    had almost no horizontal run to make. "57 min" is narrower than the region it was given, so
    the region is the wrong thing to measure.

    Two or three annotations. One is a caption and wants `stat(caption=...)`; four is a list and
    wants `rows()`.
    """
    if not 2 <= len(annotations) <= 3:
        raise ValueError(
            f"annotated_figure() takes 2 or 3 annotations, got {len(annotations)}. One is a "
            f"caption -- use stat(caption=...); four is a list -- use rows()."
        )
    anchors = [a for a, _ in annotations]
    if any(not 0.0 <= a <= 1.0 for a in anchors):
        raise ValueError(
            "annotation anchors are fractions of the figure height, 0.0-1.0"
        )
    if anchors != sorted(anchors):
        raise ValueError(
            f"annotation anchors {anchors} are out of order. They must ascend, because the "
            f"annotations are laid out top to bottom and ascending anchors are the only thing "
            f"stopping two leaders from crossing."
        )

    fg = color or ctx.accent
    gutter = points(SPACE["xl"])
    fig_w = int(box.width * 0.44)
    note_left = box.left + fig_w + gutter
    note_w = box.right - note_left

    # The figure steps down the scale until it sets on ONE line. A figure that wraps has stopped
    # being a figure and become a paragraph in large type -- the same rule `stat()` records --
    # and here it would also put the leader anchors on the wrong geometry.
    size = SCALE["display"]
    for candidate in (SCALE["mega"], SCALE["hero"], SCALE["display"], SCALE["h1"]):
        f = T.fit(
            figure,
            Box(0, 0, fig_w, points(candidate * 3)),
            ctx.brand.font_major,
            candidate,
            bold=True,
            tracking_pt=tracking_pt("display", candidate),
            line_spacing=1.0,
        )
        if f.line_count == 1 and f.width_pt <= fig_w / 12700.0:
            size = candidate
            break
        size = candidate
    ctx.assert_legible(fg, size, bold=True, what="annotated figure")

    fig_h = points(
        T.fit(
            figure,
            Box(0, 0, fig_w, points(size * 3)),
            ctx.brand.font_major,
            size,
            bold=True,
            tracking_pt=tracking_pt("display", size),
            line_spacing=1.0,
        ).height_pt
    )
    depart_x = min(
        box.left + fig_w,
        box.left
        + T.text_width_emu(
            figure,
            ctx.brand.font_major,
            size,
            bold=True,
            tracking_pt=tracking_pt("display", size),
        )
        # Clear of the last glyph, not on it. Set flush to the measured width the departure dot
        # sat inside the terminal "n" of "57 min", which reads as a printing fault.
        + points(SPACE["xs"]),
    )
    elbow_x = depart_x + (note_left - depart_x) // 2
    T.add_textbox(
        slide,
        Box(box.left, box.top, fig_w, fig_h),
        figure,
        family=ctx.brand.font_major,
        size_pt=size,
        bold=True,
        color=fg,
        tracking_pt=tracking_pt("display", size),
        line_spacing=1.0,
        anchor="top",
    )
    label_h = 0
    if label:
        label_h = points(
            T.fit(
                label,
                Box(0, 0, fig_w, points(300)),
                ctx.brand.font_minor,
                SCALE["body"],
                line_spacing=LINE_SPACING["body"],
            ).height_pt
        )
        T.add_textbox(
            slide,
            Box(box.left, box.top + fig_h + points(SPACE["xs"]), fig_w, label_h),
            label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["body"],
            color=ctx.ink,
            line_spacing=LINE_SPACING["body"],
            anchor="top",
        )

    m = len(annotations)
    slot_h = box.height // m
    hair = points(0.75)
    dot_d = points(7.0)
    leader = ctx.rule_strong

    for i, (anchor, text_str) in enumerate(annotations):
        note_top = box.top + i * slot_h
        note_fit = T.fit(
            text_str,
            Box(0, 0, note_w, points(500)),
            ctx.brand.font_minor,
            SCALE["small"],
            line_spacing=LINE_SPACING["body"],
        )
        note_h = points(note_fit.height_pt)
        note_y = note_top + (slot_h - note_h) // 2
        T.add_textbox(
            slide,
            Box(note_left, note_y, note_w, note_h),
            text_str,
            family=ctx.brand.font_minor,
            size_pt=SCALE["small"],
            color=ctx.ink,
            line_spacing=LINE_SPACING["body"],
            anchor="top",
        )

        # The leader departs from the figure's own measured height at the anchor, elbows once,
        # and arrives at the FIRST LINE of the annotation rather than at its middle -- an
        # annotation that is three lines deep is read from its top, so a leader into its centre
        # points at the second sentence.
        from_y = box.top + int(fig_h * anchor)
        to_y = note_y + points(SCALE["small"] * 0.6)
        _rect(
            slide,
            Box(depart_x, from_y - hair // 2, elbow_x - depart_x, hair),
            leader,
        )
        y0, y1 = (from_y, to_y) if from_y <= to_y else (to_y, from_y)
        _rect(slide, Box(elbow_x, y0, hair, max(hair, y1 - y0)), leader)
        _rect(
            slide,
            Box(elbow_x, to_y - hair // 2, note_left - elbow_x, hair),
            leader,
        )
        cap = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            *Box(depart_x - dot_d // 2, from_y - dot_d // 2, dot_d, dot_d).as_tuple(),
        )
        X.set_fill(cap, fg)
        X.set_line(cap, None)
        X.clear_effects(cap)
        cap.shadow.inherit = False

    return Box(
        box.left,
        box.top,
        box.width,
        max(fig_h + (points(SPACE["xs"]) + label_h if label else 0), box.height),
    )
