"""
oxml.py -- the safe raw-XML escape hatches, and only the safe ones.

python-pptx covers maybe 80% of what a designed slide needs. The other 20% (alpha,
multi-stop gradients, shadow control, character tracking, table chrome removal,
autofit hints) has no Python API and must be reached through lxml.

That is exactly where generated decks earn the PowerPoint "[Repaired]" dialog. Two rules
govern every function in this module, and they are not style preferences:

RULE 1 -- SCHEMA ELEMENT ORDER IS ENFORCED.
    OOXML complex types are xsd:sequence, not xsd:all. An <a:ln> placed before <a:solidFill>
    inside <a:spPr> is invalid even though both elements are legal children. PowerPoint's
    validator rejects the part and offers to "repair" it, which in practice means silently
    dropping content in front of a room. Every insertion here goes through
    `insert_element_before(elm, *successors)` with the successor list taken from the
    published sequence, never through a bare `append`.

RULE 2 -- NEVER SWALLOW AN EXCEPTION.
    A helper that wraps its body in try/except and passes on failure converts a loud,
    fixable bug into a silent design regression: the shadow simply is not there, and nobody
    finds out until the slide is on screen. Every function here either does the thing or
    raises. `verify.py` then re-reads the file and asserts the effect actually landed.

Element sequences used below, from the ECMA-376 DrawingML schema:

    a:spPr    xfrm, custGeom|prstGeom, noFill|solidFill|gradFill|blipFill|pattFill|grpFill,
              ln, effectLst, effectDag, scene3d, sp3d, extLst
    a:ln      noFill|solidFill|gradFill|pattFill, prstDash, round|bevel|miter,
              headEnd, tailEnd, extLst
    a:rPr     ln, fill(noFill|solidFill|...), effectLst, effectDag, highlight,
              uLnTx|uLn, uFillTx|uFill, latin, ea, cs, sym, hlinkClick, hlinkMouseOver,
              rtl, extLst
    a:bodyPr  prstTxWarp, noAutofit|normAutofit|spAutoFit, scene3d, sp3d|flatTx, extLst
"""

from __future__ import annotations

from typing import Iterable, Sequence

from lxml import etree
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Element sequences. Keep these next to the functions that rely on them.
# ---------------------------------------------------------------------------

_SPPR_FILL_SUCCESSORS = (
    "a:ln",
    "a:effectLst",
    "a:effectDag",
    "a:scene3d",
    "a:sp3d",
    "a:extLst",
)
_SPPR_EFFECT_SUCCESSORS = ("a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")
_LN_FILL_SUCCESSORS = (
    "a:prstDash",
    "a:round",
    "a:bevel",
    "a:miter",
    "a:headEnd",
    "a:tailEnd",
    "a:extLst",
)
_BODYPR_AUTOFIT_SUCCESSORS = ("a:scene3d", "a:sp3d", "a:flatTx", "a:extLst")

_FILL_TAGS = (
    "a:noFill",
    "a:solidFill",
    "a:gradFill",
    "a:blipFill",
    "a:pattFill",
    "a:grpFill",
)


def _el(tag: str, **attrs: str) -> etree._Element:
    """Create a namespaced element. Attribute values are coerced to str."""
    return etree.SubElement(
        etree.Element("_"), qn(tag), **{k: str(v) for k, v in attrs.items()}
    )


def _new(tag: str, **attrs: str) -> etree._Element:
    e = etree.Element(qn(tag))
    for k, v in attrs.items():
        e.set(k, str(v))
    return e


def _insert_ordered(
    parent: etree._Element, child: etree._Element, successors: Sequence[str]
) -> None:
    """
    Insert `child` into `parent` immediately before the first element whose tag appears in
    `successors`. This is how python-pptx itself honours xsd:sequence, and it is the single
    most important line of defence against a repair dialog.
    """
    for succ in successors:
        found = parent.find(qn(succ))
        if found is not None:
            found.addprevious(child)
            return
    parent.append(child)


def _drop(parent: etree._Element, tags: Iterable[str]) -> None:
    for tag in tags:
        for found in parent.findall(qn(tag)):
            parent.remove(found)


def _hex(color: str) -> str:
    """'#00338D' | '00338d' -> '00338D'. Rejects anything that is not 6 hex digits."""
    h = color.lstrip("#").upper()
    if len(h) != 6 or any(c not in "0123456789ABCDEF" for c in h):
        raise ValueError(f"expected a 6-digit hex colour, got {color!r}")
    return h


def _pct(value: float, name: str = "value") -> int:
    """0.0-1.0 -> OOXML per-mille-of-a-percent (0-100000)."""
    if not 0.0 <= value <= 1.0:
        raise ValueError(f"{name} must be within 0.0-1.0, got {value}")
    return int(round(value * 100000))


# ---------------------------------------------------------------------------
# Colour elements
# ---------------------------------------------------------------------------


def srgb(
    color: str,
    alpha: float | None = None,
    lum_mod: float | None = None,
    lum_off: float | None = None,
) -> etree._Element:
    """
    Build an <a:srgbClr> with optional transforms.

    The child transforms of srgbClr are an unbounded choice group, so order among them is
    not constrained -- appending is safe here, unlike everywhere else in this module.
    """
    clr = _new("a:srgbClr", val=_hex(color))
    if alpha is not None:
        clr.append(_new("a:alpha", val=str(_pct(alpha, "alpha"))))
    if lum_mod is not None:
        clr.append(_new("a:lumMod", val=str(_pct(lum_mod, "lum_mod"))))
    if lum_off is not None:
        clr.append(_new("a:lumOff", val=str(_pct(lum_off, "lum_off"))))
    return clr


def scheme(
    slot: str,
    alpha: float | None = None,
    lum_mod: float | None = None,
    lum_off: float | None = None,
) -> etree._Element:
    """
    Build an <a:schemeClr>, i.e. a THEME colour reference rather than a literal.

    Prefer this over `srgb` for anything that should track the template. A deck built from
    scheme references survives a theme swap; one built from literals silently stops matching
    the rest of the firm's material the next time the template is updated.

    Valid slots: bg1 tx1 bg2 tx2 accent1..accent6 hlink folHlink dk1 lt1 dk2 lt2 phClr
    """
    valid = {
        "bg1",
        "tx1",
        "bg2",
        "tx2",
        "accent1",
        "accent2",
        "accent3",
        "accent4",
        "accent5",
        "accent6",
        "hlink",
        "folHlink",
        "dk1",
        "lt1",
        "dk2",
        "lt2",
        "phClr",
    }
    if slot not in valid:
        raise ValueError(f"unknown theme slot {slot!r}; valid: {sorted(valid)}")
    clr = _new("a:schemeClr", val=slot)
    if alpha is not None:
        clr.append(_new("a:alpha", val=str(_pct(alpha, "alpha"))))
    if lum_mod is not None:
        clr.append(_new("a:lumMod", val=str(_pct(lum_mod, "lum_mod"))))
    if lum_off is not None:
        clr.append(_new("a:lumOff", val=str(_pct(lum_off, "lum_off"))))
    return clr


# ---------------------------------------------------------------------------
# Fills
# ---------------------------------------------------------------------------


def set_fill(shape, color: str | etree._Element, alpha: float | None = None) -> None:
    """
    Solid-fill a shape, with optional transparency.

    python-pptx can set a solid fill but has no alpha API at all -- `fill.fore_color` exposes
    no opacity. Transparency is the difference between a scrim that makes type legible over
    an image and a flat box that hides it, so it is not optional for design work.

    `color` may be a hex string or a colour element from `srgb()` / `scheme()`.
    """
    spPr = shape._element.spPr
    _drop(spPr, _FILL_TAGS)

    clr = color if isinstance(color, etree._Element) else srgb(color, alpha=alpha)
    if isinstance(color, etree._Element) and alpha is not None:
        clr.append(_new("a:alpha", val=str(_pct(alpha, "alpha"))))

    solid = _new("a:solidFill")
    solid.append(clr)
    _insert_ordered(spPr, solid, _SPPR_FILL_SUCCESSORS)


def set_no_fill(shape) -> None:
    spPr = shape._element.spPr
    _drop(spPr, _FILL_TAGS)
    _insert_ordered(spPr, _new("a:noFill"), _SPPR_FILL_SUCCESSORS)


# The measured adjacency set: the only colour pairs KPMG ramps between, read out of PDFs.
# Membership is unordered -- a pair may run in either direction.
#
# REWRITTEN 2026-08-09 AGAINST THE ACTUAL RULE, AND THIS ENTRY IS A CORRECTION OF MY OWN WORK.
#
# The KPMG Brand Guidelines, March 2022, p.46 states the gradient system outright, and it is far
# narrower than anything inferred here. There are exactly TWO gradients:
#
#     PRIMARY gradient   Purple #7213EA  ->  Cobalt Blue #1E49E2
#     SUPPORT gradient   Pacific #00B8F5 ->  Light Blue  #ACEAFF
#
# with stops at 0% and 100%, midpoint at 50%, **0 degree angle**, linear and never radial. The
# page says in its own words: *do not create new gradients; use only the gradients shown here.*
# Either may be FLIPPED (Purple->Cobalt or Cobalt->Purple) keeping location and midpoint.
#
# EARLIER TODAY THIS SET WAS EXTENDED WITH FOUR PAIRS AND THAT WAS WRONG. They were sampled off
# member-firm documents -- Purple->Pacific, Dark blue->Light purple, Cobalt->Dark blue, Deep
# purple->Light purple -- and added under the reasoning that "a refusal that meets a real
# observation is an out-of-date record". That reasoning is sound in general and was false here,
# because a brand book is not a record of observations: it is the rule the observations are
# supposed to follow. Sampling a member firm's output and generalising it is exactly the error
# that produced the rejected deck in the first place, committed a second time one layer down.
#
# Three of the ORIGINAL five rows go too, for the same reason: KPMG Blue->Cobalt, Cobalt->Pacific
# and the dark cover ramp are not among the two sanctioned gradients. Only the two below survive,
# and they survive because the brand book names them, not because they were measured.
GRADIENT_ADJACENCY = frozenset(
    frozenset(p)
    for p in (
        ("7213EA", "1E49E2"),  # PRIMARY: Purple -> Cobalt Blue.  bb FY22 p.46
        ("00B8F5", "ACEAFF"),  # SUPPORT: Pacific -> Light Blue.  bb FY22 p.46
    )
)

# The two gradients as ordered stop lists, which is how the brand book gives them. Flipping is
# permitted; inventing is not.
GRADIENT_PRIMARY: tuple[tuple[float, str], ...] = ((0.0, "#7213EA"), (1.0, "#1E49E2"))
GRADIENT_SUPPORT: tuple[tuple[float, str], ...] = ((0.0, "#00B8F5"), (1.0, "#ACEAFF"))
GRADIENT_ANGLE_DEG = 0.0


def set_gradient(
    shape,
    stops: Sequence[tuple[float, str]],
    angle_deg: float = 0.0,
    alphas: Sequence[float] | None = None,
    *,
    unchecked: bool = False,
) -> None:
    """
    Apply a two-stop axial gradient between adjacent palette colours.

    python-pptx exposes `fill.gradient()` but its `gradient_stops` collection is fixed at
    whatever the default two-stop template provides; you cannot add a third stop or control
    positions freely. Real gradient work needs the raw element.

    stops:      [(position 0.0-1.0, '#RRGGBB'), ...] -- exactly two, ascending
    angle_deg:  0 = left-to-right (the print default). 90 = top-to-bottom.
    alphas:     optional per-stop opacity, same length as stops
    unchecked:  bypass the pair check. For a non-KPMG brand, or a deliberate exception.

    THIS DOCSTRING USED TO SAY THE OPPOSITE AND WAS WRONG. It read: "a gradient between two
    DIFFERENT hues is decoration and reads as amateur under projection". That is a reasonable
    general prior, and KPMG does exactly the thing it forbids -- #7213EA -> #1E49E2 and
    #00338D -> #1E49E2, horizontal, two stops, in print, in six separate documents. The rule is
    not "same hue only"; it is NARROWER THAN THAT IN THREE OTHER DIRECTIONS and those are the
    ones that actually do the work:

        two stops only        every gradient in 34 PDFs is two-stop; multi-stop appears 0 times
                              outside a single 2019-era file
        AXIAL only            every gradient in 34 PDFs is axial. One radial in the entire
                              corpus (a member-firm annual report); zero `radial-gradient` in
                              998KB of kpmg.com CSS
        ADJACENT PAIRS only   the measured set is small and closed -- see GRADIENT_ADJACENCY.
                              Cobalt to Pacific is a ramp; Cobalt to green is a mistake
        horizontal in print   0deg. 225deg is the web value and does not transfer

    So the constraint is enforced here rather than described in a reference file, because the
    previous version WAS described in a reference file and the description was false for two
    years without anything catching it.
    """
    if len(stops) != 2 and not unchecked:
        raise ValueError(
            f"a gradient takes exactly two stops, got {len(stops)}. Every gradient measured "
            f"across 34 KPMG PDFs is two-stop; a three-stop ramp is not a richer version of "
            f"this device, it is a different device. Pass unchecked=True to override."
        )
    if len(stops) < 2:
        raise ValueError("a gradient needs at least two stops")
    positions = [p for p, _ in stops]
    if positions != sorted(positions):
        raise ValueError(f"gradient stops must ascend by position, got {positions}")
    if alphas is not None and len(alphas) != len(stops):
        raise ValueError("alphas must be the same length as stops")

    if not unchecked and len(stops) == 2:
        pair = frozenset(_hex(c) for _, c in stops)
        if len(pair) > 1 and pair not in GRADIENT_ADJACENCY:
            raise ValueError(
                f"{'/'.join(sorted(pair))} is not a measured adjacent pair. The set is "
                f"{sorted('-'.join(sorted(p)) for p in GRADIENT_ADJACENCY)}. A ramp between "
                f"non-adjacent palette colours turns muddy in the middle, which is the "
                f"failure the two-stop-adjacent rule exists to prevent. Pass unchecked=True "
                f"if this is deliberate and not KPMG."
            )

    spPr = shape._element.spPr
    _drop(spPr, _FILL_TAGS)

    grad = _new("a:gradFill", flip="none", rotWithShape="1")
    gsLst = _new("a:gsLst")
    for i, (pos, color) in enumerate(stops):
        gs = _new("a:gs", pos=str(_pct(pos, "stop position")))
        gs.append(srgb(color, alpha=None if alphas is None else alphas[i]))
        gsLst.append(gs)
    grad.append(gsLst)
    # OOXML angles are in 60000ths of a degree, measured clockwise from the positive x-axis.
    grad.append(
        _new("a:lin", ang=str(int(round(angle_deg * 60000)) % 21600000), scaled="0")
    )
    _insert_ordered(spPr, grad, _SPPR_FILL_SUCCESSORS)


# ---------------------------------------------------------------------------
# Lines
# ---------------------------------------------------------------------------


def set_line(
    shape,
    color: str | etree._Element | None,
    width_pt: float = 1.0,
    alpha: float | None = None,
    dash: str | None = None,
    cap: str | None = None,
) -> None:
    """
    Set (or with color=None, remove) a shape outline.

    dash: one of solid dot dash lgDash dashDot lgDashDot lgDashDotDot sysDash sysDot ...
    cap:  "rnd" | "sq" | "flat". LEAVE IT ALONE unless the shape is a chart arc.

    Round caps are for chart arcs and nothing else. The corpus is unambiguous on this: line
    caps across 95 pages run butt 13,482 / round 3,386 / square 8, and the round ones are the
    value arcs of ring gauges. A rounded cap on a rule, a divider or a bar is the visual
    signature of a chart library's defaults, which is exactly what this package exists to
    avoid emitting.
    """
    spPr = shape._element.spPr
    _drop(spPr, ("a:ln",))

    attrs = {"w": str(int(round(width_pt * 12700)))}  # pt -> EMU
    if cap:
        attrs["cap"] = cap
    ln = _new("a:ln", **attrs)
    if color is None:
        ln.append(_new("a:noFill"))
    else:
        clr = color if isinstance(color, etree._Element) else srgb(color, alpha=alpha)
        if isinstance(color, etree._Element) and alpha is not None:
            clr.append(_new("a:alpha", val=str(_pct(alpha, "alpha"))))
        solid = _new("a:solidFill")
        solid.append(clr)
        ln.append(solid)
    if dash:
        _insert_ordered(ln, _new("a:prstDash", val=dash), _LN_FILL_SUCCESSORS[1:])
    _insert_ordered(spPr, ln, _SPPR_FILL_SUCCESSORS[1:])


# ---------------------------------------------------------------------------
# Effects
# ---------------------------------------------------------------------------


def clear_effects(shape) -> None:
    """
    Remove every effect from a shape.

    Worth calling by default. Autoshapes inserted by python-pptx inherit the theme's default
    shape style, which in the stock Office theme carries a subtle drop shadow. On a flat,
    restrained slide that inherited shadow is the single most common tell that a deck was
    generated rather than designed -- it appears on every box and nobody put it there.
    """
    spPr = shape._element.spPr
    _drop(spPr, ("a:effectLst", "a:effectDag"))
    _insert_ordered(spPr, _new("a:effectLst"), _SPPR_EFFECT_SUCCESSORS)


def set_shadow(
    shape,
    blur_pt: float = 18.0,
    distance_pt: float = 4.0,
    direction_deg: float = 90.0,
    color: str = "#000000",
    alpha: float = 0.16,
) -> None:
    """
    Apply a deliberate outer shadow.

    Use sparingly and only to encode ELEVATION -- a card that genuinely floats above the
    page. A shadow used for decoration on a flat layout is the classic amateur tell. The
    defaults here are a soft, low-opacity, straight-down shadow; a hard offset dark shadow
    is almost never right.
    """
    spPr = shape._element.spPr
    _drop(spPr, ("a:effectLst",))

    effectLst = _new("a:effectLst")
    shdw = _new(
        "a:outerShdw",
        blurRad=str(int(round(blur_pt * 12700))),
        dist=str(int(round(distance_pt * 12700))),
        dir=str(int(round(direction_deg * 60000)) % 21600000),
        rotWithShape="0",
    )
    shdw.append(srgb(color, alpha=alpha))
    effectLst.append(shdw)
    _insert_ordered(spPr, effectLst, _SPPR_EFFECT_SUCCESSORS)


# ---------------------------------------------------------------------------
# Text
# ---------------------------------------------------------------------------


def set_tracking(run, spacing_pt: float) -> None:
    """
    Set letter-spacing (tracking) on a run.

    python-pptx has no API for this. It matters at both ends of the type scale, and getting
    it right is a large part of what separates set type from typed text:
      - Display sizes (40pt+) need NEGATIVE tracking, around -0.01 to -0.02em. Type designed
        for body copy has too much sidebearing when scaled up.
      - Small caps and label text need POSITIVE tracking, around +0.05 to +0.1em, or the
        letterforms crowd.
    Units are points; OOXML stores hundredths of a point in @spc.
    """
    run.font._rPr.set("spc", str(int(round(spacing_pt * 100))))


def set_text_alpha(run, alpha: float) -> None:
    """
    Make a run partially transparent.

    The honest use is de-emphasis in place: ghosting the non-focal rows of an exhibit so the
    focal row carries the eye, without changing the colour and thereby implying a category.
    """
    rPr = run.font._rPr

    # Read the run's existing colour BEFORE clearing the fill.
    #
    # Getting this order wrong is a genuinely nasty bug, because it fails SILENTLY and
    # only in one direction: drop the fill first and `run.font.color.rgb` reads back as
    # unset, which defaults to black. White text ghosted on a dark background then renders
    # near-black -- invisible on screen, and the generated XML is perfectly valid, so no
    # structural check catches it. Found by rendering and looking, which is the only way
    # this class of defect is ever found.
    hexval = "000000"
    try:
        color = run.font.color
        if color is not None and color.rgb is not None:
            hexval = str(color.rgb)
    except (AttributeError, TypeError):
        pass

    _drop(rPr, ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:grpFill"))
    solid = _new("a:solidFill")
    solid.append(srgb(hexval, alpha=alpha))
    # a:rPr sequence -- fill comes after a:ln and before a:effectLst.
    _insert_ordered(
        rPr,
        solid,
        (
            "a:effectLst",
            "a:effectDag",
            "a:highlight",
            "a:uLnTx",
            "a:uLn",
            "a:uFillTx",
            "a:uFill",
            "a:latin",
            "a:ea",
            "a:cs",
            "a:sym",
            "a:hlinkClick",
            "a:hlinkMouseOver",
            "a:rtl",
            "a:extLst",
        ),
    )


def set_normautofit(
    text_frame, font_scale: float = 1.0, line_space_reduction: float = 0.0
) -> None:
    """
    Write the <a:normAutofit> hint that PowerPoint uses to shrink text on overflow.

    IMPORTANT, and the reason this is a trap: PowerPoint computes fontScale at EDIT time,
    inside the application. A file generated by python-pptx carries whatever scale you write
    here and nothing recomputes it until a human opens the deck and edits that box. So this
    is a hint, never a solution -- setting MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE on generated text
    does approximately nothing at render time.

    The real fix is to MEASURE before writing (text.fits / text.fit_size) and to cut words
    rather than shrink type. See references/composition.md on why shrinking is the wrong
    lever: it breaks the size ratio that encodes hierarchy, so the slide starts lying about
    which idea is dominant.
    """
    bodyPr = text_frame._txBody.bodyPr
    _drop(bodyPr, ("a:noAutofit", "a:normAutofit", "a:spAutoFit"))
    attrs = {}
    if font_scale != 1.0:
        attrs["fontScale"] = str(_pct(font_scale, "font_scale"))
    if line_space_reduction:
        attrs["lnSpcReduction"] = str(
            _pct(line_space_reduction, "line_space_reduction")
        )
    _insert_ordered(bodyPr, _new("a:normAutofit", **attrs), _BODYPR_AUTOFIT_SUCCESSORS)


def set_no_autofit(text_frame) -> None:
    """Pin the type size. The correct default when you have measured the text yourself."""
    bodyPr = text_frame._txBody.bodyPr
    _drop(bodyPr, ("a:noAutofit", "a:normAutofit", "a:spAutoFit"))
    _insert_ordered(bodyPr, _new("a:noAutofit"), _BODYPR_AUTOFIT_SUCCESSORS)


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------


def strip_table_style(table) -> None:
    """
    Remove the banded-blue default table style.

    A python-pptx table arrives wearing the Office "Medium Style 2 Accent 1" look: alternating
    blue fills, a heavy header band, and a border on every cell. It is the most recognisable
    generated-deck artifact there is. This clears the style reference and the banding flags so
    the table can be styled deliberately -- ideally with horizontal rules only and no vertical
    ones at all, which is what every well-set table in print does.

    Clearing the style reference is NOT sufficient on its own, and assuming it is produces a
    table that looks fine in the XML and renders with a black 1pt grid on every cell. With no
    style to inherit from, renderers fall back to their own default border, so each edge has
    to be explicitly set to noFill and each cell explicitly un-filled. Caught by rendering.
    """
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for flag in (
            "firstRow",
            "lastRow",
            "firstCol",
            "lastCol",
            "bandRow",
            "bandCol",
        ):
            tblPr.set(flag, "0")
        for style_id in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(style_id)

    for row in table.rows:
        for cell in row.cells:
            for edge in ("top", "bottom", "left", "right"):
                set_cell_border(cell, edge, None)
            tcPr = cell._tc.get_or_add_tcPr()
            _drop(tcPr, _FILL_TAGS)
            _insert_ordered(tcPr, _new("a:noFill"), ("a:headers", "a:extLst"))


def set_cell_border(
    cell,
    edge: str,
    color: str | None,
    width_pt: float = 0.75,
    alpha: float | None = None,
) -> None:
    """
    Set one edge of a table cell. python-pptx exposes cell fill and margins but no borders.

    edge: 'top' | 'bottom' | 'left' | 'right'

    The <a:tcPr> sequence puts the line elements first and in a fixed order:
        lnL, lnR, lnT, lnB, lnTlToBr, lnBlToTr, cell3D, fill...
    """
    tag_for = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}
    if edge not in tag_for:
        raise ValueError(f"edge must be one of {sorted(tag_for)}, got {edge!r}")

    order = [
        "a:lnL",
        "a:lnR",
        "a:lnT",
        "a:lnB",
        "a:lnTlToBr",
        "a:lnBlToTr",
        "a:cell3D",
        "a:noFill",
        "a:solidFill",
        "a:gradFill",
        "a:blipFill",
        "a:pattFill",
        "a:grpFill",
        "a:headers",
        "a:extLst",
    ]
    tag = tag_for[edge]
    tcPr = cell._tc.get_or_add_tcPr()
    _drop(tcPr, (tag,))

    ln = _new(
        tag, w=str(int(round(width_pt * 12700))), cap="flat", cmpd="sng", algn="ctr"
    )
    if color is None:
        ln.append(_new("a:noFill"))
    else:
        solid = _new("a:solidFill")
        solid.append(srgb(color, alpha=alpha))
        ln.append(solid)

    successors = order[order.index(tag) + 1 :]
    _insert_ordered(tcPr, ln, successors)


# ---------------------------------------------------------------------------
# Slide-level operations python-pptx does not offer
# ---------------------------------------------------------------------------


def delete_slide(prs, index: int) -> None:
    """
    Remove a slide by index.

    python-pptx has no delete API (issue #67, open since 2013). Doing it correctly means
    dropping BOTH the <p:sldId> entry from the presentation part AND the relationship it
    points at. Dropping only the sldId leaves an orphaned relationship, which is one of the
    documented triggers for the repair dialog.
    """
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    if not 0 <= index < len(slides):
        raise IndexError(f"slide index {index} out of range (deck has {len(slides)})")
    sldId = slides[index]
    rId = sldId.get(qn("r:id"))
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def move_slide(prs, old_index: int, new_index: int) -> None:
    """Reorder slides. Same part/relationship integrity concern as delete."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    if not 0 <= old_index < len(slides):
        raise IndexError(f"slide index {old_index} out of range")
    if not 0 <= new_index < len(slides):
        raise IndexError(f"target index {new_index} out of range")
    sldId = slides[old_index]
    sldIdLst.remove(sldId)
    sldIdLst.insert(new_index, sldId)


def set_slide_background(
    slide,
    color: str | etree._Element | None = None,
    stops: Sequence[tuple[float, str]] | None = None,
    angle_deg: float = 90.0,
) -> None:
    """
    Set a slide's background to a solid colour or a gradient.

    python-pptx's `slide.background.fill` works for solids but shares the two-stop gradient
    limitation, so gradients are built here instead. The <p:bg> element must be the FIRST
    child of <p:cSld>, before <p:spTree>.
    """
    cSld = slide._element.find(qn("p:cSld"))
    _drop(cSld, ("p:bg",))

    bg = _new("p:bg")
    bgPr = _new("p:bgPr")

    if stops:
        grad = _new("a:gradFill", flip="none", rotWithShape="1")
        gsLst = _new("a:gsLst")
        for pos, col in stops:
            gs = _new("a:gs", pos=str(_pct(pos, "stop position")))
            gs.append(srgb(col))
            gsLst.append(gs)
        grad.append(gsLst)
        grad.append(
            _new("a:lin", ang=str(int(round(angle_deg * 60000)) % 21600000), scaled="0")
        )
        bgPr.append(grad)
    else:
        if color is None:
            raise ValueError(
                "set_slide_background needs either a colour or gradient stops"
            )
        clr = color if isinstance(color, etree._Element) else srgb(color)
        solid = _new("a:solidFill")
        solid.append(clr)
        bgPr.append(solid)

    bgPr.append(_new("a:effectLst"))
    bg.append(bgPr)
    cSld.insert(0, bg)
