"""
charts.py -- exhibits drawn from primitives.

WHY NOT NATIVE POWERPOINT CHARTS.

python-pptx can insert a real chart part, and for a chart someone will later edit that is the
right answer. For a chart nobody will edit -- which is every chart in a delivered deck -- it is
the wrong one, for three reasons that compound:

  1. A native chart arrives wearing the Office chart style: gridlines, an axis line, tick
     marks, a legend, and a title. Every one of those is chart junk in Tufte's sense, and
     removing them means reaching into the chart XML anyway, which is most of the work of
     drawing it yourself with none of the control.
  2. Legends. A native chart puts a legend somewhere and expects the reader to bounce between
     it and the marks. Direct labelling -- the label sitting on or beside the mark it names --
     is measurably faster to read and is what every well-made exhibit in print does. Native
     charts fight you on this.
  3. Renderer variance. A native chart is re-laid-out by whatever opens it, so LibreOffice,
     Google Slides and PowerPoint all produce slightly different geometry. Shapes are absolute:
     what is verified in the render loop is what appears in the room.

So these are shapes. The cost is that they are not editable as data, which for a delivered deck
is not a cost.

THE ONE RULE EVERY EXHIBIT HERE FOLLOWS: one chart, one message, and the message is the
headline. If the headline is "Revenue by region" you have labelled the chart, not stated its
finding. If it is "Three regions grew; the fourth carried the loss" the reader knows what to
look for before they look. That is the action-title convention, and it is the difference
between an exhibit and a picture of some data.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from . import oxml as X
from . import text as T
from .canvas import Box, inches, points
from .components import Ctx, _rect, hairline
from .tokens import LINE_SPACING, SCALE, SPACE, contrast_ratio, tracking_pt


@dataclass(frozen=True)
class Datum:
    """
    One mark in an exhibit.

    `emphasis` is the mechanism that makes an exhibit argue rather than merely display. Exactly
    one datum should normally carry emphasis=True: it renders in the accent while the rest
    render in a muted tone, so the eye lands on the mark the headline is about. This is the
    "ghosting" technique from consulting practice, and it is the single highest-leverage move
    available in a chart -- an exhibit where everything is equally coloured has no argument,
    only data.
    """

    label: str
    value: float
    emphasis: bool = False
    note: str | None = None


def _fmt(value: float, *, unit: str = "", decimals: int = 0) -> str:
    return f"{value:,.{decimals}f}{unit}"


def bar_chart(
    slide,
    ctx: Ctx,
    box: Box,
    data: Sequence[Datum],
    *,
    unit: str = "",
    decimals: int = 0,
    max_value: float | None = None,
    label_width_in: float = 2.4,
    show_track: bool = True,
    bar_height_pt: float = 26,
    emphasis_color: str | None = None,
):
    """
    A horizontal bar chart, direct-labelled, with no axis and no legend.

    Horizontal rather than vertical whenever the categories have names. Category labels set
    horizontally beside a bar are read at normal reading speed; the same labels under vertical
    columns either truncate or rotate, and rotated type is roughly half as fast to read. Use
    columns only when the x-axis is genuinely time.

    Anatomy, left to right:
        category label -- right-aligned, so all labels terminate at the bar origin, giving
                          the bars a single hard start line
        track          -- the full-width light grey ground the bar sits in. Doing the work an
                          axis would do (showing the maximum) without an axis line
        bar            -- accent for the emphasised datum, muted for the rest
        value          -- direct-labelled at the end of the bar. No axis needed at all
    """
    if not data:
        raise ValueError("bar_chart needs at least one datum")

    top = max(max_value or 0.0, max(d.value for d in data))
    if top <= 0:
        raise ValueError("bar_chart needs a positive maximum")

    label_w = inches(label_width_in)
    gap = points(SPACE["sm"])

    # Measure the widest value string rather than reserving a fixed inch. A fixed gutter wrapped
    # "121 marks" onto two lines in the first render -- a direct-labelled chart whose direct
    # labels wrap has lost the only advantage it had over a legend.
    value_w = max(
        inches(0.55),
        max(
            T.text_width_emu(
                _fmt(d.value, unit=unit, decimals=decimals),
                ctx.brand.font_major,
                SCALE["caption"],
                bold=True,
            )
            for d in data
        )
        # +2 x 7.2pt for the text frame's own left and right insets, which are not part of the
        # measured string width. Omitting them wrapped "121 marks" while "98 marks" fitted --
        # one extra character is exactly the margin this covers.
        + points(SPACE["sm"] + 14.4),
    )
    track_left = box.left + label_w + gap
    track_w = box.width - label_w - gap - value_w - gap

    n = len(data)
    row_h = box.height // n
    bar_h = points(bar_height_pt)
    accent = emphasis_color or ctx.accent
    # WHEN NOTHING IS EMPHASISED, THE BARS ARE PACIFIC, NOT GREY.
    #
    # A chart with no emphasis is an ordinary single-series chart, and KPMG's default fill for
    # one is Pacific #00B8F5 by a clear margin: #00B8F5 (116) > #1E49E2 (78) > #FFFFFF (50) >
    # #00338D (31). Rendering every bar in the de-emphasis grey was correct only for the
    # ghosting case; on a chart making no comparison it produced a grey chart, which is not a
    # KPMG chart and not a legible one from the back of a room.
    #
    # With emphasis present the ghosting stands: the emphasised bar takes the accent and the
    # rest go quiet, which is the mechanism that makes the exhibit argue.
    any_emphasis = any(d.emphasis for d in data)
    strong_muted = (
        (
            ctx.brand.color("faint")
            if not ctx.dark_ground
            else ctx.brand.color("on_dark_muted")
        )
        if any_emphasis
        else ctx.brand.color("chart_default")
    )

    # Fail rather than clip. A truncated category label is worse than no chart: the reader
    # cannot tell it was truncated, so they read a wrong label as a right one.
    for d in data:
        probe = Box(0, 0, label_w, max(bar_h + points(6), row_h))
        if not T.fit(d.label, probe, ctx.brand.font_minor, SCALE["caption"]).fits:
            raise ValueError(
                f"bar_chart label {d.label!r} does not fit the {label_width_in}in label "
                f"column. Shorten the label or raise label_width_in -- do not truncate it."
            )

    for i, d in enumerate(data):
        cy = box.top + i * row_h + (row_h - bar_h) // 2

        # The label box gets the FULL row height, anchored middle on the bar, not just the
        # bar's own height. A category label is often two lines ("Sign it off, then make it
        # repeatable") and a box one bar tall silently CLIPS the overflow rather than growing
        # -- so the chart renders with half a word and no error anywhere.
        # The FULL row height, not row_h minus a gap. Labels are anchored middle and the bar is
        # much shorter than the row, so adjacent labels cannot collide -- but shaving a gap off
        # left the box just under two lines, which rejected perfectly reasonable labels.
        label_box_h = max(bar_h + points(6), row_h)
        T.add_textbox(
            slide,
            Box(box.left, cy + bar_h // 2 - label_box_h // 2, label_w, label_box_h),
            d.label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.ink if d.emphasis else ctx.muted,
            bold=d.emphasis,
            align="right",
            anchor="middle",
        )

        if show_track:
            _rect(slide, Box(track_left, cy, track_w, bar_h), ctx.brand.color("track"))

        bar_w = max(points(2), int(track_w * (d.value / top)))
        _rect(
            slide,
            Box(track_left, cy, bar_w, bar_h),
            accent if d.emphasis else strong_muted,
        )

        T.add_textbox(
            slide,
            Box(track_left + track_w + gap, cy - points(3), value_w, bar_h + points(6)),
            _fmt(d.value, unit=unit, decimals=decimals),
            family=ctx.brand.font_major,
            size_pt=SCALE["caption"],
            bold=d.emphasis,
            color=ctx.ink if d.emphasis else ctx.muted,
            anchor="middle",
        )


def column_chart(
    slide,
    ctx: Ctx,
    box: Box,
    data: Sequence[Datum],
    *,
    unit: str = "",
    decimals: int = 0,
    max_value: float | None = None,
    baseline: bool = True,
    emphasis_color: str | None = None,
):
    """
    A vertical column chart. Use ONLY when the x-axis is time.

    Values sit above each column, so no y-axis is drawn at all. A single hairline baseline
    grounds the columns -- the one axis line worth keeping, because columns floating with no
    baseline read as unanchored.
    """
    if not data:
        raise ValueError("column_chart needs at least one datum")

    top_value = max(max_value or 0.0, max(d.value for d in data))
    if top_value <= 0:
        raise ValueError("column_chart needs a positive maximum")

    value_h = points(SCALE["caption"] * 1.6)
    label_h = points(SCALE["caption"] * 1.8)
    plot_h = box.height - value_h - label_h - points(SPACE["xs"])
    plot_top = box.top + value_h

    n = len(data)
    gap = points(SPACE["md"])
    accent = emphasis_color or ctx.accent
    strong_muted = (
        ctx.brand.color("faint")
        if not ctx.dark_ground
        else ctx.brand.color("on_dark_muted")
    )

    for i, d in enumerate(data):
        cell = box.row(n, i, gap=gap)
        col_h = max(points(2), int(plot_h * (d.value / top_value)))
        col_top = plot_top + plot_h - col_h

        _rect(
            slide,
            Box(cell.left, col_top, cell.width, col_h),
            accent if d.emphasis else strong_muted,
        )

        T.add_textbox(
            slide,
            Box(cell.left, col_top - value_h, cell.width, value_h),
            _fmt(d.value, unit=unit, decimals=decimals),
            family=ctx.brand.font_major,
            size_pt=SCALE["caption"],
            bold=d.emphasis,
            color=ctx.ink if d.emphasis else ctx.muted,
            align="center",
            anchor="bottom",
        )
        T.add_textbox(
            slide,
            Box(
                cell.left, plot_top + plot_h + points(SPACE["xs"]), cell.width, label_h
            ),
            d.label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.ink if d.emphasis else ctx.muted,
            bold=d.emphasis,
            align="center",
        )

    if baseline:
        hairline(
            slide,
            ctx,
            Box(box.left, plot_top + plot_h, box.width, 0),
            color=ctx.brand.color("faint"),
        )


def proportion_bar(
    slide,
    ctx: Ctx,
    box: Box,
    segments: Sequence[Datum],
    *,
    height_pt: float = 44,
    show_labels: bool = True,
):
    """
    A single stacked bar showing how a whole divides. Replaces a pie chart, always.

    A pie is a bad instrument and the reason is measurable rather than aesthetic: judging
    relative ANGLE is one of the least accurate perceptual tasks, well below judging position
    along a common scale or length. A stacked bar asks the reader to compare lengths on a
    shared baseline, which is among the most accurate. The only thing a pie does better is
    signal "these sum to a whole", and a single full-width bar signals that just as clearly.

    Labels sit below their segment when it is wide enough to hold them, and are dropped when
    it is not -- silently crowding a 3% segment with a label is worse than leaving it to the
    surrounding text.
    """
    if not segments:
        raise ValueError("proportion_bar needs at least one segment")
    total = sum(s.value for s in segments)
    if total <= 0:
        raise ValueError("proportion_bar needs a positive total")

    bar_h = points(height_pt)
    x = box.left
    label_top = box.top + bar_h + points(SPACE["sm"])

    for i, seg in enumerate(segments):
        w = (
            int(box.width * (seg.value / total))
            if i < len(segments) - 1
            else box.right - x
        )
        color = ctx.accent if seg.emphasis else ctx.brand.series(i)
        _rect(slide, Box(x, box.top, w, bar_h), color)

        if show_labels and w > inches(0.9):
            T.add_textbox(
                slide,
                Box(x, label_top, w, points(SCALE["caption"] * 3.0)),
                f"{seg.label}\n{_fmt(seg.value)}",
                family=ctx.brand.font_minor,
                size_pt=SCALE["caption"],
                color=ctx.ink if seg.emphasis else ctx.muted,
                bold=seg.emphasis,
                line_spacing=LINE_SPACING["footnote"],
            )
        x += w


def comparison(
    slide,
    ctx: Ctx,
    box: Box,
    left_title: str,
    left_items: Sequence[str],
    right_title: str,
    right_items: Sequence[str],
    *,
    left_color: str | None = None,
    right_color: str | None = None,
):
    """
    Two columns set against each other: before/after, current/target, us/them.

    The two sides are separated by a vertical hairline and nothing else -- no boxes, no fills.
    A comparison where each side is a filled panel reads as two unrelated slides side by side;
    the hairline says "these are two halves of one idea".

    Deliberately asymmetric in colour: the left side is muted and the right carries the accent,
    because a comparison almost always has a preferred side and the composition should say
    which. If genuinely neutral, pass both colours the same.
    """
    if len(left_items) > 5 or len(right_items) > 5:
        raise ValueError("comparison columns cap at five items each")

    gap = points(SPACE["xl"])
    col_w = (box.width - gap) // 2
    left_box = Box(box.left, box.top, col_w, box.height)
    right_box = Box(box.left + col_w + gap, box.top, col_w, box.height)

    _rect(
        slide,
        Box(box.left + col_w + gap // 2, box.top, points(0.75), box.height),
        ctx.rule,
    )

    for col, title, items, color in (
        (left_box, left_title, left_items, left_color or ctx.muted),
        (right_box, right_title, right_items, right_color or ctx.accent),
    ):
        # Measured, not allocated. points(SCALE["h3"] * 1.5) reserves 1.5x the POINT SIZE while
        # a line occupies size x the font's own line height, so the rule below landed on the
        # column title in the first render -- the same defect as the headline accent rule.
        title_fit = T.fit(
            title,
            Box(col.left, col.top, col.width, points(400)),
            ctx.brand.font_major,
            SCALE["h3"],
            bold=True,
            line_spacing=LINE_SPACING["h3"],
        )
        title_h = points(title_fit.height_pt)
        T.add_textbox(
            slide,
            Box(col.left, col.top, col.width, title_h),
            title,
            family=ctx.brand.font_major,
            size_pt=SCALE["h3"],
            bold=True,
            color=color,
            line_spacing=LINE_SPACING["h3"],
        )
        _rect(
            slide,
            Box(
                col.left,
                col.top + title_h + points(SPACE["xs"]),
                inches(1.0),
                points(3),
            ),
            color,
        )

        item_top = col.top + title_h + points(SPACE["xs"] + SPACE["md"])
        item_h = (col.bottom - item_top) // max(1, len(items))
        for j, item in enumerate(items):
            y = item_top + j * item_h
            if j > 0:
                hairline(
                    slide, ctx, Box(col.left, y - points(SPACE["xs"]), col.width, 0)
                )
            T.add_textbox(
                slide,
                Box(col.left, y, col.width, item_h - points(SPACE["xs"])),
                item,
                family=ctx.brand.font_minor,
                size_pt=SCALE["small"],
                color=ctx.ink,
                line_spacing=LINE_SPACING["body"],
                anchor="top",
            )


def table(
    slide,
    ctx: Ctx,
    box: Box,
    header: Sequence[str],
    body_rows: Sequence[Sequence[str]],
    *,
    col_widths: Sequence[float] | None = None,
    emphasis_row: int | None = None,
):
    """
    A table with horizontal rules only.

    Every well-set table in print does this: a rule under the header, hairlines between rows,
    and NO vertical rules at all. Columns are separated by space, which is sufficient, and
    vertical rules turn a table into a grid of cells that the eye has to climb out of. The
    default PowerPoint table -- banded blue fills, a border on every cell -- is the single most
    recognisable generated-deck artifact there is.

    Uses real shapes rather than a pptx table, for the same absolute-geometry reason as the
    charts. Use a real table only when someone must edit it later.
    """
    if not header:
        raise ValueError("table needs a header row")
    ncols = len(header)
    for i, row in enumerate(body_rows):
        if len(row) != ncols:
            raise ValueError(f"row {i} has {len(row)} cells, header has {ncols}")

    if col_widths is None:
        col_widths = [1.0 / ncols] * ncols
    if abs(sum(col_widths) - 1.0) > 0.01:
        raise ValueError(f"col_widths must sum to 1.0, got {sum(col_widths)}")

    xs, x = [], box.left
    for frac in col_widths:
        xs.append(x)
        x += int(box.width * frac)
    widths = [int(box.width * f) for f in col_widths]

    header_h = points(SCALE["caption"] * 2.2)
    for i, cell in enumerate(header):
        T.add_textbox(
            slide,
            Box(xs[i], box.top, widths[i] - points(SPACE["sm"]), header_h),
            cell,
            family=ctx.brand.font_minor,
            size_pt=SCALE["footnote"],
            bold=True,
            color=ctx.muted,
            caps=True,
            tracking_pt=SCALE["footnote"] * 0.08,
            anchor="bottom",
        )

    hairline(
        slide,
        ctx,
        Box(box.left, box.top + header_h, box.width, 0),
        color=ctx.brand.color("faint"),
        thickness_pt=1.0,
    )

    rows_top = box.top + header_h + points(SPACE["xs"])
    row_h = (box.bottom - rows_top) // max(1, len(body_rows))

    for r, row in enumerate(body_rows):
        y = rows_top + r * row_h
        emphasised = emphasis_row == r
        if r > 0:
            hairline(slide, ctx, Box(box.left, y, box.width, 0))
        for c, cell in enumerate(row):
            T.add_textbox(
                slide,
                Box(
                    xs[c],
                    y + points(SPACE["xs"]),
                    widths[c] - points(SPACE["sm"]),
                    row_h - points(SPACE["xs"]),
                ),
                cell,
                family=ctx.brand.font_minor,
                size_pt=SCALE["caption"],
                bold=emphasised,
                color=ctx.ink if emphasised else ctx.muted,
                line_spacing=LINE_SPACING["body"],
            )


# ---------------------------------------------------------------------------
# The ring gauge -- KPMG's "Key findings" device
# ---------------------------------------------------------------------------
#
# Always page 4 in all four measured CEO decks, 19-27 rings per page. It is the single most
# recognisable non-bar exhibit in the system and we did not have it.
#
# THE ONE THING THAT MAKES IT NOT A GENERIC DONUT: the value arc is ~1.9x thicker than the
# track AND SITS ON A LARGER RADIUS, so it stands proud of the ring rather than filling it.
# A uniform-width donut -- which is what every charting library emits -- reads as generic at a
# glance, and this does not. Measured:
#
#     outer hairline   stroke #E5E5E5  1.000pt   diameter 73.2pt
#     track band       stroke #E5E5E5  5.000pt   diameter 57.4pt, fill #FFFFFF
#     value arc        series colour   9.579pt   ROUND cap, on a LARGER radius than the track
#     centre           the value, display weight, IN THE ARC'S OWN COLOUR
#
# Round caps are legal here and nowhere else in the package: corpus line caps run
# butt 13,482 / round 3,386 / square 8, and the round ones are these arcs.
#
# COLOUR-CODE BY CONTENT STREAM, not by position in a series. On the measured page the rings
# run navy, cobalt, navy, purple -- grouped so that each column of findings carries one colour
# and the purple column's rings are purple. The palette breadth CODES MEANING there; it is not
# decoration, and cycling `brand.series(i)` per ring would destroy the thing it is doing.

_RING_OUTER_D = 73.2
_RING_TRACK_D = 57.4
_RING_ARC_D = 64.0  # proud of the track, inside the hairline
_RING_TRACK_W = 5.0
_RING_ARC_W = 9.579
_RING_HAIRLINE_W = 1.0


def _oval(slide, box: Box):
    from pptx.enum.shapes import MSO_SHAPE

    return slide.shapes.add_shape(MSO_SHAPE.OVAL, *box.as_tuple())


def _arc(slide, box: Box, start_deg: float, sweep_deg: float):
    """
    An open arc, as an OOXML `arc` preset with its two adjustment angles set.

    The preset measures angles in 60000ths of a degree CLOCKWISE FROM EAST, so 12 o'clock is
    270 and not 90 -- the single easiest thing to get wrong here, and it fails silently by
    drawing a correct-looking arc starting in the wrong place.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    shape = slide.shapes.add_shape(MSO_SHAPE.ARC, *box.as_tuple())
    av = shape._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    for gd in list(av):
        av.remove(gd)
    for name, deg in (("adj1", start_deg), ("adj2", start_deg + sweep_deg)):
        gd = av.makeelement(qn("a:gd"), {})
        gd.set("name", name)
        gd.set("fmla", f"val {int(round((deg % 360) * 60000))}")
        av.append(gd)
    return shape


def ring_gauge(
    slide,
    ctx: Ctx,
    box: Box,
    value: float,
    *,
    color: str | None = None,
    label: str | None = None,
    diameter_pt: float | None = None,
):
    """
    One ring gauge: hairline, track, proud value arc, and the figure in the centre.

    `value` is 0.0-1.0. `color` is the CONTENT STREAM's colour -- pass the same colour for
    every ring belonging to one group, and a different one per group. Defaults to Cobalt.

    `label` is the centre text; it defaults to the value as a whole percentage. Passing it
    explicitly is how a ring shows "1-3 yrs" or "10-20%" rather than a number it does not have.

    Returns the drawn Box, so a caller can place a gloss beside it from real geometry.
    """
    if not 0.0 <= value <= 1.0:
        raise ValueError(f"ring_gauge value must be a fraction 0.0-1.0, got {value}")

    arc_color = color or ctx.brand.color("accent")
    track = ctx.brand.color("track")

    d_outer = points(diameter_pt or _RING_OUTER_D)
    scale = d_outer / points(_RING_OUTER_D)
    d_track = int(points(_RING_TRACK_D) * scale)
    d_arc = int(points(_RING_ARC_D) * scale)

    cx = (
        box.left + min(box.width, d_outer) // 2
        if box.width < d_outer
        else box.left + d_outer // 2
    )
    cy = box.top + d_outer // 2

    def centred(diameter: int) -> Box:
        return Box(cx - diameter // 2, cy - diameter // 2, diameter, diameter)

    # 1. the outer hairline -- the mark that stops the ring floating in white space
    hair = _oval(slide, centred(d_outer))
    X.set_no_fill(hair)
    X.set_line(hair, track, width_pt=_RING_HAIRLINE_W * scale)
    X.clear_effects(hair)
    hair.shadow.inherit = False

    # 2. the track band, filled white so it reads as a physical ring on any ground
    band = _oval(slide, centred(d_track))
    X.set_fill(band, ctx.brand.color("canvas"))
    X.set_line(band, track, width_pt=_RING_TRACK_W * scale)
    X.clear_effects(band)
    band.shadow.inherit = False

    # 3. the value arc: thicker, on a larger radius, round-capped. This is the whole device.
    arc = _arc(slide, centred(d_arc), 270.0, max(value, 0.001) * 360.0)
    X.set_no_fill(arc)
    X.set_line(arc, arc_color, width_pt=_RING_ARC_W * scale, cap="rnd")
    X.clear_effects(arc)
    arc.shadow.inherit = False

    # 4. the figure, in the arc's own colour -- not the body ink, not the headline blue.
    #
    # The text frame is given the OUTER diameter, not the track's. Sized to the track it wrapped
    # "73%" onto two lines with the % sitting below the ring -- the box was narrower than the
    # string at every size that reads. The glyphs are centred and stay inside the ring; only
    # the box has to be wider than the ring, and a wider box costs nothing.
    text = label if label is not None else f"{value * 100:.0f}%"
    size = max(
        SCALE["caption"],
        min(SCALE["h2"], round((diameter_pt or _RING_OUTER_D) * 0.40)),
    )
    T.add_textbox(
        slide,
        Box(cx - d_outer // 2, cy - points(size * 0.62), d_outer, points(size * 1.4)),
        text,
        family=ctx.brand.font_major,
        size_pt=size,
        bold=True,
        color=arc_color,
        align="center",
        anchor="top",
        line_spacing=1.0,
        tracking_pt=tracking_pt("h2", size),
    )
    return centred(d_outer)


# ---------------------------------------------------------------------------
# Four arguments a bar chart cannot make
# ---------------------------------------------------------------------------
#
# Everything above this line answers "how big, compared to what". These four answer questions
# the bar has no form for, and each was added because a real page in this course had to say the
# thing in prose for want of an exhibit:
#
#   unit_grid    HOW MANY OF THEM -- a proportion as countable things rather than as a length
#   stacked_bar  WHAT IT IS MADE OF, across several cases, on one shared 100% baseline
#   waterfall    WHY THE NUMBER MOVED -- a delta decomposed into its signed contributions
#   matrix_2x2   WHERE EACH ONE SITS on two judgments at once


def unit_grid(
    slide,
    ctx: Ctx,
    box: Box,
    filled: int,
    total: int,
    *,
    columns: int = 10,
    caption: str | None = None,
    color: str | None = None,
) -> Box:
    """
    A proportion drawn as countable marks: `filled` squares in the accent out of `total`.

    USE IT WHEN THE DENOMINATOR IS PEOPLE, OR ANYTHING ELSE THE ROOM CAN PICTURE. "19% to 61%"
    is a pair of numbers to be believed; 26 squares of 140 filling up to 86 is a thing the eye
    counts, and a proportion that has been counted is one the room now owns. That is the entire
    argument for this form over a bar, and it is why the caption should name the UNIT -- "26 of
    the 140 in this room" -- rather than restate the percentage.

    It is the wrong form above roughly 200 marks (they stop being countable and become texture)
    and the wrong form for a proportion of something abstract, where 140 squares imply a
    precision that a share of revenue does not have. Squares, not dots: a dot grid reads as
    decoration and a grid of squares reads as units.

    Marks size themselves from whichever of the box's dimensions binds, so the grid always fits
    the region it was given rather than overflowing it by a row.
    """
    if total <= 0:
        raise ValueError("unit_grid needs a positive total")
    if not 0 <= filled <= total:
        raise ValueError(f"filled must be 0..{total}, got {filled}")
    if total > 200:
        raise ValueError(
            f"unit_grid was given {total} marks. Past about 200 they are texture rather than "
            f"things anyone counts, which is the only reason to use this form -- use a "
            f"proportion_bar."
        )
    if columns < 1:
        raise ValueError("unit_grid needs at least one column")

    rows_n = -(-total // columns)
    gap_ratio = (
        0.28  # a quarter-mark gutter; tighter and the grid reads as a filled block
    )

    cap_h = (
        0
        if not caption
        else points(
            T.fit(
                caption,
                Box(0, 0, box.width, points(300)),
                ctx.brand.font_minor,
                SCALE["caption"],
                line_spacing=LINE_SPACING["body"],
            ).height_pt
        )
    )
    grid_h = box.height - cap_h - (points(SPACE["md"]) if caption else 0)

    # WHICHEVER DIMENSION BINDS, BINDS. Sized from the width alone, a 14-row grid in a short
    # region runs off the bottom of the page silently -- the shapes are all valid and nothing
    # measures them.
    by_width = int(box.width / (columns + (columns - 1) * gap_ratio))
    by_height = int(grid_h / (rows_n + (rows_n - 1) * gap_ratio))
    unit = min(by_width, by_height)
    if unit <= 0:
        raise ValueError("unit_grid was given a region too small to hold a single mark")
    gap = int(unit * gap_ratio)

    on = color or ctx.accent
    # THE UNFILLED MARK MUST BE DIMMER THAN THE GROUND'S IDEA OF QUIET, AND ON DARK THAT NEEDS
    # ALPHA RATHER THAN A COLOUR. #E5E5E5 on white is a whisper; the dark-ground equivalent role
    # is #B2B2B2, which against #0C233C is BRIGHTER than the filled accent marks -- so 114 empty
    # squares out-shouted the 26 that carried the finding. Held back to 40% it recedes, which is
    # the job. On light the flat track colour is already correct.
    off = (
        ctx.brand.color("track")
        if not ctx.dark_ground
        else ctx.brand.color("on_dark_muted")
    )
    off_alpha = None if not ctx.dark_ground else 0.40

    for i in range(total):
        r, c = divmod(i, columns)
        mark = Box(
            box.left + c * (unit + gap),
            box.top + r * (unit + gap),
            unit,
            unit,
        )
        if i < filled:
            _rect(slide, mark, on)
        else:
            _rect(slide, mark, off, alpha=off_alpha)

    grid_bottom = box.top + rows_n * unit + (rows_n - 1) * gap
    if caption:
        T.add_textbox(
            slide,
            Box(box.left, grid_bottom + points(SPACE["md"]), box.width, cap_h),
            caption,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.muted,
            line_spacing=LINE_SPACING["body"],
            anchor="top",
        )
        grid_bottom += points(SPACE["md"]) + cap_h
    return Box(box.left, box.top, box.width, grid_bottom - box.top)


def stacked_bar(
    slide,
    ctx: Ctx,
    box: Box,
    categories: Sequence[tuple[str, Sequence[Datum]]],
    *,
    unit: str = "",
    decimals: int = 0,
    label_width_in: float = 2.4,
    bar_height_pt: float = 38,
) -> None:
    """
    Composition across two to five cases, each normalised to its own 100%, labelled in the bar.

    IT SHOWS SHARE AND HIDES SIZE, and that is the trade to make deliberately. Every bar runs
    the full measure, so a case worth ten times another looks the same width -- which is right
    when the claim is "the MIX changed" and a lie when it is "and it grew". If magnitude is part
    of the argument, this is the wrong exhibit.

    The segment NAMES are set once, above the first bar, and the later bars carry values only.
    That is the direct-labelling answer to a form that would otherwise need a legend: the first
    bar teaches the colours and every bar after it is read without leaving the exhibit.

    Every case must carry the same segments in the same order -- the function raises otherwise,
    because a colour that means "unattended" in one bar and "review" in the next has stopped
    coding anything. Emphasis works as it does everywhere else: mark one segment and the others
    go to the track grey, so the exhibit argues about one band rather than displaying four.

    Segments are separated by a hairline of the ground colour rather than butted together. At
    projector distance two adjacent mid-blues with no gap read as one segment.
    """
    if not 2 <= len(categories) <= 5:
        raise ValueError(
            f"stacked_bar takes 2 to 5 categories, got {len(categories)}. One is a "
            f"proportion_bar; six rows of five segments is a table."
        )
    keys = [tuple(d.label for d in segs) for _, segs in categories]
    if len(set(keys)) != 1:
        raise ValueError(
            f"every category must carry the same segments in the same order, got {keys}. A "
            f"colour that means one thing in the first bar and another in the second is not "
            f"coding anything."
        )
    if not 2 <= len(keys[0]) <= 5:
        raise ValueError(f"stacked_bar takes 2 to 5 segments, got {len(keys[0])}")

    label_w = inches(label_width_in)
    gap = points(SPACE["sm"])
    bar_left = box.left + label_w + gap
    bar_w = box.right - bar_left
    bar_h = points(bar_height_pt)
    seam = points(1.5)

    n = len(categories)
    key_h = points(SCALE["caption"] * 1.6)
    row_h = (box.height - key_h) // n
    any_emphasis = any(d.emphasis for _, segs in categories for d in segs)
    key_x: list[int] = []
    ghost = (
        ctx.brand.color("track")
        if not ctx.dark_ground
        else ctx.brand.color("on_dark_muted")
    )

    for r, (name, segs) in enumerate(categories):
        total = sum(s.value for s in segs)
        if total <= 0:
            raise ValueError(
                f"category {name!r} sums to {total}; it needs a positive total"
            )
        top = box.top + key_h + r * row_h + (row_h - bar_h) // 2

        T.add_textbox(
            slide,
            Box(box.left, top - points(4), label_w, bar_h + points(8)),
            name,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.ink,
            align="right",
            anchor="middle",
        )

        x = bar_left
        for i, seg in enumerate(segs):
            w = int(bar_w * (seg.value / total)) if i < len(segs) - 1 else box.right - x
            fill = ctx.brand.series(i) if (not any_emphasis or seg.emphasis) else ghost
            draw_w = max(points(2), w - (seam if i < len(segs) - 1 else 0))
            _rect(slide, Box(x, top, draw_w, bar_h), fill)

            # The value goes IN the segment, in whichever of white or near-black actually
            # measures against that segment's own fill. Deciding it by "is the slide dark"
            # gets it wrong on exactly the colours that matter -- Pacific is a light fill on
            # a dark slide, and white on Pacific is 2.0:1.
            value_str = _fmt(seg.value, unit=unit, decimals=decimals)
            need_w = T.text_width_emu(
                value_str, ctx.brand.font_major, SCALE["caption"], bold=True
            ) + points(SPACE["md"])
            if draw_w >= need_w:
                dark_ink, light_ink = (
                    ctx.brand.color("ink_strong"),
                    ctx.brand.color("on_dark"),
                )
                fg = (
                    light_ink
                    if contrast_ratio(light_ink, fill) >= contrast_ratio(dark_ink, fill)
                    else dark_ink
                )
                T.add_textbox(
                    slide,
                    Box(
                        x + points(SPACE["xs"]),
                        top,
                        draw_w - points(SPACE["xs"]),
                        bar_h,
                    ),
                    value_str,
                    family=ctx.brand.font_major,
                    size_pt=SCALE["caption"],
                    bold=True,
                    color=fg,
                    anchor="middle",
                )

            if r == 0:
                key_x.append(x)
            x += w

    # THE KEY NAMES EVERY SEGMENT OR THE EXHIBIT HAS A COLOUR IT NEVER EXPLAINS, and that is
    # the legend problem this form exists to avoid. Dropping a key label when its segment was
    # too narrow to hold it -- the same rule `proportion_bar` uses correctly for labels sitting
    # UNDER a segment -- left an 8% band unnamed in the first render. A key label sits ABOVE
    # the bar, in free space, so it is allowed to be wider than the segment it names: each is
    # placed at its segment's left edge, pushed right only where the one before it would
    # collide, and the last is pulled back to end at the bar's right edge.
    key_w = [
        T.text_width_emu(
            seg.label.upper(),
            ctx.brand.font_minor,
            SCALE["caption"],
            bold=seg.emphasis,
            tracking_pt=SCALE["caption"] * 0.08,
        )
        + points(SPACE["sm"])
        for seg in categories[0][1]
    ]
    if sum(key_w) > bar_w:
        raise ValueError(
            f"the segment names {[s.label for s in categories[0][1]]} need "
            f"{sum(key_w) / 914400:.2f}in of key and the bar is {bar_w / 914400:.2f}in wide. "
            f"Shorten the names -- a stacked bar with an unnamed band is a chart with a legend "
            f"nobody printed."
        )
    placed_x = list(key_x)
    for i in range(1, len(placed_x)):
        placed_x[i] = max(placed_x[i], placed_x[i - 1] + key_w[i - 1])
    placed_x[-1] = min(placed_x[-1], box.right - key_w[-1])
    for i in range(len(placed_x) - 2, -1, -1):
        placed_x[i] = min(placed_x[i], placed_x[i + 1] - key_w[i])

    for i, seg in enumerate(categories[0][1]):
        T.add_textbox(
            slide,
            Box(placed_x[i], box.top, key_w[i], key_h),
            seg.label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            bold=seg.emphasis,
            caps=True,
            tracking_pt=SCALE["caption"] * 0.08,
            color=ctx.ink if seg.emphasis or not any_emphasis else ctx.muted,
            anchor="bottom",
        )


def waterfall(
    slide,
    ctx: Ctx,
    box: Box,
    start: Datum,
    steps: Sequence[Datum],
    end: Datum,
    *,
    unit: str = "",
    decimals: int = 0,
) -> None:
    """
    A total taken apart: opening bar, the signed contributions that moved it, closing bar.

    THE ONE EXHIBIT THAT EXPLAINS A DELTA. A pair of bars says the number went from 40 to 62; a
    waterfall says which three things did it and in what proportion, which is the finding. Use
    it whenever a page's argument is "and here is WHY it changed".

    The arithmetic is checked, not assumed: start plus the steps must reach `end`, and the
    function raises if it does not. A waterfall whose floating bars do not land on its own
    closing bar is a picture that contradicts itself, and it is the standard way this exhibit
    goes wrong -- an unbalanced deck is invisible to every structural check and obvious to
    anyone in the room who adds up the labels.

    DIRECTION IS CODED BY THE SIGNED LABEL FIRST AND BY HUE SECOND, because this palette has no
    red and inventing one would be the single most visible off-brand mark on the page. Increases
    take Pacific, decreases take Purple, the two totals take the ink. A reader who does not
    learn the hue still reads "-12" and is not misled, which is the property a red/green coding
    does not have for the ~8% of any room who cannot separate them.

    Non-negative running totals only. A series that crosses zero needs a drawn zero line and
    bars on both sides of it, which is a different component.
    """
    if not 1 <= len(steps) <= 5:
        raise ValueError(f"waterfall takes 1 to 5 steps, got {len(steps)}")

    running = [start.value]
    for s in steps:
        running.append(running[-1] + s.value)
    if any(v < 0 for v in running):
        raise ValueError(
            f"waterfall running totals {running} go negative. A series crossing zero needs a "
            f"drawn zero line and bars on both sides of it -- a different exhibit."
        )
    if abs(running[-1] - end.value) > max(0.005 * abs(end.value or 1.0), 1e-9):
        raise ValueError(
            f"the steps carry {start.value:g} to {running[-1]:g}, but `end` says "
            f"{end.value:g}. A waterfall whose bars miss its own closing total is a picture "
            f"that contradicts itself."
        )

    bars = [start, *steps, end]
    top_value = max(max(running), end.value)
    if top_value <= 0:
        raise ValueError("waterfall needs a positive maximum")

    value_h = points(SCALE["caption"] * 1.6)
    label_h = points(SCALE["caption"] * 2.4)
    plot_h = box.height - value_h - label_h - points(SPACE["xs"])
    plot_top = box.top + value_h
    baseline_y = plot_top + plot_h

    def y_of(v: float) -> int:
        return baseline_y - int(plot_h * (v / top_value))

    # THE ANCHOR IS THE HEAVIEST MARK ON A LIGHT PAGE AND MUST NOT BE ON A DARK ONE. Navy on
    # white is weighty without being the loudest thing there -- the accent still wins. Its naive
    # dark-ground inverse, white, IS the loudest thing there, so the two totals shouted down the
    # emphasised step they exist to frame. #B2B2B2 is the mark that anchors without arguing.
    total_fill = (
        ctx.brand.color("on_dark_muted") if ctx.dark_ground else ctx.brand.color("ink")
    )
    up_fill = ctx.brand.color("chart_default")
    down_fill = (
        ctx.brand.color("pair_light_purple")
        if ctx.dark_ground
        else ctx.brand.color("pair_purple")
    )

    n = len(bars)
    gap = points(SPACE["md"])
    for i, d in enumerate(bars):
        cell = box.row(n, i, gap=gap)
        if i == 0 or i == n - 1:
            bar_top, bar_bottom = y_of(d.value), baseline_y
            fill, shown = total_fill, _fmt(d.value, unit=unit, decimals=decimals)
        else:
            lo, hi = sorted((running[i - 1], running[i]))
            bar_top, bar_bottom = y_of(hi), y_of(lo)
            fill = up_fill if d.value >= 0 else down_fill
            shown = f"{'+' if d.value >= 0 else '-'}{_fmt(abs(d.value), unit=unit, decimals=decimals)}"
        if d.emphasis:
            fill = ctx.accent

        _rect(
            slide,
            Box(cell.left, bar_top, cell.width, max(points(2), bar_bottom - bar_top)),
            fill,
        )
        T.add_textbox(
            slide,
            Box(cell.left, bar_top - value_h, cell.width, value_h),
            shown,
            family=ctx.brand.font_major,
            size_pt=SCALE["caption"],
            bold=True,
            color=ctx.ink,
            align="center",
            anchor="bottom",
        )
        T.add_textbox(
            slide,
            Box(cell.left, baseline_y + points(SPACE["xs"]), cell.width, label_h),
            d.label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            color=ctx.muted,
            align="center",
            line_spacing=LINE_SPACING["footnote"],
        )

        # The step leader: a hairline across the gutter at the level the previous bar left off,
        # which is what makes the floating bars read as one continuous account rather than as
        # five unconnected columns.
        if i < n - 1:
            level = y_of(running[min(i, len(running) - 1)])
            nxt = box.row(n, i + 1, gap=gap)
            _rect(
                slide,
                Box(
                    cell.right,
                    level - points(0.375),
                    nxt.left - cell.right,
                    points(0.75),
                ),
                ctx.rule_strong,
            )

    hairline(
        slide,
        ctx,
        Box(box.left, baseline_y, box.width, 0),
        color=ctx.brand.color("faint"),
    )


@dataclass(frozen=True)
class Plot:
    """
    One item on a `matrix_2x2`. `x` and `y` are 0.0-1.0, with y = 0.0 at the BOTTOM of the field.
    """

    label: str
    x: float
    y: float
    emphasis: bool = False


def matrix_2x2(
    slide,
    ctx: Ctx,
    box: Box,
    items: Sequence[Plot],
    *,
    x_axis: tuple[str, str],
    y_axis: tuple[str, str],
    quadrants: tuple[str, str, str, str] | None = None,
    dot_pt: float = 13.0,
) -> Box:
    """
    Two judgments crossed, everything placed against both, one item marked.

    The most recognisable exhibit in consulting and the most easily faked, so the rule is worth
    stating: THE AXES MUST BE THINGS THE READER CAN DISAGREE WITH. "Effort against payoff" is a
    matrix; "importance against priority" is the same axis twice and produces a diagonal line
    of dots that looks like analysis and contains none. If every item lands in one quadrant, the
    axes were chosen to produce that answer and the exhibit should be cut.

    `x_axis` and `y_axis` are each `(low, high)` and are set HORIZONTALLY, outside the field --
    the y-axis pair right-aligned into the gutter at the field's top and bottom. Rotated type is
    about half as fast to read, which is too high a price for a two-word label.

    `quadrants` names the four corners in reading order (top-left, top-right, bottom-left,
    bottom-right) and sets each at the field's OUTER corner. NAMING THEM SHRINKS THE PLOTTABLE
    AREA rather than trusting the caller to steer clear: y = 1.0 lands below the top band and
    y = 0.0 above the bottom one, so a dot cannot be drawn through a quadrant name. Reserving it
    structurally is the only version of this that holds -- in the first render "Reconciliations"
    at (0.80, 0.82) sat directly under "HAND IT OVER", and every structural check passed. Pass
    None when the quadrants have no names worth having; four invented names are worse than none.

    A LABEL STAYS INSIDE ITS OWN DOT'S QUADRANT, on whichever side of the dot has room for it.
    Labelling everything to the right runs the right-hand column's words off the field; choosing
    the side from the field's midpoint alone drew "Draft memos" at x = 0.62 straight through the
    centre divider, which is the rule-through-a-label defect in its purest form.
    """
    if not 2 <= len(items) <= 8:
        raise ValueError(
            f"matrix_2x2 takes 2 to 8 items, got {len(items)}. Past eight the labels collide "
            f"and the exhibit is a scatter plot wanting real axes."
        )
    for p in items:
        if not (0.0 <= p.x <= 1.0 and 0.0 <= p.y <= 1.0):
            raise ValueError(f"{p.label!r} is at ({p.x}, {p.y}); both must be 0.0-1.0")

    gutter = inches(1.3)
    axis_h = points(SCALE["footnote"] * 1.8)
    field = Box(
        box.left + gutter,
        box.top,
        box.width - gutter,
        box.height - axis_h - points(SPACE["xs"]),
    )

    # A FAINT FIELD, NOT A WIREFRAME. Drawn as four hairlines and a cross the exhibit measured
    # 8.0% ink -- on the density gate's floor, and it looked it: a page of thin grey lines with
    # some dots in it. The near-white canvas_alt gives the quadrants a body to sit in without
    # adding a single unit of colour, which is the difference between a matrix and a diagram of
    # one. On a dark ground there is no equivalent tint and the hairlines already read, so none
    # is drawn.
    if not ctx.dark_ground:
        _rect(slide, field, ctx.brand.color("canvas_alt"))
    for edge in (
        Box(field.left, field.top, field.width, points(0.75)),
        Box(field.left, field.bottom, field.width, points(0.75)),
        Box(field.left, field.top, points(0.75), field.height),
        Box(field.right, field.top, points(0.75), field.height),
    ):
        _rect(slide, edge, ctx.rule_strong)
    _rect(
        slide,
        Box(field.left, field.center_y, field.width, points(0.75)),
        ctx.rule_strong,
    )
    _rect(
        slide,
        Box(field.center_x, field.top, points(0.75), field.height),
        ctx.rule_strong,
    )

    def _axis_label(b: Box, text_str: str, align: str) -> None:
        T.add_textbox(
            slide,
            b,
            text_str,
            family=ctx.brand.font_minor,
            size_pt=SCALE["footnote"],
            bold=True,
            caps=True,
            tracking_pt=SCALE["footnote"] * 0.1,
            color=ctx.muted,
            align=align,
            anchor="middle",
            line_spacing=LINE_SPACING["footnote"],
        )

    # 3.4 line-heights, not 2.6. A two-word axis label wraps in a 1.3in gutter, and at 2.6 the
    # second line was 4pt past the bottom of its own box -- drawn over whatever sat beneath it.
    lab_h = points(SCALE["footnote"] * 3.4)
    y_low, y_high = y_axis
    _axis_label(
        Box(box.left, field.top, gutter - points(SPACE["sm"]), lab_h), y_high, "right"
    )
    _axis_label(
        Box(box.left, field.bottom - lab_h, gutter - points(SPACE["sm"]), lab_h),
        y_low,
        "right",
    )
    x_low, x_high = x_axis
    half = field.width // 2 - points(SPACE["sm"])
    _axis_label(
        Box(field.left, field.bottom + points(SPACE["xs"]), half, axis_h), x_low, "left"
    )
    _axis_label(
        Box(field.right - half, field.bottom + points(SPACE["xs"]), half, axis_h),
        x_high,
        "right",
    )

    inset = points(SPACE["sm"])
    q_h = points(SCALE["footnote"] * 1.5)
    band = (q_h + 2 * inset) if quadrants else points(SPACE["md"])
    plot_top = field.top + band
    plot_h = field.height - 2 * band
    if plot_h <= 0:
        raise ValueError("matrix_2x2 was given a region too short to plot in")

    if quadrants:
        q_w = field.width // 2 - 2 * inset
        for text_str, (qx, qy, align) in zip(
            quadrants,
            (
                (field.left + inset, field.top + inset, "left"),
                (field.center_x + inset, field.top + inset, "right"),
                (field.left + inset, field.bottom - inset - q_h, "left"),
                (field.center_x + inset, field.bottom - inset - q_h, "right"),
            ),
        ):
            T.add_textbox(
                slide,
                Box(qx, qy, q_w, q_h),
                text_str,
                family=ctx.brand.font_minor,
                size_pt=SCALE["footnote"],
                bold=True,
                caps=True,
                tracking_pt=SCALE["footnote"] * 0.1,
                color=ctx.faint,
                align=align,
                anchor="middle",
            )

    label_w = inches(1.7)
    for p in items:
        d = points(dot_pt * (1.7 if p.emphasis else 1.0))
        cx = field.left + int(field.width * p.x)
        cy = plot_top + plot_h - int(plot_h * p.y)
        dot = _oval(slide, Box(cx - d // 2, cy - d // 2, d, d))
        X.set_fill(dot, ctx.accent if p.emphasis else ctx.rule_strong)
        X.set_line(dot, None)
        X.clear_effects(dot)
        dot.shadow.inherit = False

        # The label goes on whichever side of the dot has room INSIDE THAT DOT'S OWN QUADRANT.
        # Right by preference, because a label read left-to-right away from its mark is the
        # faster read; left when the quadrant's right edge is too close.
        outward = points(SPACE["xs"]) + d // 2
        q_left = field.left if p.x <= 0.5 else field.center_x
        q_right = field.center_x if p.x <= 0.5 else field.right
        if cx + outward + label_w <= q_right:
            lb_left, align = cx + outward, "left"
        elif cx - outward - label_w >= q_left:
            lb_left, align = cx - outward - label_w, "right"
        else:
            # Neither side fits the quadrant: take the roomier one and let the label run to the
            # field edge rather than through the divider.
            right_room, left_room = q_right - cx - outward, cx - outward - q_left
            if right_room >= left_room:
                lb_left, align = cx + outward, "left"
            else:
                lb_left, align = max(q_left, cx - outward - label_w), "right"
        lb = Box(
            lb_left,
            cy - points(SCALE["caption"]),
            label_w,
            points(SCALE["caption"] * 2.4),
        )
        T.add_textbox(
            slide,
            lb,
            p.label,
            family=ctx.brand.font_minor,
            size_pt=SCALE["caption"],
            bold=p.emphasis,
            color=ctx.ink if p.emphasis else ctx.muted,
            align=align,
            anchor="middle",
            line_spacing=LINE_SPACING["footnote"],
        )
    return field
