"""
deck.py -- the public API. A closed set of slide archetypes.

THE CONTRACT THIS MODULE OFFERS: you choose an archetype and supply content. You do not choose
a layout, a colour, a size, or a position. Those are already decided, consistently, for every
slide of that kind in the deck.

That is a real restriction and it is the entire point. A caller who can position anything will,
over 37 slides, position things forty different ways -- each defensible in isolation, none
consistent with the rest. A closed archetype set means slide 4 and slide 31 of the same kind
are structurally identical without anyone having to remember to make them so.

WHEN AN ARCHETYPE REFUSES CONTENT, THAT IS THE FEATURE. `rows()` caps at five items and
`headline()` raises rather than shrinking. Both refusals are editorial findings wearing an
error message: a six-item list is a grouping that wants a level, and a headline that will not
fit is usually carrying two ideas. The fix is upstream in the writing, which is where a
designer would push it back to as well.

    from kpmg_deck import Deck, Brand
    from kpmg_deck.charts import Datum

    deck = Deck(Brand.load("kpmg"), owner="Your Name")
    deck.cover("Set up a clean workspace", "KPMG Lakehouse | September 2026")
    deck.section(1, "Set up a clean workspace",
                 "An agent that cannot find your files will invent them.")
    deck.exhibit("Three regions grew; the fourth carried the loss",
                 lambda s, ctx, box: bar_chart(s, ctx, box, data),
                 source="Source: internal analysis, 2026")
    deck.save("out.pptx")
"""

from __future__ import annotations

import datetime as _dt
from dataclasses import dataclass, field
from typing import Callable, Sequence

from pptx import Presentation
from pptx.util import Emu

from . import components as C
from . import imagery as IM
from . import oxml as X
from . import text as T
from . import theme as TH
from .canvas import GRID, Box, Grid, inches, points
from .charts import (
    Datum,
    bar_chart,
    column_chart,
    comparison,
    proportion_bar,
    ring_gauge,
    table,
)
from .components import Ctx, Ground
from .tokens import LINE_SPACING, SCALE, SPACE, Brand, contrast_ratio, tracking_pt

SlideBuilder = Callable[[object, Ctx, Box], None]

# Box allocation is measured against the metric fallback face. Allocation only needs to be
# right to within a line, and every face in the fallback chain has similar vertical metrics;
# the per-string fit check that follows uses the deck's ACTUAL face.
_ALLOC_FONT = "Arial"


# ---------------------------------------------------------------------------
# Group 3 -- the no-photograph layouts, brand book pp.112-115
# ---------------------------------------------------------------------------
#
# The book's selection rule is mechanical, p.101 verbatim: "If you have a vertical photo, use
# group 1. If you have a singular photo, use group 2. If you do not want to use a photo, use
# group 3." We hold no licensed KPMG photography, so every cover and divider this package
# builds is a Group 3 page and there is no judgment call in that.
#
# Two variants, both stated in words and drawn on p.112:
#   G3-A  10% vertical element left, 90% solid field right. All type in the field.
#   G3-B  40% wide vertical element left, 60% field right. All type INSIDE the element,
#         and the 60% is empty. The empty 60% is the design -- see group3_b's docstring.
#
# p.115 prints the pair matrix, and its rule is the one that condemned our previous deck:
# "under no circumstances can the same color be used twice in one layout". That refusal lives
# in palette.py, not here; these archetypes ask for a pair and are handed a legal one.
G3_A_ELEMENT = 0.10
G3_B_ELEMENT = 0.40

# The display measure as a fraction of the type column. See the long note in `_group3`, which
# lists the five pages it was measured on and argues why the landscape PowerPoint witness wins
# over the portrait book pages.
G3_DISPLAY_MEASURE = 0.55

# THE AA RIDER, WHICH IS OURS AND NOT KPMG'S. The book says white logo and type, full stop,
# and against the 2015 hexes that was nearly safe -- the whole 2015 set bottoms out at 3.11:1
# against white. The 2022 palette does not: Light Purple #B497FF (2.38), Pacific #00B8F5 and
# Green #00C0AE (2.29 each) put white BELOW the WCAG floor at every size, so a name-for-name
# mapping of the book's matrix onto current values produces unreadable pages that satisfy
# every rule the book states.
#
# Exactly one surface per Group 3 layout carries type, which is what makes the fix cheap: in
# G3-A the 10% element carries nothing, in G3-B the 60% field carries nothing. So the rider
# binds one colour and leaves the other free, which is where the bright three are strongest
# anyway -- the book itself sets a huge display word on Light Blue at 3.46 and never sets body
# copy there.
G3_TYPE_MIN_BODY = 4.5  # the type-bearing surface, where it carries body copy
G3_TYPE_MIN_DISPLAY = 3.0  # the type-bearing surface, display type only


def _lines(role: str, count: float = 2.0) -> int:
    """
    The height of `count` lines of type at `role`, in EMU.

    Exists because every archetype needs it and because the alternative -- a magic multiple
    like points(SCALE["h2"] * 1.8) -- is wrong in a way that is invisible until it fires. That
    expression allocates 1.8x the POINT SIZE, but a line occupies size x leading, so at h2 it
    reserves 61pt for two lines that actually need 73pt. Every two-line headline then failed
    the fit check against a box that was never big enough, and the error blamed the text.

    The 1.04 is slack for the text frame's own top and bottom insets.

    Two lines is the default deliberately: a slide headline longer than two lines is not a
    headline. If it does not fit in two, the sentence is carrying more than one idea.
    """
    return points(
        T.line_height_pt(
            _ALLOC_FONT, SCALE[role], LINE_SPACING.get(role, 1.1), bold=True
        )
        * count
        * 1.06
    )


@dataclass
class Deck:
    """
    A brand-driven presentation.

    brand        the Brand object. Data, not code -- see tokens.Brand.
    licensed     whether this deck may carry the brand owner's marks.

                 FALSE (the default) is the correct setting for anyone who is not producing
                 official material for the brand in question. It emits the palette, the type
                 scale, the grid and the editorial register -- all of which are style, and
                 style is neither owned nor restricted -- but no logo, no brand line, and no
                 corporate copyright assertion. The deck reads as belonging in the room
                 without claiming to be the organisation speaking.

                 TRUE emits the full corporate furniture: classification marker, the member
                 firm's copyright and entity sentence, and the logo if one is supplied. Set it
                 only when the deck genuinely IS the organisation's own material, in which case
                 it goes through that organisation's normal brand review like anything else.

                 The distinction is not pedantic. An unlicensed third party shipping a deck
                 with the marks on it has produced a branded artifact that never went through
                 review -- which is a rejection trigger in a way that a merely on-palette deck
                 never is. The risk runs opposite to intuition: being too on-brand is the
                 problem, not being insufficiently so.

    owner        your own name or mark, used on the cover in unlicensed mode so the deck is
                 legible as yours.
    member_firm  e.g. "KPMG LLP"; with firm_form, builds the copyright line in licensed mode.

    doc_title    the short running title in the footer, e.g. "How to become an AI engineer".
                 Present on 100% of KPMG's non-cover pages and absent from ours; see `sections`.
    sections     the section names for the top nav strip. Six or eight; a strip needs at least
                 two. Leave empty and no strip is drawn.
    nav          draw the nav strip at all. UNVERIFIED IN THE ROOM -- see components.nav_strip.
                 Set False to keep the footer and drop the strip, which is the fallback if 10pt
                 furniture turns out to be illegible from the back.
    standing_line
                 the deck's own footer line, unlicensed mode. Never KPMG's copyright.
    """

    brand: Brand
    licensed: bool = False
    owner: str | None = None
    member_firm: str | None = None
    firm_form: str | None = None
    logo_path: str | None = None
    grid: Grid = GRID
    doc_title: str | None = None
    sections: Sequence[str] = ()
    nav: bool = True
    standing_line: str | None = None
    year: int = field(default_factory=lambda: _dt.date.today().year)

    prs: Presentation = field(init=False, repr=False)
    _page: int = field(init=False, default=0, repr=False)
    _log: list[str] = field(init=False, default_factory=list, repr=False)
    _panel_side: str | None = field(init=False, default=None, repr=False)
    _dark_range: tuple[int, int] | None = field(init=False, default=None, repr=False)
    _section: int | None = field(init=False, default=None, repr=False)

    def __post_init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Emu(self.grid.slide_w)
        self.prs.slide_height = Emu(self.grid.slide_h)

        major, minor = self._resolve_fonts()
        report = TH.apply(
            self.prs,
            self.brand.theme_colors,
            major=major,
            minor=minor,
            scheme_name=self.brand.name,
        )
        self._log.append(
            f"theme: {report['color_schemes']} colour scheme(s), "
            f"{report['font_schemes']} font scheme(s), "
            f"{report['effect_styles_flattened']} effect style(s) flattened"
        )

        if self.licensed and not (self.member_firm and self.firm_form):
            raise ValueError(
                "licensed=True needs member_firm and firm_form so the copyright line can be "
                "built exactly. These vary by member firm and must not be guessed -- e.g. "
                "member_firm='KPMG LLP', firm_form='a Delaware limited liability partnership'."
            )

        warning = self.brand.provenance_warning()
        if warning:
            self._log.append("PROVENANCE: " + warning)

    def _resolve_fonts(self) -> tuple[str, str]:
        """
        Use the licensed display face only if it is genuinely installed.

        Naming a font the machine lacks is worse than not using it: PowerPoint substitutes
        silently, the substitute has different advance widths, and every text box measured
        against the declared font re-flows -- on the presenting machine, after every check has
        passed on the building one.
        """
        major = self.brand.font_major
        preferred = getattr(
            self.brand, "font_preferred_major", None
        ) or self.brand.provenance.get("font_preferred_major")
        if preferred and T.font_available(preferred, bold=True):
            major = preferred
            self._log.append(f"fonts: using licensed display face {preferred!r}")
        elif preferred:
            self._log.append(
                f"fonts: {preferred!r} not installed, falling back to {major!r} (correct behaviour)"
            )
        return major, self.brand.font_minor

    # -- infrastructure -----------------------------------------------------

    def _new(
        self, ground: Ground = Ground.LIGHT, *, section: int | None = None
    ) -> tuple[object, Ctx]:
        self._panel_side = None
        self._dark_range = None
        self._section = section
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        ctx = Ctx(self.brand, ground, self.grid)
        C.background(slide, ctx)
        self._page += 1
        return slide, ctx

    def _panel(
        self, slide, ctx: Ctx, *, side: str, fraction: float, color: str | None = None
    ) -> tuple[Box, Box]:
        """
        Draw a panel and REMEMBER WHERE IT IS, so the chrome can colour itself against it.

        The remembering is the point. A panel bleeding to y=0 runs under the nav strip and a
        panel bleeding right runs under the footer's page number; both then need reversed ink.
        Without this the gate catches it as a contrast failure at best (the page number at
        1.18:1 on five slides, which it did) and at worst it ships invisible.

        THIS PANEL'S WIDTHS AND THE GROUP 3 SPLITS ARE TWO DIFFERENT SYSTEMS AND THEY CONFLICT.
        Recorded here rather than resolved, because resolving it by picking one would throw away
        a real measurement or a real published rule:

          `C.PANEL_DEFAULT` 29.511% / `C.PANEL_WIDE` 33.770% are FREQUENCIES measured off 95
          rendered pages of the CEO-Outlook report family -- 283.3pt occurs six times across
          three decks, 324.2pt five times. Nobody recovered a rule behind them.
          `G3_A_ELEMENT` 10% / `G3_B_ELEMENT` 40% are STATED IN WORDS in the brand book,
          pp.113-114, drawn on p.112, and measured at 10/90 on eight of the book's own section
          dividers.

        Neither is 29.5% or 33.8%, and the book contains no 29.5% anything. The split is by
        GENRE: the measured widths are correct for the report-style interior page, which is what
        they were measured from and which the book's nine layouts do not cover; the Group 3
        splits are correct for covers and dividers, which the book specifies and the report
        corpus does not contain. `split()` uses these; `group3_a()`/`group3_b()` use those.
        A page that is trying to be both is a page that has not decided what it is.
        """
        pbox, rest = C.panel(slide, ctx, side=side, fraction=fraction, color=color)
        self._panel_side = side
        self._dark_range = (pbox.left, pbox.right)
        return pbox, rest

    def _furniture(self, slide, ctx: Ctx, *, chrome: bool = True) -> None:
        """
        Emit the standing per-slide furniture: nav strip, footer, and the licensed extras.

        THIS USED TO EMIT A PAGE NUMBER AND NOTHING ELSE -- 6.6pt of ink, a single digit,
        against KPMG's 845pt of document furniture on every one of 95 pages. That gap is the
        cheapest signal there is that an artifact is a slide template rather than a published
        document, and closing it costs two components and no editorial judgment.

        The unlicensed footer is NOT a reduced version of the licensed one. It carries the
        deck's own standing line and the running title; what it never carries is KPMG's
        copyright assertion or the marks.

        Called LAST by every archetype, after any panel, so a panel bleeding to y=0 sits UNDER
        the chrome rather than over it -- and `_dark_range` tells both components which of
        their own runs are now standing on colour.
        """
        if not chrome or self._page <= 1:
            # Page 1 is the cover. 0 of 4 measured covers carry a nav strip or a footer.
            return

        if self.nav and len(self.sections) >= 2:
            C.nav_strip(
                slide, ctx, self.sections, self._section, dark_range=self._dark_range
            )

        C.footer(
            slide,
            ctx,
            self.doc_title or "",
            self._page,
            legal=self.standing_line if not self.licensed else None,
            dark_range=self._dark_range,
        )

        if not self.licensed:
            return

        legal = self.brand.legal
        if legal.get("classification"):
            C.classification(slide, ctx, legal["classification"])

        template = legal.get("copyright_template")
        if template and self.member_firm and self.firm_form:
            line = template.format(
                year=self.year, member_firm=self.member_firm, firm_form=self.firm_form
            )
            g = self.grid
            T.add_textbox(
                slide,
                Box(
                    g.margin_x + inches(0.75),
                    g.slide_h - points(46),
                    int(g.slide_w * 0.62),
                    points(32),
                ),
                line,
                family=self.brand.font_minor,
                size_pt=6.5,
                color=ctx.faint,
                line_spacing=1.1,
            )

    def notes(self, slide, text: str) -> None:
        """
        Speaker notes.

        Notes are free and slides are not, which is the whole basis of the split: the CLAIM
        goes on the slide, the ARGUMENT goes here. Every content slide should have them --
        it is the pressure valve that lets a slide stay at one idea.
        """
        slide.notes_slide.notes_text_frame.text = text

    def _head(
        self,
        slide,
        ctx: Ctx,
        title: str,
        *,
        top: int | None = None,
        eyebrow: str | None = None,
        role: str = "h2",
        floor: str = "h3",
        span: int = 10,
    ) -> int:
        """
        THE FULL-WIDTH HEAD: eyebrow over a two-tone headline. Returns the y below it.

        Used by the archetypes whose content genuinely needs the full measure -- a wide chart,
        a five-across process. Everything else should use `_head_col1`, which is the layout 60%
        of real KPMG pages use.

        NO ACCENT RULE. It used to draw one under every headline and it appeared on 28 of our
        29 light pages, every one starting at exactly x = 8.1%. A mark that never varies has
        stopped being a mark: it is a template field, and it was the third-fastest tell in a
        side-by-side against real pages. KPMG's only equivalent rule is 30.91pt long, sits under
        the SECTION NUMBER, and appears on section openers only -- see `C.section_number`.

        The measure-then-place discipline stays and is load-bearing: positioning anything at
        `top + allocated_box_height` coincides with the bottom of a two-line headline and draws
        through the words, which is what the first render of this package did on four slides
        out of ten. Advance from `Placed.bottom`, always.
        """
        g = self.grid
        y = g.margin_top if top is None else top

        if eyebrow:
            C.eyebrow(slide, ctx, g.col(0, span, top=y, height=points(26)), eyebrow)
            y += points(30)

        placed = C.two_tone_headline(
            slide,
            ctx,
            g.col(0, span, top=y, height=_lines(role, 2)),
            title,
            role=role,
            on_panel=ctx.dark_ground,
            allow_shrink_to=floor,
        )
        return placed.bottom + points(SPACE["lg"])

    def _head_col1(
        self,
        slide,
        ctx: Ctx,
        title: str,
        *,
        standfirst: str | None = None,
        number: int | str | None = None,
        eyebrow: str | None = None,
        role: str = "display",
        floor: str = "h2",
        body_span: int = 6,
        divider: bool = True,
    ) -> Box:
        """
        THE PAGE ARCHITECTURE: headline in column one, content BESIDE it. Returns that region.

        This is the layout 60% of measured KPMG pages use and the one we did not have at all.
        Ours put every headline across the top at `col(0, span=10)` with the body underneath --
        which is the default slide layout of every template ever made, and it is why our pages
        read as a template with correct colours. A large headline in the left column with two
        columns of body beside it reads as a document page; the same words stacked vertically
        do not.

        Geometry, measured span-by-span off D1 p5 and TMT p5 (identical to 0.1pt):

            "01"        Bold 26.36pt #1E49E2      at (56.69, 80.70)
            rule        1.465pt #1E49E2, 30.91pt long, x 56.69 -> 87.61, at y 110.63
            headline    x 55.40, first line top 111.77, leading 66.00 (0.88 x 75)
                        line 1 #00338D, remaining lines #1E49E2 -- FLUSH LEFT, no indent
            standfirst  13pt Roman #1E49E2, x 54.82, leading 15.60, width 423
            divider     0.5pt vertical rule at x 490.22, y 127.38 -> 489.99
            body        10.5pt Light #00338D, TWO columns at x 515.08 / 719.16, 184.24 wide

        THE HEADLINE COLUMN IS col(0,6) = 413.4pt, NOT the 269pt three-column module, and that
        is a deliberate deviation with arithmetic behind it. KPMG-Bold carries 1.57x the
        characters of Arial Bold in the same width at the same visual weight (condensation
        index 0.352-0.447 against 0.582-0.653, at cap heights 15-30% smaller). At 413pt an
        Arial Bold h1 fits ~16 characters per line, which is KPMG's own line length. Chasing
        KPMG-Bold instead is not an option -- it is not licensed to us, and naming a font the
        delivery machine lacks re-flows every measured box on the presenting machine.

        Whether 413pt plus C8's tighter leading closes ENOUGH of that 1.57x gap has been
        rendered but not measured against `img/crop-kpmg-divider-headline.png`. If it has not,
        the next lever is the column, not the face.
        """
        g = self.grid
        y = g.margin_top

        if number is not None:
            y = C.section_number(slide, ctx, number) + points(SPACE["md"])
        elif eyebrow:
            C.eyebrow(slide, ctx, g.col(0, 6, top=y, height=points(26)), eyebrow)
            y += points(32)

        head_box = g.col(0, 6, top=y, height=_lines(role, 3))
        placed = C.two_tone_headline(
            slide,
            ctx,
            head_box,
            title,
            role=role,
            on_panel=ctx.dark_ground,
            allow_shrink_to=floor,
        )
        y = placed.bottom

        if standfirst:
            # THE STANDFIRST HANGS FROM THE CONTENT BOTTOM LINE, not from the headline.
            #
            # On D1 p5 the headline occupies y 111-600 and the standfirst runs y 695-900, with
            # its last line landing on the content bottom at 491.81pt (= 900/1080). The gap
            # between them is whatever is left over, and it is large. Placing the standfirst
            # directly under the headline instead -- which is the obvious thing, and what the
            # first version of this did -- leaves the bottom third of the column empty on every
            # page whose headline came out short, and the composition reads as top-heavy.
            #
            # Anchoring it to a line that never moves is also what makes two section openers
            # with different headline lengths look like siblings.
            sf_bottom = g.slide_h - g.margin_bottom
            sf_color = (
                ctx.accent if not ctx.dark_ground else ctx.brand.color("on_dark_accent")
            )
            sf_top = max(y + points(SPACE["lg"]), sf_bottom - points(150))
            C.body(
                slide,
                ctx,
                Box(head_box.left, sf_top, head_box.width, sf_bottom - sf_top),
                standfirst,
                role="lead",
                color=sf_color,
                anchor="bottom",
                strict=False,
            )

        # The region beside the headline: col(6,6), then inset by the measured 24.86pt gap
        # between the divider and the first body column. That inset makes our body block
        # [514.78, 903.31] against KPMG's measured [515.08, 903.40] -- 0.3pt.
        beside = g.col(body_span, 12 - body_span)
        if divider:
            C.column_divider(slide, ctx, beside.left, points(127.38), points(489.99))
        inset = points(24.86)
        return Box(
            beside.left + inset,
            g.margin_top
            + points(42.34),  # body top 127.38 against headline cap-top 85.04
            beside.width - inset,
            points(489.99) - g.margin_top - points(42.34),
        )

    def _body_columns(
        self,
        slide,
        ctx: Ctx,
        box: Box,
        text: str,
        *,
        columns: int = 2,
        role: str = "body",
    ) -> int:
        """
        Set prose in two columns inside a region. RETURNS THE MEASURED BOTTOM OF THE TEXT.

        The return value is the whole point and its absence was a real defect. The first version
        returned None, so callers placed the page's third element at a FRACTION of the region --
        62% down, 44% down -- and on every page whose body ran longer than the guess, the stat
        or the contrast pair was drawn straight through the prose. Ten of 37 pages shipped that
        way in the first render of this rebuild.

        This is the same measure-then-place defect `Placed` was introduced to kill, reintroduced
        by new code that did not use it. The rule is in components.Placed's docstring and it has
        no exceptions: ADVANCE FROM WHAT THE TEXT MEASURED, never from the box it was offered.

        Two columns of 184.24pt with a 19.843pt gutter is what every measured KPMG content page
        does, and it is not decoration: at 10.5pt a single 388pt measure runs to ~75 characters
        a line, which is past the point where the eye reliably finds the next line. At our 20pt
        it would be worse.

        Splits on a blank line if the text has one, otherwise at the sentence nearest the
        midpoint. It never hyphenates and never balances by character count alone -- a column
        break mid-sentence in a two-column block reads as a typesetting accident, which is
        exactly what it would be.
        """
        g = self.grid

        def _fit_role(chunks: list[str], col_w: int) -> str:
            """
            The largest body step at which every chunk fits the column. STEP DOWN, NEVER OVERFLOW.

            `strict=False` was silencing a real fit failure: a body two lines too long simply
            ran past the content bottom line and into the footer band, which the structural
            gate cannot see -- a text frame that overflows is still a valid text frame, and
            LibreOffice draws the extra lines quite happily. The scale has three body steps for
            exactly this, and dropping one is invisible where a collision with the footer is not.

            It stops at `caption` (16pt). Below that the deviation this package documents -- a
            20pt body for a 140-person room against KPMG's 10.5 -- would have been given away
            silently, so an overflow past that point is allowed to stand and be SEEN instead.
            """
            for candidate in (role, "small", "caption"):
                if all(
                    T.fit(
                        chunk,
                        Box(0, 0, col_w, box.height),
                        ctx.brand.font_minor,
                        SCALE[candidate],
                        line_spacing=LINE_SPACING["body"],
                    ).fits
                    for chunk in chunks
                ):
                    return candidate
            return "caption"

        if columns < 2:
            return C.body(
                slide, ctx, box, text, role=_fit_role([text], box.width), strict=False
            ).bottom

        paras = [p.strip() for p in text.split("\n\n") if p.strip()]
        if len(paras) >= 2:
            half = len(paras) // 2 + len(paras) % 2
            chunks = ["\n\n".join(paras[:half]), "\n\n".join(paras[half:])]
        else:
            sentences = [
                s.strip()
                for s in text.replace(". ", ".\x00").split("\x00")
                if s.strip()
            ]
            if len(sentences) < 2:
                return C.body(
                    slide,
                    ctx,
                    box,
                    text,
                    role=_fit_role([text], box.width),
                    strict=False,
                ).bottom
            target = len(text) / 2
            best, run = 1, 0
            for i, s in enumerate(sentences[:-1], start=1):
                run += len(s) + 1
                if abs(run - target) < abs(
                    sum(len(x) + 1 for x in sentences[:best]) - target
                ):
                    best = i
            chunks = [" ".join(sentences[:best]), " ".join(sentences[best:])]

        gap = g.gutter
        col_w = (box.width - gap) // 2

        role = _fit_role(chunks, col_w)

        bottom = box.top
        for i, chunk in enumerate(chunks):
            placed = C.body(
                slide,
                ctx,
                Box(box.left + i * (col_w + gap), box.top, col_w, box.height),
                chunk,
                role=role,
                strict=False,
            )
            bottom = max(bottom, placed.bottom)
        return bottom

    # -- archetypes ---------------------------------------------------------

    def cover(
        self,
        title: str,
        subtitle: str | None = None,
        *,
        eyebrow: str | None = None,
        date: str | None = None,
        ground: Ground = Ground.DARK,
    ):
        """
        SUPERSEDED. Its gradient is one FY22 does not permit. Use `window()`.

        =====================================================================================
        This raises. It is kept, not deleted, because the measurement behind it is real and
        because deleting it would delete the record of how it stopped being usable.
        =====================================================================================

        The ramp here is KPMG Blue #00338D -> Cobalt #1E49E2, read off four real KPMG report
        covers. FY22 p.46 sanctions exactly two gradients -- Purple->Cobalt and
        Pacific->Light Blue -- and says "do not create new gradients". This is not one of them,
        so a faithful reproduction of a real cover is nonetheless a page the current brand
        forbids. That is the same lesson S7 recorded about the four sampled adjacency pairs,
        arriving a second time from the other direction: SAMPLING OUTPUTS TELLS YOU WHAT
        SOMEONE DID; ONLY THE RULEBOOK TELLS YOU WHAT IS ALLOWED.

        IT USED TO FAIL, BUT NOT LIKE THIS, AND THAT IS THE POINT OF THIS CHANGE. It failed
        deep inside `set_gradient` with "00338D/1E49E2 is not a measured adjacent pair... pass
        unchecked=True if this is deliberate". A reader lands on that message believing they
        have hit an over-tight guard, and the two obvious repairs -- widen
        GRADIENT_ADJACENCY, or pass `unchecked=True` -- both re-commit the exact error the
        guard exists to prevent. A refusal that misdirects the next author is worse than no
        refusal.

        THE MEASUREMENT IS STILL TRUE AND STILL USEFUL, which is why the note below survives:
        KPMG covers measure 97% ink and the ground is a full-bleed photographic or gradient
        field, never a flat brand colour -- #1E49E2 and #00338D both measure 0.00% as flat
        fills on all four measured covers. `window()` and `window_image()` carry that finding
        forward on a gradient the book actually permits.

        We have no photography, AND THAT IS NOT A COMPROMISE KPMG WOULD RECOGNISE AS OFF-BRAND.
        In the three 2025 CEO Outlook sector reports the only image in the whole document is
        the cover, and two of ten covers in the wider corpus contain no photograph at all --
        Global Tech Report 2026 and Pulse of Fintech. Both are a flat two-stop gradient field
        with an empty left 47-60% and a rendered object breaking a window frame. That is the
        composition built here, from measurements of those two.

            1 FIELD    full bleed, horizontal two-stop gradient, 0deg
                       #00338D at x=0 -> #1E49E2 at x=960
                       [GT26 cover TL #03328E -> TR #153DC0; Pulse of Fintech #00338E -> #113FB9]
            2 WINDOW   rect(557.4, 69.8, 282.0, 402.0), aspect 1:1.4255
                       radius 0, no stroke, no shadow, axis-aligned
                       fill: horizontal two-stop #0C233C -> #102C6E
            3 STREAKS  three 15%-white hairlines top-aligned to the window, running LEFT to
                       x=302, plus one 13%-white panel from x 302 -> 557.4
            4 TYPE     staircase, three steps: x 72.0 / 119.3 / 300.8
                       eyebrow #ACEAFF, title #FFFFFF (largest), series #ACEAFF
            5 BYLINE   owner and date at x = 72.0

        THE WINDOW'S ASPECT RATIO IS NOT SETTLED AND FIVE MEASUREMENTS DISAGREE: 1.416 (~sqrt2),
        1.4295 (which statistically EXCLUDES sqrt2 at t=+3.9, fitting 10:7), 1.433, 1.405. The
        spread is method, not noise. 1.4255 sits inside every measured interval and anything in
        1.41-1.44 is indistinguishable at projection. The load-bearing facts are not in dispute:
        hard edges, zero radius, axis-aligned, top y~69-70, bottom y~467-473, and the subject
        breaking the frame. Do NOT build it at the logo's box ratio -- that was claimed from a
        raster and refuted by parsing logo.svg's own path arithmetic (1:1.37509, 8:11 to within
        0.007%); the two are 3.96% apart, which is 16pt of height at 282pt wide.

        THE COVER MARGIN IS 72.0pt, NOT the body's 56.693. Logo, staircase and tagline all sit
        on 72.0 across four covers. `grid.cover_margin` holds it.

        Licensing is unchanged and must stay that way: the logo at (72, 72) and the brand line
        at (71.9, 456.4) are emitted only when `licensed=True`. In the default mode the cover
        carries the field, the window, the staircase and the OWNER's byline -- which is the
        correct footer for a deck that is not KPMG's own material.
        """
        raise ValueError(
            "cover() is superseded and its gradient is not one FY22 permits. Its ramp is "
            "KPMG Blue #00338D -> Cobalt #1E49E2, measured off four real KPMG report covers; "
            "p.46 sanctions exactly two gradients (Purple->Cobalt, Pacific->Light Blue) and "
            "says 'do not create new gradients'.\n\n"
            "  Use window() for a type-holding cover, or window_image() for an image-holding "
            "one. Both are on a permitted ramp and both carry the same measurement that made "
            "this archetype right in the first place -- KPMG covers run 97% ink on a "
            "full-bleed field, never a flat brand colour.\n\n"
            "  DO NOT 'fix' this by widening oxml.GRADIENT_ADJACENCY or passing "
            "unchecked=True. That set was narrowed to the book's two ramps on purpose, after "
            "four pairs sampled from real member-firm documents were added here with a "
            "confident argument and all four turned out to be disallowed. Sampling outputs "
            "tells you what someone did; only the rulebook tells you what is allowed."
        )
        slide, ctx = self._new(ground)
        g = self.grid
        cm = g.cover_margin

        # 1. THE FIELD. Two stops, adjacent palette colours, horizontal. oxml.set_gradient
        #    enforces all three now.
        field = C._rect(slide, g.full, self.brand.color("canvas_brand"))
        X.set_gradient(
            field,
            [
                (0.0, self.brand.color("canvas_brand")),
                (1.0, self.brand.color("accent")),
            ],
            angle_deg=0.0,
        )

        # 2. THE WINDOW. Pinned top and bottom, sliding horizontally -- D2's sits 19.1pt left
        #    of D3/D4's for exactly that reason, so x is the give and y never is.
        win = Box(points(557.4), points(69.8), points(282.0), points(402.0))
        window = C._rect(slide, win, self.brand.color("canvas_dark"))
        X.set_gradient(
            window,
            [
                (0.0, self.brand.color("canvas_dark")),
                (1.0, self.brand.color("cover_window_far")),
            ],
            angle_deg=0.0,
        )

        # 3. THE STREAKS. Measured on two covers: three hairlines top-aligned to the window
        #    running left, plus one translucent panel from x 302 to the window's edge. They are
        #    what stops the empty left 60% reading as an empty left 60%.
        streak_left = points(302.0)
        for i in range(3):
            C._rect(
                slide,
                Box(
                    streak_left,
                    win.top + points(i * 7.4),
                    win.left - streak_left,
                    points(1.6),
                ),
                self.brand.color("on_dark"),
                alpha=0.15,
            )
        C._rect(
            slide,
            Box(
                streak_left,
                win.top + points(26),
                win.left - streak_left,
                points(154),
            ),
            self.brand.color("on_dark"),
            alpha=0.13,
        )

        # 3b. THE FRAME-BREAK. Horizontal luminous beams starting INSIDE the window's middle
        #     third and continuing PAST its right edge onto the flat field, out to x ~ 920.
        #
        #     This is the move, and without it the window is a dark rectangle sitting on a blue
        #     page -- which is what the first render of this cover was. KPMG does it on the
        #     Global Tech Report 2026 cover, and the proportion to target is the logo's own:
        #     in the KPMG mark the thing inside the frame OVERSHOOTS it by 38.9% of its own
        #     height. A device that merely touches its frame is not doing what the logo does.
        #
        #     Pacific and Light blue, because those are the two tints that read as luminous on
        #     a blue field; white would read as a scratch.
        beam_specs = (
            (0.36, 3.6, self.brand.color("accent_bright"), 1.0),
            (0.46, 1.8, self.brand.color("on_dark_accent"), 0.85),
            (0.55, 5.2, self.brand.color("accent_bright"), 0.9),
            (0.63, 1.8, self.brand.color("on_dark_accent"), 0.7),
        )
        beam_x0 = win.left + int(win.width * 0.34)
        for frac, thickness, colour, alpha in beam_specs:
            C._rect(
                slide,
                Box(
                    beam_x0,
                    win.top + int(win.height * frac),
                    points(920.0) - beam_x0,
                    points(thickness),
                ),
                colour,
                alpha=alpha,
            )

        # 4. THE STAIRCASE. Three steps, right-stepping, POSITIONAL colours -- and this is the
        #    only place in the package where the step is correct. Interior headlines are flush
        #    left; see components.two_tone_headline.
        steps_x = (cm, points(119.3), points(300.8))
        lines: list[tuple[int, str, float, str]] = []
        if eyebrow:
            lines.append(
                (steps_x[0], eyebrow, SCALE["h3"], self.brand.color("on_dark_accent"))
            )
        lines.append((steps_x[1], title, SCALE["display"], self.brand.color("on_dark")))
        if subtitle:
            lines.append(
                (steps_x[2], subtitle, SCALE["h2"], self.brand.color("on_dark_accent"))
            )

        y = points(150.0)
        for x, text_str, size, color in lines:
            box_w = points(560.0) - (x - cm)
            # THE STAIRCASE STEP IS FITTED, NOT ASSUMED. A cover title is the one string an
            # author is least likely to shorten, and at a fixed size a long one ran 76pt past
            # its own box and over the byline. Two lines is the cap: a staircase step that
            # takes three has stopped being a step.
            for candidate in sorted(
                (v for v in set(SCALE.values()) if v <= size), reverse=True
            ):
                probe = T.fit(
                    text_str,
                    Box(x, y, box_w, points(candidate * 2.8)),
                    self.brand.font_major,
                    candidate,
                    bold=True,
                    tracking_pt=tracking_pt("hero", candidate),
                    line_spacing=LINE_SPACING["hero"],
                )
                if probe.line_count <= 2:
                    size = candidate
                    break
            lh = T.line_height_pt(
                self.brand.font_major, size, LINE_SPACING["hero"], bold=True
            )
            box = Box(x, y, box_w, points(lh * 2.6))
            fit = T.fit(
                text_str,
                box,
                self.brand.font_major,
                size,
                bold=True,
                tracking_pt=tracking_pt("hero", size),
                line_spacing=LINE_SPACING["hero"],
            )
            T.add_textbox(
                slide,
                box,
                text_str,
                family=self.brand.font_major,
                size_pt=size,
                bold=True,
                color=color,
                tracking_pt=tracking_pt("hero", size),
                line_spacing=LINE_SPACING["hero"],
                anchor="top",
            )
            y += points(fit.height_pt) + points(SPACE["sm"])

        # 5. THE BYLINE. Owner and date, on the cover margin, below the staircase.
        mark = []
        if self.owner and not self.licensed:
            mark.append(self.owner)
        if date:
            mark.append(date)
        if mark:
            T.add_textbox(
                slide,
                Box(
                    cm,
                    max(y + points(SPACE["md"]), points(456.4)),
                    points(560.0),
                    _lines("body", 1.5),
                ),
                "  |  ".join(mark),
                family=self.brand.font_minor,
                size_pt=SCALE["body"],
                color=self.brand.color("on_dark_accent"),
            )

        if self.licensed and self.brand.legal.get("brand_line"):
            T.add_textbox(
                slide,
                Box(points(71.9), points(456.4), points(400), points(20)),
                self.brand.legal["brand_line"],
                family=self.brand.font_minor,
                size_pt=12.0,
                bold=True,
                color=self.brand.color("on_dark"),
            )

        self._furniture(slide, ctx, chrome=False)
        return slide

    # -- Group 3: the two no-photograph layouts -----------------------------

    def _g3_pair(
        self,
        *,
        type_bearing: str,
        element: str | None = None,
        field: str | None = None,
        body: bool = False,
        index: int | None = None,
    ) -> tuple[str, str]:
        """
        Choose the two colours for a Group 3 page. Returns (element_role, field_role).

        The pair comes from `palette.py`, which owns the p.115 matrix and the refusal that
        makes "the same colour twice" unreachable. Pass both `element` and `field` to name a
        pair (it is still validated); pass neither and the rotation picks one that has not been
        used recently, which is the book's own instruction on p.112: "try to use as many
        combinations as you can over a period of time."

        THE IMPORT IS DEFERRED AND THAT IS DELIBERATE. `palette.py` is a sibling wave's file.
        A module-level import would make every deck in the package fail to build if it were
        absent; a deferred one fails loudly, at the call, only for the caller who needed it,
        and naming explicit roles does not need it at all. It is not stubbed and there is no
        fallback: a missing palette must be an error, because inventing a pair here is exactly
        the thing the matrix exists to prevent.

        `type_bearing` says which of the two surfaces carries the type -- "field" for G3-A,
        "element" for G3-B -- AND THE TWO LAYOUTS ARE INVERTED ON THIS, WHICH IS THE ONE PLACE
        THE PALETTE'S API AND THIS ONE DO NOT LINE UP.

        `palette.pair(element, field, body=)` applies the AA rider to its SECOND argument. That
        is exactly right for G3-A, where the 90% field carries everything, and exactly wrong
        for G3-B, where the type is inside the 40% element and the 60% field is empty. So this
        method does not hand the two colours over positionally: it works in terms of the
        BEARING surface and the FREE one, feeds the bearing surface into the slot the palette
        constrains, and maps the answer back to (element, field) on the way out. Feeding them
        positionally instead would check the empty field's legibility and let an illegible
        element through, on the one layout where the element is all there is to read.

        A pair is legal in either order -- the matrix is a permutation table and the ordering
        carries only which colour is the narrow one -- so nothing is lost by the inversion.
        """
        from .palette import pair, rotate  # noqa: PLC0415 -- see docstring

        bears_field = type_bearing == "field"

        if element is None and field is None:
            # rotate() walks the matrix in the palette's own (element, field) order, so its
            # SECOND value is the type-safe one. G3-B wants that value as the element.
            walked_element, walked_field = rotate(
                self._page if index is None else index, body=body
            )
            bearing, free = walked_field, walked_element
        elif element is None or field is None:
            raise ValueError(
                "a Group 3 layout takes two colours or none. Naming one and leaving the "
                "other to chance produces a pair nobody chose; pass both, or pass neither "
                "and let the rotation walk the matrix."
            )
        else:
            bearing, free = (field, element) if bears_field else (element, field)
            free, bearing = pair(free, bearing, body=body)

        # A SECOND LOCK ON THE ONE PROPERTY THIS LAYOUT CANNOT SURVIVE LOSING. The palette has
        # already refused an illegible bearing surface; this repeats the check because the two
        # modules disagree about which surface that is, and a silent re-tiering on the other
        # side of that seam would produce white type on Pacific rather than an exception.
        floor = G3_TYPE_MIN_BODY if body else G3_TYPE_MIN_DISPLAY
        measured = contrast_ratio(
            self.brand.color("on_dark"), self.brand.color(bearing)
        )
        if measured < floor:
            raise ValueError(
                f"white type on {bearing!r} measures {measured:.2f}:1, under the {floor}:1 "
                f"floor for the surface that carries the type in this layout. The book says "
                f"white type on Group 3 without qualification, which was safe on the 2015 "
                f"hexes and is not on the 2022 ones. The bright colours belong in the "
                f"surface that carries nothing."
            )

        return (free, bearing) if bears_field else (bearing, free)

    def _logo(self, slide, left: int, top: int, *, width: int | None = None) -> int:
        """
        Place the licensed logo and return the y below it, INCLUDING its clear space.

        Returns `top` unchanged when nothing is drawn, so a caller can always advance from the
        return value. Nothing is drawn unless the deck is licensed AND a logo file was supplied
        -- the default unlicensed build carries no marks at all, and p.49 changes nothing about
        that. This method decides WHERE a logo goes, never WHETHER.

        p.49 IS THREE RULES, NOT ONE, and the callers pass the position that follows from which:

          cover slide of a presentation      upper left, ABOVE THE HEADLINE
          flat colour + NARROW element (G3-A) upper left OF THE AREA TO THE RIGHT of the element
          flat colour + WIDE element (G3-B)   upper left WITHIN THE ELEMENT

        p.50 gives the fourth: interior pages and slides, BOTTOM LEFT. Confirmed current --
        `S1-brandbook.md` measured the 2025 covers' logo at (72, 72), which is the cover margin
        this package already holds as `grid.cover_margin`.

        THE SIZE IS THE ONE UNMEASURED NUMBER HERE. The book's own G3 pages put the mark at
        14-16% of the composition width (p.100's divider, p.113's mock), but those are A4
        PORTRAIT pages: the same fraction of a 16:9 slide's width is a substantially larger
        mark against the page as a whole, because the page got wider without getting taller.
        11.5% is the reduction, and it is a judgment rather than a measurement -- there is no
        KPMG PowerPoint in the corpus to check it against. The height is never specified: it is
        read from the image file's own aspect ratio, which cannot be got wrong.
        """
        if not (self.licensed and self.logo_path):
            return top
        w = width or int(self.grid.slide_w * 0.115)
        pic = slide.shapes.add_picture(self.logo_path, left, top, width=w)
        # Clear space. The book requires "at least the minimum amount" and specifies that
        # minimum in a section this package has not measured, so half the mark's own height is
        # used -- generous rather than exact, which is the safe direction to be wrong in.
        return pic.top + pic.height + pic.height // 2

    def _g3_baseblock(
        self,
        slide,
        ctx: Ctx,
        left: int,
        width: int,
        bottom: int,
        *,
        date: str | None,
        summary: str | None,
        url: str | None,
    ) -> int:
        """
        The lower block of a Group 3 page: date, hairline, one-sentence summary, URL.
        Returns the y this block STARTS at, so the caller can prove the page did not collide.

        Bottom-anchored, because that is what the book draws and because it is what makes two
        dividers with different headline lengths read as siblings -- the same reasoning already
        written into `_head_col1`'s standfirst. p.113 annotates the space between the headline
        and this block as "lots of empty space gives the design a clean, simple aesthetic": the
        gap is the composition, so nothing may be allowed to grow into it.

        The summary is the only body copy on the page and it is clamped to the four-of-six
        measure (p.99). On G3-A the available width is 750pt and the clamp takes it to 558pt;
        on G3-B the element is 271pt wide and the clamp does nothing, which is the correct
        behaviour for a rule expressed as a maximum.
        """
        g = self.grid
        white = self.brand.color("on_dark")
        gap = points(SPACE["sm"])
        rule_h = points(1.0)

        def measure(text: str, role: str, w: int) -> int:
            return points(
                T.fit(
                    text,
                    Box(0, 0, w, points(400)),
                    self.brand.font_minor,
                    SCALE[role],
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )

        summary_w = g.max_body_width(width)
        blocks: list[tuple[str, str, int, int]] = []  # (text, role, width, height)
        if date:
            blocks.append((date, "small", width, measure(date, "small", width)))
        if summary:
            blocks.append(
                (summary, "small", summary_w, measure(summary, "small", summary_w))
            )
        if url:
            blocks.append((url, "footnote", width, measure(url, "footnote", width)))
        if not blocks:
            return bottom

        total = sum(h for _, _, _, h in blocks) + gap * (len(blocks) - 1)
        if date and summary:
            total += rule_h + gap  # the hairline sits between the date and the summary

        top = bottom - total
        y = top
        for i, (text, role, w, h) in enumerate(blocks):
            C.body(
                slide,
                ctx,
                Box(left, y, w, h),
                text,
                role=role,
                color=white,
                strict=False,
            )
            y += h + gap
            if i == 0 and date and summary:
                C._rect(slide, Box(left, y - gap // 2, points(60), rule_h), white)
                y += rule_h + gap
        return top

    def _group3(
        self,
        *,
        variant: str,
        headline: str,
        subhead: str | None,
        series: str | None,
        date: str | None,
        summary: str | None,
        url: str | None,
        element: str | None,
        field: str | None,
        numeral: str | None,
        floor: str,
        chrome: bool,
        section_index: int | None,
    ):
        """The shared build. `group3_a` and `group3_b` differ only in the four values below."""
        g = self.grid
        wide = variant == "B"
        split = G3_B_ELEMENT if wide else G3_A_ELEMENT

        element_role, field_role = self._g3_pair(
            type_bearing="element" if wide else "field",
            element=element,
            field=field,
            body=bool(summary),
        )

        slide, ctx = self._new(Ground.BRAND, section=section_index)
        elem_box = Box(0, 0, int(g.slide_w * split), g.slide_h)
        field_box = Box(elem_box.right, 0, g.slide_w - elem_box.right, g.slide_h)
        # THE RAMP REPLACES THE SPLIT; IT DOES NOT DECORATE IT. THIS WAS BUILT THE OTHER WAY
        # FIRST AND THE RENDER KILLED IT, which is the whole reason the render step exists.
        #
        # The first version kept the 10% element and ramped the 90% field. It passed every gate
        # and looked terrible: a light-purple element sitting beside a field that ramps TO light
        # purple reads as one smear with a dark blob in the middle, because the element and the
        # far stop are the same colour by construction. Nothing structural was wrong with it. It
        # was simply not a thing KPMG does, and the reason it is not is now obvious.
        #
        # Going back to the evidence: the 2026 KPMG Luxembourg dividers ramp across the WHOLE
        # page and have NO separate vertical element. The 2015 book's pages have a hard 10/90
        # split and NO ramp. These are two forms of the same two-colour idea, eleven years apart
        # -- they are ALTERNATIVES, and combining them was an invention with no page behind it.
        #
        # So: a pair with a measured ramp renders as one full-bleed gradient (the 2026 form); a
        # pair without one renders as the hard split (the 2015 form). Both are two different
        # brand colours on one page, which is the rule both eras share.
        from .palette import gradient_stops  # noqa: PLC0415 -- sibling module, see _g3_pair

        stops = gradient_stops(element_role, field_role, body=bool(summary))
        if stops is not None:
            # 0 degrees = left to right, the only angle measured in print across 34 PDFs. The
            # stop order is the AA rider extended to a ramp: the colour `pair()` certified for
            # white type sits at position 0, under the type, and the ramp runs away from it --
            # exactly as the 2026 divider sets its type over the dark end and leaves the pale
            # end empty.
            X.set_gradient(
                C._rect(slide, g.full, self.brand.color(field_role)),
                stops,
                angle_deg=0.0,
            )
        else:
            C._rect(slide, elem_box, self.brand.color(element_role))
            C._rect(slide, field_box, self.brand.color(field_role))
        self._dark_range = (0, g.slide_w)

        # The type column. p.113 and p.114 annotate the inset in inches on an A4 portrait page
        # -- 0.5in on G3-A, 0.6in on G3-B -- which is 6.0% and 7.3% of that page's width. Our
        # own measured body margin, 56.693pt, is 5.9% of a 960pt slide: the same inset arrived
        # at from a different page, so it is used rather than a converted inch.
        home = elem_box if wide else field_box
        left = home.left + g.margin_x
        width = home.width - 2 * g.margin_x

        # THE DISPLAY MEASURE, AND IT IS NARROWER THAN THE COLUMN IT SITS IN. This is the single
        # thing that separates a KPMG divider from a coloured banner: the headline sets as three
        # or four SHORT lines, not two long ones, and the resulting tall stack is what p.68 means
        # by "accentuates the vertical nature of the KPMG Font". Letting the headline run the full
        # width of the field produces a wide two-line banner -- structurally correct, generically
        # ugly, and the shape the first rebuild shipped.
        #
        # MEASURED, on the white-type extent inside the coloured field:
        #     bb p.113 G3-A mock          41% of the field width
        #     bb p.100 G3-A real divider  37%
        #     bb p.114 G3-B mock          31%
        #     LU Impact Report 2026 p40   24% of page width (landscape)
        #     KPMG India 2023 deck p6     57% of page width (landscape, PowerPoint)
        #
        # 0.55 is taken, and the choice between those numbers is the argument. The portrait book
        # pages cluster at 37-41%, but A4 portrait is a different aspect: the same measure on a
        # 960pt landscape slide would leave the headline in a column narrow enough to break short
        # words. THE WITNESS WHOSE MEDIUM MATCHES OURS IS THE 2023 KPMG INDIA DECK -- landscape,
        # authored in PowerPoint, 57% -- so that is the one followed, and 0.55 sits just inside
        # it. G3-B is exempt: its type column is already only 40% of the slide, so narrowing it
        # again would do the damage this constant exists to prevent.
        #
        # This is a DISPLAY measure and has nothing to do with `BODY_MEASURE_FRACTION` (4 of 6
        # columns, bb p.99), which governs body copy and is a legibility rule about line length.
        # Reusing that number here would be a coincidence dressed as a derivation.
        headline_width = width if wide else int(width * G3_DISPLAY_MEASURE)

        white = self.brand.color("on_dark")

        # Logo: inside the element when it is wide, in the field's upper left when it is
        # narrow. p.49, both cases. Vertically on the cover margin, which is the 1in the book
        # annotates and the (72, 72) measured on the 2025 covers.
        y = self._logo(slide, left, g.cover_margin)

        # THE SECTION NUMERAL, LARGE, ABOVE THE HEADLINE. Both current-era witnesses lead their
        # dividers with it and both make it a display element rather than a label: the 2026 KPMG
        # Luxembourg report sets "01" over "Governance" over "A culture of care", and the 2023
        # KPMG India PowerPoint deck sets "01 Everything is digital", "02 Delivering digital
        # experience" and so on. The 2015 book does not do this, so it is a current-era addition
        # rather than a brand-book rule -- which is exactly why it is worth having: it is the
        # most visible single difference between our divider and theirs.
        #
        # It is set at the subhead's weight and the headline's face, not at headline size. A
        # numeral competing with the headline for the eye would inverts the hierarchy p.68 sets
        # out, where the headline is "much larger than all other typographic elements".
        if numeral:
            h = points(
                T.fit(
                    numeral,
                    Box(0, 0, headline_width, points(200)),
                    self.brand.font_major,
                    SCALE["h3"],
                    bold=True,
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            T.add_textbox(
                slide,
                Box(left, y, headline_width, h),
                numeral,
                family=self.brand.font_major,
                size_pt=SCALE["h3"],
                bold=True,
                color=white,
                line_spacing=LINE_SPACING["body"],
            )
            y += h + points(SPACE["sm"])

        placed = C.headline(
            slide,
            ctx,
            Box(left, y, headline_width, _lines("display", 4)),
            headline,
            role="display",
            color=white,
            allow_shrink_to=floor,
        )
        y = placed.bottom + points(SPACE["md"])

        # Subhead and series line: bold over regular, directly under the headline. The book's
        # own example is "Survey outlook" / "Sector and themes" and it is the deck's series
        # identity rather than a standfirst -- it does not get the standfirst's size.
        for text, bold in ((subhead, True), (series, False)):
            if not text:
                continue
            h = points(
                T.fit(
                    text,
                    Box(0, 0, width, points(200)),
                    self.brand.font_minor,
                    SCALE["body"],
                    bold=bold,
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            T.add_textbox(
                slide,
                Box(left, y, width, h),
                text,
                family=self.brand.font_minor,
                size_pt=SCALE["body"],
                bold=bold,
                color=white,
                line_spacing=LINE_SPACING["body"],
            )
            y += h

        base_top = self._g3_baseblock(
            slide,
            ctx,
            left,
            width,
            g.slide_h - g.cover_margin,
            date=date,
            summary=summary,
            url=url,
        )

        # THE EMPTY BAND IS LOAD-BEARING, SO IT IS ASSERTED RATHER THAN HOPED FOR.
        #
        # The head stack grows DOWN from the cover margin and the base block grows UP from the
        # page's bottom margin; nothing in between negotiates. On the 750pt-wide G3-A field
        # they never meet, but the G3-B element is 271pt, where a three-line headline plus a
        # subhead plus a two-line summary is genuinely close -- and "close" here means the
        # summary is drawn through the subhead, which no structural check can see (an
        # overflowing frame is a valid frame; that is why `check_overflow` exists at all).
        #
        # It refuses instead of shrinking because p.113 annotates the gap in words -- "lots of
        # empty space gives the design a clean, simple aesthetic" -- so a page that closes it
        # has stopped being this layout. The fix is upstream, in the sentence.
        if y > base_top:
            raise ValueError(
                f"a Group 3 page has {(y - base_top) / 12700:.0f}pt more type than it has "
                f"room for: the headline block runs to {y / 12700:.0f}pt and the date/summary "
                f"block starts at {base_top / 12700:.0f}pt. Cut the headline or the summary. "
                f"The empty band between them is specified (bb p.113) -- filling it is not a "
                f"denser page, it is a different layout."
            )

        self._furniture(slide, ctx, chrome=chrome)
        return slide

    def group3_a(
        self,
        headline: str,
        *,
        subhead: str | None = None,
        series: str | None = None,
        date: str | None = None,
        summary: str | None = None,
        url: str | None = None,
        element: str | None = None,
        field: str | None = None,
        numeral: str | None = None,
        chrome: bool = False,
        section_index: int | None = None,
    ):
        """
        SUPERSEDED by `Deck.window` -- this is the 2015 system. See palette.py's header and
        S7-fy22-brandbook.md. FY22 (March 2022) has no layout groups; it has the window on one
        of two gradients. Kept for 2015-era work and for the record, not for new decks.

        G3-A: a 10% vertical element and a 90% field, in two different colours. bb p.113.

        The book's words: "One color will be in a vertical element on the left side of the
        composition, 10% in width. The other 90% will be another solid color. A white logo and
        text will be left-aligned." Every one of the eight Group 3 dividers in the book's own
        pages measures 10/90 at the sampling resolution, and no two of them repeat a pair.

        This is the archetype for SECTION DIVIDERS and for the COVER. The element carries
        nothing: it is a colour, not a column, and putting anything in it is the commonest way
        to get this layout wrong.

        THE EMPTINESS IS SPECIFIED. p.113 annotates the gap between the subhead and the date as
        "lots of empty space gives the design a clean, simple aesthetic", and p.100 -- the
        book's own Layout divider, built as a G3-A -- is one word and a logo on an otherwise
        untouched page. That is the opposite of the near-empty page this package's gates exist
        to catch, and the difference is worth being precise about: `check_rhythm` catches a
        page carrying 6% ink and no argument, while this page carries 100% ink and one word
        deliberately. A divider states the section; the argument goes on the pages after it.

        `chrome` is FALSE by default and that is a deliberate departure from every other
        archetype here. The nav strip and footer colour their ink from a single `dark_range`,
        which cannot express a page that is two different colours side by side: the footer's
        left run sits at x=56.7pt, inside the 10% element, and the element is the one surface
        the AA rider does NOT constrain. The book's own dividers carry no running chrome
        either. Pass `chrome=True` only for a pair whose element is dark.
        """
        return self._group3(
            variant="A",
            headline=headline,
            subhead=subhead,
            series=series,
            date=date,
            summary=summary,
            url=url,
            element=element,
            field=field,
            numeral=numeral,
            floor="h2",
            chrome=chrome,
            section_index=section_index,
        )

    def group3_b(
        self,
        headline: str,
        *,
        subhead: str | None = None,
        series: str | None = None,
        date: str | None = None,
        summary: str | None = None,
        url: str | None = None,
        element: str | None = None,
        field: str | None = None,
        numeral: str | None = None,
        chrome: bool = False,
        section_index: int | None = None,
    ):
        """
        SUPERSEDED by `Deck.window` -- this is the 2015 system. See palette.py's header and
        S7-fy22-brandbook.md. FY22 (March 2022) has no layout groups; it has the window on one
        of two gradients. Kept for 2015-era work and for the record, not for new decks.

        G3-B: a 40% vertical element carrying everything, and a 60% field carrying nothing.
        bb p.114.

        The book's words: "This layout consists of a wide vertical element on the left side of
        the composition, taking up 40% of the page... To the right of the vertical element,
        taking up 60% of the page is a different KPMG primary or secondary color."

        THE EMPTY 60% IS THE DESIGN, NOT A GAP TO FILL. Look at p.114: logo, headline, subhead,
        date and summary all sit inside the 40%, and the other 60% is one flat colour with
        nothing on it at all. This archetype therefore takes no content that could go there,
        which is the only reliable way to stop it being filled -- an optional right-hand
        parameter would be used, and the layout would quietly become a two-column page with a
        coloured left margin, which is a different layout that the book does not have.

        The consequence for type is real and worth stating: the type column is 271pt wide
        against G3-A's 750pt, so the same headline sets at a smaller size or across more lines.
        That is why this variant shrinks to `h3` rather than stopping at `h2`. It suits a short
        title; a long one belongs in G3-A.

        The type-bearing surface here is the ELEMENT, not the field -- inverted from G3-A -- so
        the AA rider binds the 40% and leaves the 60% free. That is the right way round: a
        bright Pacific or Light Purple 60% field is the most striking thing this package can
        put on a screen, and it is exactly where those colours are safe.
        """
        return self._group3(
            variant="B",
            headline=headline,
            subhead=subhead,
            series=series,
            date=date,
            summary=summary,
            url=url,
            element=element,
            field=field,
            numeral=numeral,
            floor="h3",
            chrome=chrome,
            section_index=section_index,
        )

    def section(
        self,
        number: int | str,
        title: str,
        standfirst: str,
        body: str,
        *,
        aside: Callable | None = None,
        section_index: int | None = None,
    ):
        """
        A section opener, built as KPMG builds one: A WHITE WORKING PAGE.

        THE 46%-WIDE COBALT PANEL THIS REPLACES WAS THE CATEGORY ERROR THAT SURVIVED THE LAST
        REBUILD. Only the hue changed then; the mistake did not. Measured across the corpus,
        KPMG's full-colour fields are reserved for FOUR things -- the cover, an attributed
        pull-quote, the methodology page, and the closing -- and all five of D1's section
        openers are white. A colour field marks a CHANGE OF VOICE, never a change of section.
        We used it for section boundaries and nothing else, which is exactly inverted.

        The failure case is on disk: `img/evidence-rebuilt-divider-0407.png` is a 46% panel
        that is 76% EMPTY, with the headline collapsed to ~24pt cap against KPMG's 51pt. A
        panel with nothing in it is not a composition, and the more colour it has the more
        obviously so.

        What actually carries a section opener is TYPE AS THE GRAPHIC: number, rule, a giant
        two-tone headline in the left column, a Cobalt standfirst under it, two columns of the
        section's real argument beside it, and about 60% of the page left white. No ornament at
        all. `body` is REQUIRED for that reason -- a section opener with nothing to say is the
        near-empty page this whole rebuild exists to eliminate, and making the argument optional
        is what made it easy to omit.

        `aside` is an optional builder `(slide, ctx, box)` for a right-hand statistic or ring,
        drawn below the body columns. It is the third element on a page that already has two.
        """
        slide, ctx = self._new(
            Ground.LIGHT,
            section=section_index
            if section_index is not None
            else self._index_of(number),
        )
        g = self.grid

        beside = self._head_col1(
            slide, ctx, title, standfirst=standfirst, number=number
        )

        # The aside starts below what the BODY MEASURED, never at a fraction of the region.
        # A fraction drew the section's minute-count straight through its own argument on every
        # THE ASIDE SITS BESIDE THE BODY, NOT UNDER IT.
        #
        # Stacking them vertically was the first attempt and every opener dropped its figure:
        # the two-column body fills all 363pt of the region, so there was never anything left.
        # That is not a content problem to solve by cutting the argument -- KPMG's own answer
        # is a third COLUMN. D1 p4 runs headline / stats / rings / rings across the page and
        # D3 p12 runs headline / body / figure, both with the figure in its own narrow column.
        #
        # So with an aside the body goes to one column and the figure takes the right third.
        if aside is None:
            self._body_columns(slide, ctx, beside, body)
        else:
            aside_w = int(beside.width * 0.34)
            body_w = beside.width - aside_w - g.gutter
            self._body_columns(
                slide,
                ctx,
                Box(beside.left, beside.top, body_w, beside.height),
                body,
                columns=1,
            )
            aside(
                slide,
                ctx,
                Box(
                    beside.left + body_w + g.gutter,
                    beside.top,
                    aside_w,
                    beside.height,
                ),
            )

        self._furniture(slide, ctx)
        return slide

    def _index_of(self, number: int | str) -> int | None:
        try:
            i = int(number) - 1
        except (TypeError, ValueError):
            return None
        return i if 0 <= i < len(self.sections) else None

    # -- the three-element content archetypes -------------------------------
    #
    # C1 of the build spec, and the item everything else waits on: EVERY PAGE CARRIES THE CLAIM
    # PLUS TWO SUPPORTING ELEMENTS. Our previous deck called `statement(claim)` with no second
    # argument thirty times, producing a median of 69 characters per page against KPMG's
    # 2,275-2,422, and 28 of 37 pages under 8% ink where ~2% of 62 real KPMG pages are (and
    # their one instance is a back cover).
    #
    # The target is THREE ELEMENTS, NOT MORE WORDS. Copying KPMG's 2,400-character density
    # would be wrong: that is right for a desk-read InDesign report and wrong for 140 people in
    # a room. The second and third elements are drawn from a closed list -- a chart or small
    # exhibit, a large stat figure with a one-line gloss, an attributed pull-quote, a
    # two-column contrast, a labelled diagram, a worked terminal example.
    #
    # These archetypes exist so that the three-element page is the CHEAP path. An archetype
    # that requires its supporting content is the only reliable way to stop it being omitted.

    def window(
        self,
        headline: str,
        *,
        subhead: str | None = None,
        numeral: str | None = None,
        support_ground: bool = False,
        section_index: int | None = None,
    ):
        """
        A gradient ground carrying THE WINDOW, with the headline set inside it. bb FY22 pp.46, 72.

        THIS IS THE DEVICE THE DECK WAS MISSING, and it is the answer to "still just a blue
        theme". The FY22 brand guidelines (March 2022) replace the 2015 book's nine-layout system
        with a graphic motif called **the window** -- a rectangle derived from the KPMG logo's
        four boxes -- used as a type-holding or image-holding shape on a gradient ground. Every
        KPMG PowerPoint slide printed on p.113 that carries display type uses it.

        The book's own words on p.72: *"our window allows us to highlight the insight/opportunity
        by guiding the viewer's focus to the action... it is not intended to be used decoratively
        or manipulated to create alternate shapes."* So it is a device with a job, not a frame.

        GEOMETRY IS STATED, NOT INFERRED. p.72: *"size proportions are 7:10 for vertical and 10:7
        for horizontal."* A 16:9 slide takes the horizontal window. Square corners -- the shape
        comes from the logo's boxes, which are square.

        COLOUR IS STATED TOO, AND IT IS NARROW. p.46 gives exactly two gradients and forbids any
        other: PRIMARY Purple->Cobalt, SUPPORT Pacific->Light Blue, stops at 0% and 100%,
        midpoint 50%, 0 degree angle, linear never radial. The two example compositions on that
        page put ONE gradient on the ground and THE OTHER in the window, which is what this
        archetype does -- `support_ground` picks which way round.

        TYPE COLOUR FOLLOWS THE WINDOW, NOT THE GROUND, because the type sits inside the window.
        White on the Purple/Cobalt window measures 6.76-7.01:1 across the ramp; KPMG Blue on the
        Pacific/Light Blue window measures 4.93-8.63:1. Both clear 4.5:1 at every point, so the
        AA rider is satisfied along the whole ramp rather than at one end -- which is what a
        gradient needs and what the flat-field rider could not express.
        """
        g = self.grid
        slide, ctx = self._new(Ground.BRAND, section=section_index)

        ground = X.GRADIENT_SUPPORT if support_ground else X.GRADIENT_PRIMARY
        win_stops = X.GRADIENT_PRIMARY if support_ground else X.GRADIENT_SUPPORT
        on_window = (
            self.brand.color("on_dark") if support_ground else self.brand.color("ink")
        )

        X._noop = (
            None  # (keeps linters quiet about the module-level import being "unused")
        )
        X.set_gradient(
            C._rect(slide, g.full, ground[0][1]),
            list(ground),
            angle_deg=X.GRADIENT_ANGLE_DEG,
        )
        self._dark_range = (0, g.slide_w)

        # GEOMETRY IS DENOMINATED IN LOGO HEIGHTS, AND SIZED OFF HEIGHT, NOT WIDTH.
        #
        # The first version of this sized the window at 62% of the CONTENT WIDTH and hung it near
        # the top of the left margin. Both were wrong, and the width one was wrong in the most
        # deceptive way: 62% of width happens to land within a point of the book's own horizontal
        # window, so it looked right and could not produce a VERTICAL window at all -- the 7:10
        # form, which is equally sanctioned and better for a short headline.
        #
        # EVERY STATED SIZE RULE IN THE BOOK IS A HEIGHT, and every one is expressed as a multiple
        # of the KPMG logo's height rather than as a page fraction (bb FY22 pp.72-89):
        #
        #     layout margin / minimum window-to-edge clearance   1x logo height
        #     window-to-logo gap                                 2x logo height, 1x when tight
        #     type inset on every edge the type touches          0.75x logo height
        #     headline cap height                                ~2x logo height
        #
        # The logo's own height is NOT given in that section, which makes all of it undefined
        # until you have it. Ours is measured: 89.6 x 36.0pt at (72, 72) on four current KPMG
        # covers. 36 / 540 = 6.7% of slide height, against 6.2-6.5% measured independently off
        # the book's own window examples -- two different sources, one point apart, so the unit
        # is sound.
        #
        # PLACEMENT: p.79 gives exactly two positions -- to the RIGHT of the logo, or BELOW it.
        # The old left-margin/near-top placement was neither, and sat flush to the live-area edge,
        # which p.95 bans outright. This takes "below the logo", which keeps the type left-aligned
        # on the same margin as the rest of the deck.
        logo_h = points(36.0)  # measured, 89.6 x 36.0pt at (72,72) on four KPMG covers
        win_top = (
            g.cover_margin + logo_h + logo_h
        )  # logo line + the tight-layout 1x gap
        win_bottom = g.slide_h - logo_h  # 1x clearance from the live-area edge
        win_h = win_bottom - win_top
        win_w = int(win_h * 10 / 7)
        win = Box(g.cover_margin, win_top, win_w, win_h)
        X.set_gradient(
            C._rect(slide, win, win_stops[0][1]),
            list(win_stops),
            angle_deg=X.GRADIENT_ANGLE_DEG,
        )

        inner = win.inset(int(logo_h * 0.75))
        y = inner.top

        if numeral:
            T.add_textbox(
                slide,
                Box(inner.left, y, inner.width, points(30)),
                numeral,
                family=self.brand.font_minor,
                size_pt=SCALE["body"],
                bold=True,
                color=on_window,
                line_spacing=LINE_SPACING["body"],
            )
            y += points(34)

        # Measured and placed directly rather than through `C.headline`, for the same reason
        # `blocks` does it: the component asserts legibility against the CONTEXT's ground, and
        # this context describes the gradient GROUND while the type sits inside the WINDOW. It
        # would read KPMG Blue on KPMG Blue and refuse, correctly, being unable to see the
        # window. The legibility guarantee is not lost -- it is discharged in the docstring
        # above, against both ends of both sanctioned ramps, which is stricter than the
        # single-colour check this bypasses.
        # RESERVE THE SUBHEAD'S SPACE BEFORE SIZING THE HEADLINE, or the headline expands into
        # it and the subhead is placed into whatever is left -- which is routinely nothing.
        # Measured before this was fixed: the shortfall moved PERVERSELY with headline length,
        # because a longer headline trips the shrink loop below and ends up occupying LESS total
        # height than a short one set at full size. So "shorten the headline" made the subhead
        # fit less often, which is the opposite of the advice the failure invites. Reserving the
        # space up front makes the two independent and the shrink loop's job well-posed.
        sub_reserve = 0
        if subhead:
            sub_reserve = points(
                T.fit(
                    subhead,
                    Box(inner.left, y, inner.width, inner.height),
                    self.brand.font_minor,
                    SCALE["small"],
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            ) + points(SPACE["sm"])
        head_box = Box(inner.left, y, inner.width, inner.bottom - y - sub_reserve)
        # THE HEADLINE IS SIZED FROM THE LOGO, NOT FROM THE TYPE SCALE, and it starts far above
        # the scale's ceiling. bb FY22 puts the headline cap height at ~2x logo height; Arial's
        # cap is 0.716 of its em, so 2 x 36pt of cap wants roughly a 100pt font. The deck's own
        # `display` role tops out at 60pt, which is inside the book's STATED 48-60pt desktop band
        # but well under what its own slides MEASURE -- 74-96pt, found independently by a second
        # reader off the p.113 gallery. Two sources agree the stated band understates the practice.
        #
        # Starting high and shrinking is what makes the headline FILL the window the way KPMG's
        # does. Starting at 60 left two lines floating in a box sized for four, which reads as a
        # box with some words in it rather than as a type-holding shape.
        size = max(SCALE["display"], int(logo_h / 12700 * 2 / 0.716))
        while (
            size > SCALE["h3"]
            and not T.fit(
                headline,
                head_box,
                self.brand.font_major,
                size,
                bold=True,
                line_spacing=LINE_SPACING["display"],
            ).fits
        ):
            size -= 2
        fitted = T.fit(
            headline,
            head_box,
            self.brand.font_major,
            size,
            bold=True,
            line_spacing=LINE_SPACING["display"],
        )
        T.add_textbox(
            slide,
            head_box.resize(height=points(fitted.height_pt)),
            headline,
            family=self.brand.font_major,
            size_pt=size,
            bold=True,
            color=on_window,
            line_spacing=LINE_SPACING["display"],
        )
        placed = head_box.resize(height=points(fitted.height_pt))

        if subhead:
            # THE HEADLINE IS SIZED TO FILL THE WINDOW (see the shrink loop above), so whatever
            # it leaves behind is what the subhead gets -- and at four lines of display type that
            # remainder can be NEGATIVE. Subtracting to a negative height used to produce a text
            # box whose bottom sits above its own top: valid XML, invisible to every structural
            # check except the geometry one, and drawn by PowerPoint straight over the window's
            # lower edge. Refuse it instead, the same way `headline()` refuses a headline that
            # will not fit. A subhead that does not fit is a finding about the headline.
            sub_top = placed.bottom + points(SPACE["sm"])
            avail = inner.bottom - sub_top
            sub_box = Box(inner.left, sub_top, inner.width, avail)
            sub_fit = (
                T.fit(
                    subhead,
                    sub_box,
                    self.brand.font_minor,
                    SCALE["small"],
                    line_spacing=LINE_SPACING["body"],
                )
                if avail > 0
                else None
            )
            if sub_fit is None or not sub_fit.fits:
                short = (
                    abs(avail) / 12700
                    if avail <= 0
                    else sub_fit.height_pt - avail / 12700
                )
                raise ValueError(
                    f"window(): the headline leaves no room for the subhead -- short by "
                    f"{short:.1f}pt. The headline is sized to fill the window, so this is a "
                    f"finding about the headline rather than about the subhead: cut the "
                    f"headline (real KPMG slide headlines run 2-12 words, median 3), shorten "
                    f"the subhead, or drop the subhead. Shrinking the headline is not offered "
                    f"-- it would break the size ratio that makes the window read as a "
                    f"type-holding shape.\n  headline: {headline!r}\n  subhead:  {subhead!r}"
                )
            T.add_textbox(
                slide,
                sub_box.resize(height=points(sub_fit.height_pt)),
                subhead,
                family=self.brand.font_minor,
                size_pt=SCALE["small"],
                color=on_window,
                line_spacing=LINE_SPACING["body"],
            )

        self._furniture(slide, ctx, chrome=False)
        return slide

    def blocks(
        self,
        text: str,
        standfirst: str,
        items: Sequence[tuple[str, str]],
        *,
        eyebrow: str | None = None,
        section_index: int | None = None,
    ):
        """
        A claim, then two or three full-width COLOUR BLOCKS carrying the body copy.

        THIS ARCHETYPE EXISTS BECAUSE OF A MEASUREMENT, AND IT IS THE ONE THAT ANSWERS "STILL
        JUST A BLUE THEME". Across the whole deck, ours measured 54% white against the KPMG
        India 2023 PowerPoint deck's 23%. The hue ratio was never the problem -- our purple
        share, 17.6% of all coloured pixels, sits BETWEEN the two current KPMG documents
        (23.3% and 13.7%). The problem is where the colour lives.

        KPMG PUTS BODY COPY INSIDE COLOURED BLOCKS. WE PUT BODY COPY ON WHITE. Look at
        `sources/kpmg2026/in-ux-2023/p-07.png`: a short white intro, then three full-width
        blocks stacked -- dark navy, cobalt, purple -- each carrying white body text and an
        illustration, each a different colour from the one above it. That is the brand book's
        two-colour rule generalised from the page to the paragraph.

        ON A WHITE PAGE THE ONLY COLOUR IS THE TYPE, and blue type on white is what reads as a
        generic corporate template no matter which blue it is. Fixing the seven Group 3 pages
        left thirty-six white pages untouched, and thirty-six is the deck.

        THE BLOCK COLOURS ARE THE AA RIDER, NOT A PREFERENCE. A block carries body copy, so it
        must clear 4.5:1 against white -- which is exactly `palette.BODY_FIELDS`: KPMG Blue,
        Cobalt, Dark blue, Deep purple, Purple. Five colours, and the KPMG page uses three of
        them. Pacific and Light purple are excluded here for the same reason they cannot be a
        Group 3 field: white body text on them is illegal at any size.

        Consecutive blocks never repeat a colour, which is the p.115 rule applied down the page
        rather than across it. `section_index` offsets the starting colour so consecutive
        SECTIONS do not all open with the same block.
        """
        from .palette import BODY_FIELDS  # noqa: PLC0415 -- sibling module, see _g3_pair

        if not 2 <= len(items) <= 3:
            raise ValueError(
                f"a block page takes two or three blocks, got {len(items)}. One block is a "
                f"panel and belongs on `split`; four will not carry legible body copy in the "
                f"space a 16:9 slide has below a headline."
            )

        slide, ctx = self._new(Ground.LIGHT, section=section_index)
        g = self.grid
        y = self._head(slide, ctx, text, eyebrow=eyebrow, span=8)

        if standfirst:
            # The standfirst sits directly under the headline here, NOT hung from the content
            # bottom line the way `_head_col1` hangs it. That anchoring exists to stop a short
            # headline leaving the bottom of a tall single column empty; on this page the blocks
            # fill everything below, so hanging it would open a gap in the middle instead.
            sf = C.body(
                slide,
                ctx,
                g.col(0, 7, top=y, height=_lines("body", 2)),
                standfirst,
                role="lead",
                color=ctx.accent,
                strict=False,
            )
            y = sf.bottom + points(SPACE["md"])

        region = Box(g.content.left, y, g.content.width, g.content.bottom - y)
        gap = points(SPACE["sm"])
        block_h = (region.height - gap * (len(items) - 1)) // len(items)

        start = 0 if section_index is None else section_index
        white = self.brand.color("on_dark")

        for i, (label, body) in enumerate(items):
            role = BODY_FIELDS[(start + i) % len(BODY_FIELDS)]
            top = region.top + i * (block_h + gap)
            box = Box(region.left, top, region.width, block_h)
            C._rect(slide, box, self.brand.color(role))

            # Asymmetric inset: full padding at the sides, tighter top and bottom. Side padding
            # is what keeps the measure off the block edge and it is doing real work; vertical
            # padding is only air, and at `md` both ends it cost 6pt of overflow at the 14pt
            # floor -- i.e. it was spending legibility on whitespace.
            inner = box.inset(points(SPACE["md"]), points(SPACE["sm"]))
            lh = points(
                T.fit(
                    label,
                    Box(0, 0, inner.width, points(120)),
                    self.brand.font_minor,
                    SCALE["body"],
                    bold=True,
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            T.add_textbox(
                slide,
                Box(inner.left, inner.top, inner.width, lh),
                label,
                family=self.brand.font_minor,
                size_pt=SCALE["body"],
                bold=True,
                color=white,
                line_spacing=LINE_SPACING["body"],
            )
            # T.fit + T.add_textbox rather than C.body, and the reason is not convenience.
            # `C.body` asserts legibility against the CONTEXT's ground, and this page's context
            # is the white page -- so it reads white-on-white and refuses, correctly, because it
            # cannot see the block underneath. This is the same path `panel_builder` takes for
            # the same reason. The legibility guarantee is not lost, it moves: the block colour
            # comes from `BODY_FIELDS`, every member of which is >= 4.5:1 against white by
            # construction, and `verify.check_pairs` re-checks it on the built file.
            body_box = Box(
                inner.left,
                inner.top + lh + points(SPACE["xs"]),
                inner.width,
                inner.height - lh - points(SPACE["xs"]),
            )
            # THE BLOCK'S HEIGHT IS WHATEVER THE HEADLINE AND STANDFIRST LEFT OVER, DIVIDED BY
            # THE BLOCK COUNT, so a tall headline can drive it below the label's own line height
            # and this subtraction goes NEGATIVE. That produced a text box whose bottom sat above
            # its top: valid XML, invisible to every check but the geometry one, and drawn by
            # PowerPoint over the block beneath it -- white on white, so it read as a bar with a
            # smear through it. Refuse it here, the same way the 14pt floor below refuses type
            # that is too small, and for the same reason: the fix is fewer words, never a
            # geometry the renderer has to interpret.
            if body_box.height <= 0:
                raise ValueError(
                    f"blocks(): block {i + 1} of {len(items)} has no room for its body -- the "
                    f"headline and standfirst leave {body_box.height / 12700:.1f}pt after the "
                    f"label. Use fewer blocks, a shorter headline (real KPMG slide headlines "
                    f"run 2-12 words, median 3), or a shorter standfirst.\n"
                    f"  label: {label!r}"
                )
            # 14pt is the FLOOR, not 12: `verify.check_type_floor` calls anything under 14
            # illegible for a large room, and this deck is shown to 140 people. If the body
            # cannot fit at 14 the block is overfull and the fix is fewer words or fewer
            # blocks -- not smaller type, which just moves the failure somewhere no gate looks.
            size = SCALE["small"]
            while (
                size > 14
                and not T.fit(
                    body,
                    body_box,
                    self.brand.font_minor,
                    size,
                    line_spacing=LINE_SPACING["body"],
                ).fits
            ):
                size -= 1
            # THE LOOP ABOVE STOPS AT THE FLOOR AND USED TO DRAW ANYWAY, which made the comment
            # it sits under false: it promises that a body which cannot fit at 14pt is an
            # overfull block, and then the code shipped the overfull block. What that looks like
            # is not a missing line -- it is a line SLICED IN HALF by the block's lower edge,
            # because the frame's last line falls outside the coloured rectangle while staying
            # inside its own text box. `check_overflow` measures the text box and so reports
            # nothing; only the render shows it. Assert the promise instead.
            final = T.fit(
                body,
                body_box,
                self.brand.font_minor,
                size,
                line_spacing=LINE_SPACING["body"],
            )
            if not final.fits:
                raise ValueError(
                    f"blocks(): block {i + 1} of {len(items)} is overfull -- its body needs "
                    f"{final.height_pt:.0f}pt at the {size}pt floor and the block gives it "
                    f"{body_box.height / 12700:.0f}pt. Cut the body, use fewer blocks, or "
                    f"shorten the headline. Type is not reduced below 14pt: that is the "
                    f"legibility floor for a large room.\n  label: {label!r}\n  body:  "
                    f"{body!r}"
                )
            T.add_textbox(
                slide,
                body_box,
                body,
                family=self.brand.font_minor,
                size_pt=size,
                color=white,
                line_spacing=LINE_SPACING["body"],
            )

        self._furniture(slide, ctx)
        return slide

    def statement(
        self,
        text: str,
        standfirst: str,
        body: str,
        *,
        support: Callable | None = None,
        eyebrow: str | None = None,
        section_index: int | None = None,
    ):
        """
        A claim, with its argument beside it and one supporting element under that.

        THE DENSE WHITE PAGE -- 46-53.6% of KPMG's corpus and the target for 18 of our 37. It
        is the same architecture as a section opener without the number: two-tone headline and
        Cobalt standfirst in column one, body beside it, a supporting element below the body.

        `standfirst` and `body` are REQUIRED, and that is the whole mechanism. The previous
        signature was `statement(text, support=None)`, so a claim on its own was not merely
        possible but the shortest thing to type -- and it was typed thirty times. Making the
        second and third elements positional arguments does not improve anyone's judgment; it
        just makes the near-empty page impossible to reach by accident.

        `support` is a builder `(slide, ctx, box)`: a contrast pair, a code block, a stat, a
        small chart. Pass `None` only when the body genuinely runs long enough to fill the
        region, which is rarer than it feels.
        """
        slide, ctx = self._new(Ground.LIGHT, section=section_index)

        beside = self._head_col1(
            slide, ctx, text, standfirst=standfirst, eyebrow=eyebrow
        )

        # THE SUPPORT ELEMENT'S ROOM IS RESERVED BEFORE THE BODY IS SET, and this is the third
        # arrangement of these two blocks. The first placed the support at a FRACTION of the
        # region and drew it through the prose. The second measured the body first and gave the
        # support whatever was left -- which is collision-free but produced 20pt ring gauges
        # with their own labels clipped, because "whatever was left" was 60pt.
        #
        # Neither works, because both let one block decide before knowing about the other. So
        # the body is capped at 45% of the region and fitted DOWN into that cap, and the
        # support gets a floor of 150pt. An element small enough to be illegible is not a
        # cheaper version of that element; it is a worse page than not having it.
        #
        # Two columns, not one, and this was inverted at first: a single column is TALLER for
        # the same words, so the page with something to put underneath was the one given the
        # tallest possible body.
        if support is None:
            self._body_columns(slide, ctx, beside, body, columns=2)
        else:
            SUPPORT_MIN = points(168)
            cap = min(int(beside.height * 0.45), beside.height - SUPPORT_MIN)
            body_bottom = self._body_columns(
                slide,
                ctx,
                Box(beside.left, beside.top, beside.width, cap),
                body,
                columns=2,
            )
            top = max(body_bottom, beside.top + cap) + points(SPACE["lg"])
            support(
                slide, ctx, Box(beside.left, top, beside.width, beside.bottom - top)
            )

        self._furniture(slide, ctx)
        return slide

    def split(
        self,
        headline: str,
        standfirst: str,
        body: str | None = None,
        *,
        panel_build: Callable,
        body_build: Callable | None = None,
        side: str | None = None,
        fraction: float = C.PANEL_DEFAULT,
        eyebrow: str | None = None,
        section_index: int | None = None,
    ):
        """
        THE WORKHORSE: white content page plus ONE full-bleed colour panel. 29.8% of the corpus.

        Headline, standfirst and body on white in column one and two; one Cobalt panel bleeding
        off a vertical edge carrying a quote, a statistic, or two or three key points.

        `side` alternates automatically so consecutive panels mirror -- which is what the real
        decks do and what stops a long deck feeling static. The alternation logic was already
        right on `section()`; it belongs here now, because `section()` no longer has a panel.
        The corpus default is RIGHT (17 of 24 non-flood panels bleed right, 7 left).

        `panel_build(slide, panel_ctx, box)` draws into the panel's inner region, which is
        already inset by the measured 56.69pt on every edge including the bleeding one.

        `body_build` replaces the prose body with a builder -- a chart, a contrast, a code
        block. It exists so that a page can be a split AND an exhibit, which is what the deck's
        second page has to be: rule 1 of the measured sequencing is that colour arrives by page
        2 in every one of D1-D4, and an exhibit that had to be an all-white page would push it
        to page 3.
        """
        slide, ctx = self._new(Ground.LIGHT, section=section_index)
        g = self.grid

        if side is None:
            side = "left" if (self._page % 2 == 0) else "right"

        panel_ctx = Ctx(
            self.brand, Ground.BRAND, g
        )  # legal foregrounds for a colour field
        pbox, _ = self._panel(
            slide,
            panel_ctx,
            side=side,
            fraction=fraction,
            color=self.brand.color("accent"),
        )

        # The white region is everything the panel does not occupy, and the headline column is
        # sized against THAT rather than against the page -- otherwise a 33.8% panel overlaps a
        # headline placed on the full-page col(0,6).
        text_left = g.margin_x if side == "right" else pbox.right + g.margin_x
        text_right = (
            (pbox.left - g.margin_x) if side == "right" else (g.slide_w - g.margin_x)
        )
        avail = text_right - text_left

        y = g.margin_top
        if eyebrow:
            C.eyebrow(slide, ctx, Box(text_left, y, avail, points(26)), eyebrow)
            y += points(32)

        bottom = points(489.99)

        if body_build is not None:
            # THE EXHIBIT VARIANT: headline across the full white measure, figure beneath it.
            #
            # Splitting the white area into a headline column and a body column is right for
            # prose and wrong for a figure. On the house 283.3pt panel the body column is 225pt
            # -- narrower than a bar chart's own label gutter, so `bar_chart` computed a
            # NEGATIVE track width and emitted six shapes at -402,839 EMU. They rendered as
            # nothing and the gate caught them as non-positive geometry, which is the one thing
            # standing between that and a page that is silently missing its exhibit.
            #
            # A figure needs the measure, so it gets it, and the headline goes above.
            placed = C.two_tone_headline(
                slide,
                ctx,
                Box(text_left, y, avail, _lines("h2", 3)),
                headline,
                role="h2",
                allow_shrink_to="h3",
            )
            sf_top = placed.bottom + points(SPACE["sm"])
            sf_h = points(
                T.fit(
                    standfirst,
                    Box(text_left, sf_top, avail, points(400)),
                    self.brand.font_minor,
                    SCALE["lead"],
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            C.body(
                slide,
                ctx,
                Box(text_left, sf_top, avail, sf_h),
                standfirst,
                role="lead",
                color=ctx.accent,
                strict=False,
            )
            body_build(
                slide,
                ctx,
                Box(
                    text_left,
                    sf_top + sf_h + points(SPACE["lg"]),
                    avail,
                    bottom - sf_top - sf_h - points(SPACE["lg"]),
                ),
            )
        else:
            # 0.52, not 0.60. Widening the headline column to 0.60 fixed a headline that
            # would not fit and starved the body column doing it -- 205pt of measure, in which
            # a 350-character argument runs 265pt past the bottom of the page at every size
            # down to the floor. `two_tone_headline` searches five line counts and four sizes;
            # it does not need the extra width, and the body does.
            head_w = int(avail * 0.52)
            placed = C.two_tone_headline(
                slide,
                ctx,
                Box(text_left, y, head_w, _lines("h1", 3)),
                headline,
                role="h1",
                allow_shrink_to="h3",
            )

            # Same rule as `_head_col1`: the standfirst hangs from the content bottom line, so
            # two split pages with different headline lengths still read as siblings.
            sf_top = max(placed.bottom + points(SPACE["lg"]), bottom - points(140))
            C.body(
                slide,
                ctx,
                Box(text_left, sf_top, head_w, bottom - sf_top),
                standfirst,
                role="lead",
                color=ctx.accent,
                anchor="bottom",
                strict=False,
            )

            body_left = text_left + head_w + g.gutter
            C.column_divider(
                slide, ctx, body_left - g.gutter // 2, points(127.38), bottom
            )
            if body:
                C.body(
                    slide,
                    ctx,
                    Box(
                        body_left,
                        points(127.38),
                        text_right - body_left,
                        bottom - points(127.38),
                    ),
                    body,
                    strict=False,
                )

        panel_build(slide, panel_ctx, C.panel_inner(panel_ctx, pbox))

        self._furniture(slide, ctx)
        return slide

    def quote_panel(
        self,
        quote: str,
        name: str,
        role: str,
        *,
        side: str = "left",
        headline: str | None = None,
        support: Callable | None = None,
        section_index: int | None = None,
    ):
        """
        An attributed pull-quote on a wide Cobalt panel. Measured complete on D1 p6 and TMT p6.

            panel        #1E49E2  rect(0, 0, 558.4, 540)   full bleed left/top/bottom
            quote glyph  #00B8F5  bbox 60.02 x 49.41pt at (56.8, 158.2)   <- PACIFIC, not white
            quote        #FFFFFF  Bold 14 / 17pt, x 56.8, w 423.8
            name         #FFFFFF  Bold 12 / 14.4
            role         #FFFFFF  Light 12 / 14.4
            closing mark typed as a " at the end of the text, never a second glyph

        THIS IS ONE OF THE FOUR PLACES A COLOUR FIELD IS CORRECT, because a quotation IS a
        change of voice -- literally, someone else is speaking. That is the rule the section
        divider was breaking.

        `support` fills the white side: a chart, a stat, a contrast. Without it the page is a
        quote and a lot of white, which is the density defect wearing better clothes.
        """
        slide, ctx = self._new(Ground.LIGHT, section=section_index)
        g = self.grid

        panel_ctx = Ctx(self.brand, Ground.BRAND, g)
        pbox, rest = self._panel(
            slide,
            panel_ctx,
            side=side,
            fraction=C.PANEL_QUOTE,
            color=self.brand.color("accent"),
        )

        inner = C.panel_inner(panel_ctx, pbox, top=points(158.2))
        C.quote_glyph(slide, panel_ctx, inner.left, points(158.2))

        y = points(158.2) + points(88)
        qbox = Box(inner.left, y, min(inner.width, points(423.8)), points(230))
        qfit = T.fit(
            f"“{quote}”",
            qbox,
            self.brand.font_minor,
            SCALE["lead"],
            bold=True,
            line_spacing=LINE_SPACING["h3"],
        )
        T.add_textbox(
            slide,
            qbox,
            f"“{quote}”",
            family=self.brand.font_minor,
            size_pt=SCALE["lead"],
            bold=True,
            color=self.brand.color("on_dark"),
            line_spacing=LINE_SPACING["h3"],
            anchor="top",
        )

        att_top = y + points(qfit.height_pt) + points(SPACE["md"])
        T.add_textbox(
            slide,
            Box(inner.left, att_top, qbox.width, points(30)),
            name,
            family=self.brand.font_minor,
            size_pt=SCALE["small"],
            bold=True,
            color=self.brand.color("on_dark"),
            line_spacing=LINE_SPACING["body"],
        )
        T.add_textbox(
            slide,
            Box(inner.left, att_top + points(28), qbox.width, points(90)),
            role,
            family=self.brand.font_minor,
            size_pt=SCALE["small"],
            color=self.brand.color("on_dark_accent"),
            line_spacing=LINE_SPACING["body"],
        )

        if headline:
            C.two_tone_headline(
                slide,
                ctx,
                Box(rest.left, points(127.38), rest.width, _lines("h2", 3)),
                headline,
                role="h2",
                allow_shrink_to="h3",
            )
        if support is not None:
            top = (
                points(127.38)
                + (_lines("h2", 3) if headline else 0)
                + points(SPACE["md"])
            )
            support(slide, ctx, Box(rest.left, top, rest.width, points(489.99) - top))

        self._furniture(slide, ctx)
        return slide

    def key_findings(
        self,
        headline: str,
        findings: Sequence[tuple[str, str, str]],
        *,
        standfirst: str | None = None,
        section_index: int | None = None,
    ):
        """
        The ring-gauge page. Always page 4 in all four measured CEO decks, 19-27 rings a page.

        Each finding is (value_label, gloss, group). `value_label` is what goes in the ring --
        "73%", "3h", "1-3 yrs". `group` is a key; every finding sharing a group gets the same
        ring colour.

        COLOUR-CODE BY CONTENT STREAM, NOT BY POSITION. On the measured page the rings run
        navy, cobalt, navy, purple, with the purple column's rings purple -- the palette breadth
        CODES MEANING there. Cycling `brand.series(i)` per ring, which is what a chart library
        would do, destroys exactly the thing the colour is doing.

        Capped at six. A page of 19 rings is a report page read at a desk; six is what carries
        from the back of a room, and the value_label has to be legible or the ring is ornament.
        """
        if not 2 <= len(findings) <= 4:
            raise ValueError(
                f"key_findings takes 2 to 4 rings, got {len(findings)}. KPMG runs 19-27 on a "
                f"page read at a desk; a projected page cannot carry that and stay legible.\n"
                f"  The cap fell from six to four when the two-column version was measured: a "
                f"96pt ring in a 194pt column leaves 58pt for its gloss, which is four "
                f"characters a line. A ring nobody can read the gloss of is ornament."
            )
        slide, ctx = self._new(Ground.LIGHT, section=section_index)
        g = self.grid

        beside = self._head_col1(slide, ctx, headline, standfirst=standfirst)

        groups: list[str] = []
        for _, _, group in findings:
            if group not in groups:
                groups.append(group)
        # THE RING COLOUR ALSO SETS THE VALUE LABEL, WHICH IS TYPE, so it has to clear contrast
        # against this page's ground -- and the dataviz palette is not chosen for that. Pacific
        # #00B8F5 is a perfectly good chart colour and measures 2.29:1 on white, so cycling the
        # palette unfiltered put a 34pt numeral on the page at a ratio the package refuses
        # everywhere else. Filter to the members that are legal here, keeping palette ORDER so
        # the stream-coding survives, and fall back to the headline colour rather than emitting
        # something illegible if a deck ever has more streams than legal colours.
        ground_hex = self.brand.color("canvas")
        legal = [
            hexv
            for hexv in self.brand.dataviz
            if contrast_ratio(hexv, ground_hex) >= 4.5
        ]
        if not legal:
            legal = [self.brand.color("headline")]
        stream = {grp: legal[i % len(legal)] for i, grp in enumerate(groups)}

        # ONE FINDING PER ROW, FULL WIDTH. Two columns of rings is what the reference page
        # does at 10.5pt; at our 18pt gloss it leaves 58pt of measure and the gloss becomes a
        # vertical stack of half-words. The ring keeps its size and the gloss gets the rest of
        # the column, which is the same trade the whole package makes against KPMG's density.
        n = len(findings)
        row_h = beside.height // n
        for i, (value_label, gloss, group) in enumerate(findings):
            cell = Box(
                beside.left,
                beside.top + i * row_h,
                beside.width,
                row_h,
            )
            d = min(points(96), row_h - points(SPACE["md"]))
            ring_gauge(
                slide,
                ctx,
                Box(cell.left, cell.top, d, d),
                1.0,
                color=stream[group],
                label=value_label,
                diameter_pt=d / 12700,
            )
            C.body(
                slide,
                ctx,
                Box(
                    cell.left + d + g.gutter,
                    cell.top,
                    cell.width - d - g.gutter,
                    max(row_h, beside.bottom - cell.top),
                ),
                gloss,
                role="small",
                strict=False,
            )

        self._furniture(slide, ctx)
        return slide

    def full_field(
        self,
        headline: str,
        body: str | None = None,
        *,
        standfirst: str | None = None,
        color: str | None = None,
        attribution: str | None = None,
        section_index: int | None = None,
    ):
        """
        A full-bleed colour page. 11.9% of the corpus, 6 Cobalt to 4 KPMG Blue.

        Built to `D1-ceo-private-p-21.png`, which is the measured shape of one: a SHORT headline
        top-left at display size, a LARGE standfirst indented and running most of the page
        width, and small body copy beneath it. Three sizes, three indents, nothing else.

        The first version of this stacked headline / lead / body all on col(0,7) at 24pt and
        it did not survive contact with real text: with a five-line headline the body box came
        out THIRTEEN POINTS TALL, and 266pt of prose was drawn down the page over everything
        under it. The reference page is why -- KPMG gives the standfirst the width, not the
        headline, so the headline stays short and the long text has 700pt of measure to run in.

        Cobalt is the default. #00338D reads as the heavier, more formal one and is correct for
        a closing or an "about" page -- which is what "The journey ahead" and "About the
        authors" are in the measured decks.

        The deck ENDS in a colour crescendo and then drops to white for back matter: D1 closes
        `C C B B D D`, D4 closes `B D B B B D`.
        """
        # THE SOLID GROUND IS KPMG BLUE, NOT COBALT, AND THIS DEFAULT USED TO BE WRONG.
        #
        # Measured off the nine real KPMG slides the FY22 book prints on p.113: both solid-ground
        # slides are #00338D. Cobalt appears there only as a GRADIENT ENDPOINT and as a chart
        # ring -- never as a flat field. Cobalt-as-ground is the brand book's own divider styling
        # (pp.110, 125), which is a different medium: a printed guideline page, not a slide.
        #
        # This was almost certainly the single largest colour error in the build, and it is
        # exactly the one the operator kept naming. Every flat coloured page in the deck was
        # Cobalt, which is both the brightest blue in the palette and the wrong one for a field.
        ground_color = color or self.brand.color("canvas_brand")
        slide, ctx = self._new(Ground.BRAND, section=section_index)
        g = self.grid
        C._rect(slide, g.full, ground_color)
        self._dark_range = (0, g.slide_w)

        bottom = points(489.99)

        # 1. The headline: short, top-left, display size, white into pale blue.
        placed = C.two_tone_headline(
            slide,
            ctx,
            g.col(0, 7, top=g.margin_top, height=_lines("display", 2)),
            headline,
            role="display",
            on_panel=True,
            allow_shrink_to="h2",
        )
        y = placed.bottom + points(SPACE["xl"])

        # 2. The standfirst: indented, wide, and the biggest block of reading on the page.
        #    Indented to col(2) because that is what the reference does -- the headline hangs
        #    left of everything else, which is what makes it read as a title rather than as the
        #    first line of the paragraph.
        # THIS PAGE VIOLATES THE FOUR-OF-SIX BODY MEASURE AND IS LEFT VIOLATING IT, DELIBERATELY,
        # WITH THE MEASUREMENT RECORDED RATHER THAN THE RULE QUIETLY DROPPED.
        #
        # bb p.99: "do not set body copy over more than four columns of the six-column grid, as
        # this will affect legibility." col(2,10) is 702.2pt against a 557.8pt limit -- 144pt
        # over -- and at 18pt that is ~78 characters a line. `closing()` routes through here, so
        # it is the deck's LAST page that reads worst. `grid.max_body_width()` is the clamp and
        # applying it here is one line.
        #
        # It is not applied because THE CLAMP ALONE MOVES THE DEFECT RATHER THAN FIXING IT: with
        # it, `course.py` slide 35 runs 37pt of prose past the bottom of its own box, which
        # `check_overflow` reports as an ERROR. The narrower measure is right and the page is
        # then simply carrying more text than the layout has room for. There are exactly two
        # honest fixes and both are decisions this method cannot make alone:
        #   1. cut the copy on that page (a content change, in examples/content.py), or
        #   2. set the body in TWO columns inside the clamped width, which satisfies p.99 by
        #      measure rather than by block width -- this package's own `_body_columns`
        #      docstring already makes that argument -- and changes what a full-field page
        #      looks like.
        # Until one is chosen, a 702pt measure that is legible beats a 558pt one drawn through
        # the footer.
        indent = g.col(2, 10)
        if standfirst:
            sf = C.body(
                slide,
                ctx,
                Box(indent.left, y, indent.width, bottom - y),
                standfirst,
                role="h3",
                color=self.brand.color("on_dark"),
                strict=False,
            )
            y = sf.bottom + points(SPACE["lg"])

        # 3. The body: small, under the standfirst, in the pale tint so the hierarchy holds.
        if body:
            C.body(
                slide,
                ctx,
                Box(indent.left, y, indent.width, bottom - y),
                body,
                role="small",
                color=self.brand.color("on_dark_accent"),
                strict=False,
            )

        if attribution:
            T.add_textbox(
                slide,
                Box(indent.left, points(455), indent.width, points(34)),
                attribution,
                family=self.brand.font_minor,
                size_pt=SCALE["body"],
                color=self.brand.color("on_dark_accent"),
            )

        self._furniture(slide, ctx)
        return slide

    def exhibit(
        self,
        title: str,
        build: SlideBuilder,
        *,
        source: str | None = None,
        takeaway: str | None = None,
        eyebrow: str | None = None,
        ground: Ground = Ground.LIGHT,
    ):
        """
        The workhorse: an action title over one figure, with a source line.

        `title` must STATE THE FINDING, not name the chart. "Revenue by region" has labelled
        the exhibit; "Three regions grew; the fourth carried the loss" has told the reader what
        to look for before they look. This is the single highest-leverage rule in consulting
        deck craft and it costs nothing but the discipline of writing a sentence.

        `build(slide, ctx, box)` draws the figure into the region it is given. Anything in
        charts.py has that signature.

        `takeaway` adds a callout in the right margin -- use it when the "so what" is not
        already the title, and leave it off when it is.
        """
        slide, ctx = self._new(ground)
        g = self.grid

        plot_top = self._head(slide, ctx, title, eyebrow=eyebrow)
        plot_bottom = g.slide_h - g.margin_bottom - points(SPACE["md"])
        span = 12 if not takeaway else 8
        plot_box = g.col(0, span, top=plot_top, height=plot_bottom - plot_top)

        build(slide, ctx, plot_box)

        if takeaway:
            C.callout(
                slide,
                ctx,
                g.col(9, 3, top=plot_top, height=plot_bottom - plot_top),
                takeaway,
                role="body",
            )

        if source:
            C.source_line(slide, ctx, source)

        self._furniture(slide, ctx)
        return slide

    def points_slide(
        self,
        title: str,
        items: Sequence[tuple[str, str]],
        *,
        eyebrow: str | None = None,
        ground: Ground = Ground.LIGHT,
    ):
        """
        Label/description rows under a headline. The honest replacement for a bullet list.

        The rows block is clamped to the four-of-six body measure (bb p.99). col(0,9) is 630pt
        and each row's description is 18pt prose set across the whole of it; 557.8pt is the
        limit. The labels narrow with it, which is correct -- a label wider than the prose it
        heads is a label that has stopped indexing anything.
        """
        slide, ctx = self._new(ground)
        g = self.grid

        body_top = self._head(slide, ctx, title, eyebrow=eyebrow)
        body_bottom = g.slide_h - g.margin_bottom - points(SPACE["sm"])
        region = g.col(0, 9, top=body_top, height=body_bottom - body_top)
        C.rows(slide, ctx, region.resize(width=g.max_body_width(region.width)), items)

        self._furniture(slide, ctx)
        return slide

    def process(
        self,
        title: str,
        items: Sequence[tuple[str, str]],
        *,
        ground: Ground = Ground.LIGHT,
    ):
        """A numbered horizontal sequence: 2 to 5 steps."""
        slide, ctx = self._new(ground)
        g = self.grid

        top = self._head(slide, ctx, title)
        bottom = g.slide_h - g.margin_bottom - points(SPACE["md"])
        C.steps(slide, ctx, g.col(0, 12, top=top, height=bottom - top), items)

        self._furniture(slide, ctx)
        return slide

    def metrics(
        self,
        title: str,
        stats: Sequence[tuple[str, str, str | None]],
        *,
        source: str | None = None,
        ground: Ground = Ground.LIGHT,
    ):
        """
        Two to four oversized figures. Each stat is (value, label, caption|None).

        Capped at four, because beyond that no figure is large enough to be the point and the
        slide becomes a table with delusions.
        """
        if not 2 <= len(stats) <= 4:
            raise ValueError(f"metrics takes 2 to 4 figures, got {len(stats)}")

        slide, ctx = self._new(ground)
        g = self.grid

        top = self._head(slide, ctx, title)
        bottom = g.slide_h - g.margin_bottom - points(SPACE["md"])
        region = g.col(0, 12, top=top, height=bottom - top)

        n = len(stats)

        # Measure every label up front and use the tallest, so all captions share one baseline.
        cell_w = region.row(n, 0, gap=points(SPACE["xl"])).width
        probe = Box(0, 0, cell_w, points(400))
        label_h = max(
            points(
                T.fit(
                    label,
                    probe,
                    self.brand.font_minor,
                    SCALE["body"],
                    line_spacing=LINE_SPACING["body"],
                ).height_pt
            )
            for _, label, _ in stats
        )

        for i, (value, label, caption) in enumerate(stats):
            cell = region.row(n, i, gap=points(SPACE["xl"]))
            if i > 0:
                C._rect(
                    slide,
                    Box(
                        cell.left - points(SPACE["xl"]) // 2,
                        cell.top,
                        points(0.75),
                        cell.height,
                    ),
                    ctx.rule,
                )
            C.stat(
                slide, ctx, cell, value, label, caption=caption, label_height=label_h
            )

        if source:
            C.source_line(slide, ctx, source)
        self._furniture(slide, ctx)
        return slide

    def versus(
        self,
        title: str,
        left_title: str,
        left_items: Sequence[str],
        right_title: str,
        right_items: Sequence[str],
        *,
        ground: Ground = Ground.LIGHT,
    ):
        """Two columns set against each other."""
        slide, ctx = self._new(ground)
        g = self.grid

        top = self._head(slide, ctx, title)
        bottom = g.slide_h - g.margin_bottom - points(SPACE["md"])
        comparison(
            slide,
            ctx,
            g.col(0, 12, top=top, height=bottom - top),
            left_title,
            left_items,
            right_title,
            right_items,
        )

        self._furniture(slide, ctx)
        return slide

    def quote(self, text: str, attribution: str, *, ground: Ground = Ground.LIGHT):
        """
        A pull quote on white. The quieter of the two quote treatments.

        Use `quote_panel()` for the house form -- a wide Cobalt field with the Pacific glyph.
        This one is for a quotation that is supporting evidence rather than a change of voice
        in its own right.

        The ground defaults to LIGHT, not ALT. `canvas_alt` #F5F6FA does not appear in the
        top-28 flat colours on any of 95 measured pages and could not be isolated as a
        deliberate ground on any of 320; the real light neutral in this system is #E5E5E5 and
        the real ground is white.
        """
        slide, ctx = self._new(ground)
        g = self.grid
        C.quote(
            slide,
            ctx,
            g.col(0, 9, top=int(g.slide_h * 0.30), height=int(g.slide_h * 0.46)),
            text,
            attribution,
        )
        self._furniture(slide, ctx)
        return slide

    def closing(
        self, title: str, lines: Sequence[str] = (), *, color: str | None = None
    ):
        """
        The final statement, on a full colour field.

        Routed through `full_field()` because that is the measured shape of a KPMG closing: the
        deck ends in a colour crescendo and then drops to white for back matter. The default
        here is #00338D rather than Cobalt -- of ten full-bleed pages in the corpus, 6 are
        Cobalt and 4 are KPMG Blue, and the four are the formal ones ("The journey ahead",
        "About the authors"). A closing is formal.
        """
        return self.full_field(
            title,
            standfirst=lines[0] if lines else None,
            body="\n\n".join(lines[1:]) if len(lines) > 1 else None,
            color=color or self.brand.color("canvas_brand"),
        )

    # -- imagery ------------------------------------------------------------
    #
    # THE ENGINE HAD NO IMAGERY AT ALL, AND THAT WAS THE LARGEST SINGLE REASON A CORRECT PALETTE
    # STILL READ AS "A BLUE THEME". bb FY22 gives thirteen pages to image treatment (pp.97-109)
    # and four more to how the window holds one (pp.90-93); six of the nine real KPMG slides on
    # p.113 carry a photograph. Before these three archetypes, `add_picture` appeared in this
    # package exactly once, for the logo.
    #
    # THE TREATMENT IS BAKED IN PILLOW AND ARRIVES AS A FLAT PNG, because OOXML shape fills have
    # no blend modes and the book's treatment IS blend modes. `imagery.py` holds the arithmetic
    # and its provenance; everything here is geometry, the contrast measurement that baking
    # makes possible, and the rights record.
    #
    # ONE WINDOW AT A TIME (p.96). `window_image` is therefore a whole page and never composes
    # with `window`, and no archetype here draws a second frame.

    def _image_block(
        self,
        image,
        box: Box,
        *,
        treatment: str,
        neutralise: bool = False,
        dpi: int = 144,
    ):
        """
        Crop the source to the box, apply the treatment, and record where the pixels came from.

        CROP BEFORE TREAT, NEVER AFTER, and the reason is p.46 rather than efficiency: the
        gradient's stops sit at 0% and 100% OF THE SHAPE. Treat a 3:2 source and then crop it to
        a 7:10 window and the ramp starts and ends off-frame -- still two stops, still linear,
        still 0 degrees, and no longer the specified gradient. Nothing in the file would show it
        and no checker could.

        144 DPI is 1920x1080 for a full-bleed page: enough that no edge in the source shows a
        stair-step at projection, small enough that a 40-slide deck still opens.

        THE PROVENANCE LINE IS NOT OPTIONAL AND IT IS WHY THIS METHOD EXISTS RATHER THAN A BARE
        `add_picture`. We hold no licensed KPMG photography. A deck built on generated imagery
        that does not SAY it is built on generated imagery is a rights claim made by silence, so
        every image writes its own origin into the build log the first time it is used -- once
        per distinct source, because the same field used on six pages is one fact, not six.
        """
        px_w = max(1, int(round(box.width / points(72) * dpi)))
        px_h = max(1, int(round(box.height / points(72) * dpi)))
        treated = IM.treat(
            IM.crop_to_fill(image, px_w, px_h),
            mode=treatment,
            neutral=neutralise,
        )
        line = f"imagery: {IM.provenance(image)}"
        if line not in self._log:
            self._log.append(line)
        return treated

    def _place_image(self, slide, image, box: Box):
        """Insert a baked image. Drawn FIRST by every archetype here, so nothing needs reordering."""
        return slide.shapes.add_picture(
            IM.to_png_bytes(image),
            box.left,
            box.top,
            width=box.width,
            height=box.height,
        )

    def _window_type(self, slide, box: Box, text_str: str, color: str, floor: str):
        """
        Set display type into a box on a gradient, measuring first and raising if it will not fit.

        Placed through `text` rather than through `components.headline` for the same reason
        `window()` does it: the component asserts legibility against the CONTEXT's ground, and
        every context here describes a flat canvas colour while the type is standing on a
        gradient the Ctx cannot see. It would compare KPMG Blue with KPMG Blue and refuse,
        correctly and uselessly. The legibility guarantee is not lost -- it is discharged in the
        archetype docstrings, against BOTH ends of the relevant ramp, which is stricter than the
        single-colour check this bypasses.

        Shrinking stops at `floor` and then raises, rather than continuing down. A headline that
        will not fit a window at 34pt is carrying two ideas, and the fix is in the writing.
        """
        # bb FY22 puts the headline cap height at ~2x logo height, and Arial's cap is 0.716 of
        # its em, so 2 x 36pt of cap wants roughly a 100pt font -- far above the scale's 60pt
        # display ceiling, and above the 74-96pt measured off the p.113 gallery. Starting high
        # and stepping down is what makes the headline FILL the shape rather than float in it.
        size = max(SCALE["display"], int(points(36.0) / 12700 * 2 / 0.716))
        while (
            size > SCALE[floor]
            and not T.fit(
                text_str,
                box,
                self.brand.font_major,
                size,
                bold=True,
                line_spacing=LINE_SPACING["display"],
            ).fits
        ):
            size -= 2
        fitted = T.fit(
            text_str,
            box,
            self.brand.font_major,
            size,
            bold=True,
            line_spacing=LINE_SPACING["display"],
        )
        if not fitted.fits:
            raise ValueError(
                f"headline does not fit its window even at the {floor} floor "
                f"({SCALE[floor]}pt): {fitted.describe()}\n  text: {text_str!r}\n"
                f"  CUT WORDS. bb FY22's own slide headlines run 2-12 words, median 3."
            )
        T.add_textbox(
            slide,
            box.resize(height=points(fitted.height_pt)),
            text_str,
            family=self.brand.font_major,
            size_pt=size,
            bold=True,
            color=color,
            line_spacing=LINE_SPACING["display"],
        )
        return box.resize(height=points(fitted.height_pt))

    def window_image(
        self,
        headline: str,
        *,
        image,
        subhead: str | None = None,
        support_ground: bool = False,
        neutralise: bool = False,
        section_index: int | None = None,
    ):
        """
        Window style 3: the window holds an IMAGE, the headline sits on the gradient ground.

        bb FY22 p.72 names four window styles and only one of them puts type inside the shape.
        `window()` builds that one. A deck using only style 2 is leaning on the least
        characteristic use of the device -- the book's own words, p.72: "don't overemphasize the
        type relationship with the window; first and foremost the window should interact with
        objects and people."

        THE IMAGE IS NOT GRADIENT-OVERLAID HERE, AND THAT IS THE WHOLE COMPOSITION. p.92's
        four-step ramp ends at "opaque gradient -- used for style 3, the window as holding shape
        for imagery; the gradient at full opacity is a perfect contrast to the light-colored
        photo", and p.90 states the principle: "the contrast between the foreground (neutral
        image) and background (color) means that the color and neutral image complement each
        other and provide focus in the layout". Putting the colour over the image AND behind it
        collapses that contrast and produces one purple page with a slightly different purple
        rectangle on it, which is what the naive reading of "apply the treatment" gives you.

        `support_ground` IS NOT A FREE CHOICE OF PALETTE. p.46: support "never appears without
        primary". With the support ramp on the ground and an untreated image in the window, the
        primary gradient would be absent from the page entirely, so this flag also moves the
        image onto the primary treatment at Overlay 65% -- which is separately sanctioned
        (p.92 step 2) and restores the rule. Passing `support_ground=True` with an untreatable
        image is the one combination that raises.

        GEOMETRY IS DENOMINATED IN LOGO HEIGHTS AND SIZED OFF HEIGHT (pp.72-89), as in
        `window()`: 1x clearance to the live-area edge, 2x between window and type, 0.75x type
        inset. Ours is the VERTICAL 7:10 form and it sits on the RIGHT, which is p.79's first
        sanctioned placement ("to the right of the logo") and gives the headline a real column
        rather than the strip `window()`'s horizontal 10:7 form leaves.

        TYPE COLOUR IS MEASURED AGAINST BOTH ENDS OF THE RAMP IT STANDS ON, not against one
        flat colour: white on Purple->Cobalt measures 7.01 and 6.76, and KPMG Blue on
        Pacific->Light Blue measures 4.94 and 8.59. All four clear 4.5:1, so the AA rider holds
        along the whole gradient rather than at one end of it.
        """
        g = self.grid
        slide, ctx = self._new(Ground.BRAND, section=section_index)

        ground = X.GRADIENT_SUPPORT if support_ground else X.GRADIENT_PRIMARY
        treatment = "overlay" if support_ground else "neutral"
        on_ground = (
            self.brand.color("ink") if support_ground else self.brand.color("on_dark")
        )

        X.set_gradient(
            C._rect(slide, g.full, ground[0][1]),
            list(ground),
            angle_deg=X.GRADIENT_ANGLE_DEG,
        )
        self._dark_range = (0, g.slide_w)

        logo_h = points(36.0)  # measured, 89.6 x 36.0pt at (72,72) on four KPMG covers
        win_top = g.cover_margin
        win_h = g.slide_h - 2 * g.cover_margin
        win_w = int(win_h * 7 / 10)  # p.72: 7:10 vertical
        win = Box(g.slide_w - g.cover_margin - win_w, win_top, win_w, win_h)
        self._place_image(
            slide,
            self._image_block(image, win, treatment=treatment, neutralise=neutralise),
            win,
        )

        # The type column: from the layout margin to 2x logo height clear of the window.
        col = Box(
            g.cover_margin,
            win_top,
            win.left - 2 * logo_h - g.cover_margin,
            win_h,
        )
        head_h = col.height if not subhead else int(col.height * 0.68)
        placed = self._window_type(
            slide,
            Box(col.left, col.top, col.width, head_h),
            headline,
            on_ground,
            floor="h2",
        )

        if subhead:
            T.add_textbox(
                slide,
                Box(
                    col.left,
                    placed.bottom + points(SPACE["md"]),
                    col.width,
                    col.bottom - placed.bottom - points(SPACE["md"]),
                ),
                subhead,
                family=self.brand.font_minor,
                size_pt=SCALE["lead"],
                color=on_ground,
                line_spacing=LINE_SPACING["body"],
            )

        self._furniture(slide, ctx, chrome=False)
        return slide

    def image_field(
        self,
        headline: str,
        *,
        image,
        standfirst: str | None = None,
        body: str | None = None,
        treatment: str = "softlight_multiply",
        neutralise: bool = False,
        section_index: int | None = None,
    ):
        """
        The full-bleed photographic page: a treated image edge to edge, type over its lower left.

        This is p.113's photographic slide, which is 2 of that gallery's 9 and the one page type
        this package could not build at all. The default treatment is p.93's fully worked
        example -- Soft Light at 100%, then the same gradient again in Multiply at 40% -- because
        that is the recipe the book demonstrates step by step and the one whose 40% was
        independently confirmed by measuring the book's own three panels.

        CONTRAST HERE IS MEASURED OFF THE PIXELS, NOT ASSUMED. A picture is not a filled shape,
        so nothing in `verify` can see what type over an image is standing on -- `_filled_rects`
        reads `a:solidFill` and `a:gradFill`, and a `p:pic` has neither, which means every run
        on a page like this would otherwise be checked against a slide background the reader
        cannot see. So the type region is measured in the baked image at the bias that is worst
        for the palest ink on the page (the brightest quarter of the region, not its mean), and
        the brand gradient is ramped up in Multiply from the bottom edge until it clears.

        THE FOREGROUND THAT SETS THE TARGET IS #ACEAFF, NOT WHITE, AND GETTING THAT WRONG COST A
        REAL FAILURE. Sizing the scrim for white type left the page's own furniture -- the
        navigation strip and the footer, both set in the pale tint -- at 4.43:1 on a ground that
        gave white a comfortable 5.82. The rule is general: measure against the LEAST contrasting
        ink the page carries, not the one the headline happens to use, because the headline is
        never the run that fails first.

        THE DECLARED GROUND IS COBALT, WHICH IS A PALETTE COLOUR AND A CONSERVATIVE ONE. Writing
        the raw measured hex back as the slide background told the verifier the truth and cost a
        palette-conformance warning for a colour that exists nowhere in the brand -- accurate and
        useless. Instead the scrim is driven until the region is at least as dark as Cobalt
        #1E49E2, the treatment gradient's own endpoint, and Cobalt is then declared as the
        ground. That declaration UNDERSTATES the real contrast in every case, so it can only ever
        fail safe, and the measured value goes into the build log for anyone who wants it.

        An image that cannot reach the target raises.

        THE TYPE IS ANCHORED LOW, WHICH IS A DELIBERATE DEPARTURE FROM THE DECK'S TOP MARGIN.
        The upper two-thirds of the frame stay clear, so the image reads as the page rather than
        as a texture behind a slide; and the scrim that makes the type legal runs from the
        bottom edge, where a darkening ramp reads as light falling off rather than as a panel
        laid over a photograph.
        """
        g = self.grid
        slide, ctx = self._new(Ground.DARK, section=section_index)
        self._dark_range = (0, g.slide_w)

        treated = self._image_block(
            image, g.full, treatment=treatment, neutralise=neutralise
        )

        # WHERE THE TYPE STARTS FOLLOWS HOW MUCH OF IT THERE IS. A fixed 0.44 anchored the block
        # low, which is right for a headline alone and ran 19pt of body copy through the bottom
        # of the page once a standfirst and body were both present. The stack is measured
        # downward from here, so the start has to move up as it grows -- and a page carrying all
        # three is a page whose image is doing less work anyway.
        share = 0.56 if not (standfirst or body) else (0.48 if not body else 0.38)
        type_top = int(g.slide_h * share)
        col = g.col(0, 7)
        fg = self.brand.color("on_dark")
        # The palest ink on the page, which is what the scrim has to satisfy -- see the docstring.
        palest = self.brand.color("on_dark_accent")
        standin = self.brand.color("accent")  # Cobalt, the declared ground
        # The region the type will occupy, in image fractions. Deliberately a little wider and
        # taller than the boxes below it: the measurement should cover the glyphs' surroundings,
        # not just their bounding boxes.
        treated, measured = IM.deepen_for_type(
            treated,
            region=(0.0, type_top / g.slide_h, 0.62, 1.0),
            fg=palest,
            min_ratio=contrast_ratio(palest, standin),
            edge="bottom",
            feather=0.22,
        )
        self._place_image(slide, treated, g.full)
        X.set_slide_background(slide, color=standin)
        self._log.append(
            f"imagery: page {self._page} type region measures {measured} "
            f"({contrast_ratio(palest, measured):.2f}:1 against {palest}); "
            f"declared as {standin}"
        )

        y = type_top
        placed = C.headline(
            slide,
            ctx,
            Box(col.left, y, col.width, _lines("display", 2)),
            headline,
            role="display",
            color=fg,
            allow_shrink_to="h2",
        )
        y = placed.bottom + points(SPACE["md"])
        bottom = g.slide_h - g.margin_bottom

        if standfirst:
            sf = C.body(
                slide,
                ctx,
                Box(col.left, y, col.width, bottom - y),
                standfirst,
                role="h3",
                color=fg,
                strict=False,
            )
            y = sf.bottom + points(SPACE["sm"])

        if body:
            C.body(
                slide,
                ctx,
                Box(col.left, y, col.width, bottom - y),
                body,
                role="small",
                color=self.brand.color("on_dark_accent"),
                strict=False,
            )

        self._furniture(slide, ctx)
        return slide

    def image_split(
        self,
        headline: str,
        standfirst: str,
        body: str,
        *,
        image,
        treatment: str = "overlay",
        neutralise: bool = False,
        section_index: int | None = None,
    ):
        """
        A working page: the argument in the left seven columns, a treated image block bleeding right.

        THIS IS THE ONE THAT WILL CARRY MOST INTERIOR PAGES, so it is built as a working page
        rather than as a statement: a white ground, a real headline, a standfirst that states the
        claim and body copy that argues it. All three are required arguments, for the same reason
        `statement()` requires its three -- a page with an image and a title and nothing else is
        a page that has not said anything.

        THE TREATMENT DEFAULTS TO OVERLAY AT 65%, NOT TO THE HEAVIER RECIPE, and the choice is
        p.92's: Overlay at 65% is step 2 of its four-step ramp, "moderate transparency", and the
        book's note on step 3 is the argument -- "if used at full opacity, photo details
        disappear". A block image beside argued copy is support imagery, and support imagery that
        has lost its detail is a coloured rectangle with a photograph's file size.

        THE IMAGE BLEEDS RIGHT AND IS INSET TOP AND BOTTOM, WHICH IS A CORRECTNESS DECISION AS
        MUCH AS A COMPOSITIONAL ONE. Bleeding all four edges is the stronger look, and it puts
        the nav strip and the footer's page number on top of a light treated image at roughly
        3.1:1 -- a real contrast failure on every page built this way, and one the reversed-ink
        mechanism (`_dark_range`) cannot fix because the image is light. Stopping at y=48 and
        y=500 keeps both chrome bands on white, where they are legal by construction. The right
        edge still bleeds, which is where the move earns its keep.

        The column split is the grid's own: type on col(0,7), image from col(7) to the slide
        edge, one gutter between them.
        """
        g = self.grid
        slide, ctx = self._new(Ground.LIGHT, section=section_index)

        # Clear of the nav strip (y 12.4-36.0) above and the footer band (y 513.7-529.1) below.
        img_left = g.col(7, 5).left
        img_box = Box(img_left, points(48.0), g.slide_w - img_left, points(452.0))
        self._place_image(
            slide,
            self._image_block(
                image, img_box, treatment=treatment, neutralise=neutralise
            ),
            img_box,
        )

        col = g.col(0, 7)
        bottom = g.slide_h - g.margin_bottom
        y = g.margin_top

        placed = C.headline(
            slide,
            ctx,
            Box(col.left, y, col.width, _lines("h1", 2)),
            headline,
            role="h1",
            allow_shrink_to="h3",
        )
        y = placed.bottom + points(SPACE["md"])

        sf = C.body(
            slide,
            ctx,
            Box(col.left, y, col.width, bottom - y),
            standfirst,
            role="lead",
            color=ctx.accent,
            strict=False,
        )
        y = sf.bottom + points(SPACE["md"])

        C.body(
            slide,
            ctx,
            Box(col.left, y, col.width, bottom - y),
            body,
            role="body",
            strict=False,
        )

        self._furniture(slide, ctx)
        return slide

    def custom(self, ground: Ground = Ground.LIGHT):
        """
        An escape hatch: a blank slide plus its Ctx, for something the archetypes cannot express.

        Use it rarely and deliberately. Every custom slide is a slide that will not
        automatically match the others, so the consistency the archetypes buy has to be
        re-established by hand. If you find yourself reaching for this repeatedly, the right
        move is to add an archetype rather than to keep improvising.
        """
        slide, ctx = self._new(ground)
        self._furniture(slide, ctx)
        return slide, ctx

    # -- output -------------------------------------------------------------

    def save(self, path: str) -> str:
        self.prs.save(path)
        return path

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides._sldIdLst)

    def build_log(self) -> list[str]:
        """What the builder did, including any provenance warning. Print it at build time."""
        return list(self._log)
