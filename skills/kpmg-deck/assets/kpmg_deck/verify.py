"""
verify.py -- the mechanical gates.

There are two kinds of defect in a generated deck, and confusing them is why most generated
decks ship broken.

MACHINE-CHECKABLE: the file is structurally invalid, a colour is off-palette, type is below the
legibility floor, a shape hangs off the slide, contrast fails, sibling slides disagree. Every one
of these is a fact about the file. This module checks all of them, and an agent with no design
judgment whatsoever can run it and act on the result.

EYE-ONLY: two elements collide, a headline has an ugly rag, the composition is bottom-heavy, an
exhibit's emphasis lands on the wrong bar. NOTHING in this module can see any of these. They are
found by rendering the deck to images and LOOKING at them, which is why scripts/render.sh exists
and why the skill makes that step mandatory rather than optional.

The honest framing, and it is the whole reason this file has a docstring this long: passing every
check here means the deck is not BROKEN. It does not mean the deck is GOOD. Treating a green
verify as sign-off is the exact mistake this package is built to prevent -- in its own first
render, every one of five slides with a rule struck through its headline passed every structural
check cleanly, because valid XML in the wrong place is still valid XML.
"""

from __future__ import annotations

import posixpath
import re
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from typing import Iterable

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .canvas import EMU_PER_POINT, Grid
from .tokens import MIN_READABLE_PT, Brand, contrast_ratio, contrast_verdict

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class Finding:
    level: str  # "ERROR" -- would ship broken; "WARN" -- worth a human look
    check: str
    slide: int | None
    detail: str

    def __str__(self) -> str:
        where = f"slide {self.slide}" if self.slide else "deck"
        return f"[{self.level}] {self.check} ({where}): {self.detail}"


@dataclass
class Report:
    findings: list[Finding] = field(default_factory=list)
    checked: list[str] = field(default_factory=list)

    @property
    def errors(self) -> list[Finding]:
        return [f for f in self.findings if f.level == "ERROR"]

    @property
    def warnings(self) -> list[Finding]:
        return [f for f in self.findings if f.level == "WARN"]

    @property
    def ok(self) -> bool:
        return not self.errors

    def add(self, level: str, check: str, slide: int | None, detail: str) -> None:
        self.findings.append(Finding(level, check, slide, detail))

    def render(self) -> str:
        lines = [
            f"{'PASS' if self.ok else 'FAIL'} -- {len(self.errors)} error(s), "
            f"{len(self.warnings)} warning(s) across {len(self.checked)} checks",
            "",
        ]
        lines += [f"  {f}" for f in self.errors]
        lines += [f"  {f}" for f in self.warnings]
        if not self.findings:
            lines.append("  no structural defects found")
        lines += [
            "",
            "  Structural checks only. Collisions, ragged line breaks, bad emphasis and",
            "  compositional problems are INVISIBLE here -- render the deck and look at it.",
        ]
        return "\n".join(lines)


# ---------------------------------------------------------------------------
# Package-level integrity -- the "[Repaired]" gate
# ---------------------------------------------------------------------------


def check_package(path: str, report: Report) -> None:
    """
    Assert the .pptx is a structurally sound OOXML package.

    This is the check that stands between a build and a "[Repaired]" dialog in front of a room.
    PowerPoint offers to repair a file when it finds a defect its schema validation trips over,
    and the failure mode is brutal: it typically "repairs" by DISCARDING the offending content,
    so the deck opens with slides silently missing.

    The four causes checked here account for essentially every repair prompt a generator
    produces:
      1. a corrupt zip member
      2. a part that is not well-formed XML
      3. a relationship id referenced from a slide that its .rels file does not define
         (orphaned rId -- the classic result of deleting a slide by removing only the sldId)
      4. a part present in the package with no [Content_Types].xml entry
    """
    report.checked.append("package integrity")

    with zipfile.ZipFile(path) as z:
        bad = z.testzip()
        if bad:
            report.add("ERROR", "package integrity", None, f"corrupt zip member: {bad}")
            return

        names = set(z.namelist())

        if "[Content_Types].xml" not in names:
            report.add(
                "ERROR", "package integrity", None, "[Content_Types].xml is missing"
            )
            return

        for name in sorted(names):
            if name.endswith((".xml", ".rels")):
                try:
                    etree.fromstring(z.read(name))
                except etree.XMLSyntaxError as exc:
                    report.add(
                        "ERROR",
                        "package integrity",
                        None,
                        f"{name} is not well-formed XML: {exc}",
                    )

        # Content-type coverage: every part needs a Default (by extension) or an Override.
        try:
            ct = etree.fromstring(z.read("[Content_Types].xml"))
            defaults = {
                d.get("Extension", "").lower() for d in ct if d.tag.endswith("Default")
            }
            overrides = {o.get("PartName") for o in ct if o.tag.endswith("Override")}
            for name in names:
                if name == "[Content_Types].xml" or name.endswith("/"):
                    continue
                ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                if ext not in defaults and f"/{name}" not in overrides:
                    report.add(
                        "ERROR",
                        "package integrity",
                        None,
                        f"{name} has no content-type Default or Override",
                    )
        except etree.XMLSyntaxError:
            pass

        # Relationship integrity, per slide.
        for name in sorted(names):
            if not re.match(r"ppt/slides/slide\d+\.xml$", name):
                continue
            idx = int(re.search(r"slide(\d+)\.xml", name).group(1))
            rels_name = posixpath.join(
                posixpath.dirname(name), "_rels", posixpath.basename(name) + ".rels"
            )
            defined = set()
            if rels_name in names:
                try:
                    rels = etree.fromstring(z.read(rels_name))
                    defined = {r.get("Id") for r in rels}
                except etree.XMLSyntaxError:
                    report.add(
                        "ERROR",
                        "package integrity",
                        idx,
                        f"{rels_name} is not well-formed",
                    )
                    continue
            try:
                doc = etree.fromstring(z.read(name))
            except etree.XMLSyntaxError:
                continue
            for el in doc.iter():
                for key, value in el.attrib.items():
                    if (
                        key.endswith("}id")
                        or key.endswith("}embed")
                        or key.endswith("}link")
                    ):
                        if value and value.startswith("rId") and value not in defined:
                            report.add(
                                "ERROR",
                                "package integrity",
                                idx,
                                f"references undefined relationship {value}",
                            )

        # 5. THE SKELETON PARTS, AND THIS IS THE ONE THE OTHER FOUR CANNOT SEE.
        #
        # The four checks above all test something that is PRESENT and wrong. This tests for
        # something ABSENT, which is a different failure and an invisible one: a package can be
        # a valid zip, hold only well-formed XML, declare every content type and define every
        # relationship it uses, and STILL make PowerPoint show the repair dialog -- because a
        # part it expects unconditionally is simply not there.
        #
        # Nothing upstream catches this. python-pptx writes these parts when it builds from its
        # default template, so the defect only appears once a generator starts assembling or
        # rewriting a package itself, which is exactly what this one does to the theme. And the
        # repair is SILENT in the worst way: PowerPoint regenerates the missing part with its
        # own defaults, so a rewritten theme is quietly replaced by the stock Office theme and
        # the deck opens in the wrong colours with no error anywhere.
        #
        # The list is the set of parts a PowerPoint-authored .pptx always contains, and each is
        # here because its absence has a named consequence rather than because it looked tidy.
        for part, consequence in (
            ("ppt/presentation.xml", "the presentation itself"),
            ("ppt/_rels/presentation.xml.rels", "slide order and theme binding"),
            (
                "ppt/theme/theme1.xml",
                "every brand colour and font -- repaired to stock Office",
            ),
            ("ppt/presProps.xml", "presentation-level defaults"),
            ("ppt/viewProps.xml", "view state"),
            ("ppt/tableStyles.xml", "table defaults"),
            ("docProps/core.xml", "title and author metadata"),
            ("docProps/app.xml", "application metadata"),
        ):
            if part not in names:
                report.add(
                    "ERROR",
                    "package integrity",
                    None,
                    f"{part} is missing -- PowerPoint will repair the file silently and "
                    f"regenerate it, losing {consequence}",
                )


# ---------------------------------------------------------------------------
# Design-system conformance
# ---------------------------------------------------------------------------


def _iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        yield shape
        if (
            getattr(shape, "shape_type", None) is not None and shape.shape_type == 6
        ):  # GROUP
            yield from _iter_shapes(shape.shapes)


def check_deck(
    path: str, brand: Brand, report: Report, *, grid: Grid | None = None
) -> None:
    """Palette conformance, type floor, geometry, contrast and per-slide sanity."""
    prs = Presentation(path)
    grid = grid or Grid(slide_w=prs.slide_width, slide_h=prs.slide_height)

    report.checked += [
        "slide size",
        "palette conformance",
        "type floor",
        "geometry",
        "empty shapes",
    ]

    if (prs.slide_width, prs.slide_height) != (12192000, 6858000):
        report.add(
            "WARN",
            "slide size",
            None,
            f"{prs.slide_width}x{prs.slide_height} EMU is not the 16:9 standard "
            f"12192000x6858000; confirm this is deliberate",
        )

    legal = {
        c.upper().lstrip("#")
        for c in list(brand.theme_colors.values())
        + list(brand.roles.values())
        + list(brand.dataviz)
    }
    legal.update({"FFFFFF", "000000"})

    for index, slide in enumerate(prs.slides, start=1):
        # -- off-palette literals -------------------------------------------
        xml = slide._element
        for srgb in xml.iter(qn("a:srgbClr")):
            val = (srgb.get("val") or "").upper()
            if val and val not in legal:
                report.add(
                    "WARN",
                    "palette conformance",
                    index,
                    f"#{val} is not in the {brand.name} palette",
                )

        # -- type below the legibility floor --------------------------------
        #
        # RUNNING CHROME IS EXEMPT BY POSITION, NOT BY SIZE. The nav strip and the footer are
        # document furniture: nobody reads them and their job is to signal that this is a
        # document rather than a slide. A 10pt nav tab is correct and a 10pt body line is not,
        # and the only thing that distinguishes them is WHERE THEY ARE -- so the exemption is
        # keyed on the band, and the 14pt floor is not weakened anywhere else. Keying it on
        # size instead would have exempted every undersized run in the deck.
        nav, foot = grid.nav_band, grid.footer_band
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            try:
                top, height = shape.top, shape.height
            except (AttributeError, TypeError):
                continue
            if None in (top, height):
                continue
            in_chrome = (
                top >= nav.top and top + height <= nav.bottom + EMU_PER_POINT * 8
            ) or (top >= foot.top - EMU_PER_POINT * 6)
            for rpr in shape._element.iter(qn("a:rPr")):
                sz = rpr.get("sz")
                if sz is None:
                    continue
                pt = int(sz) / 100
                if pt >= MIN_READABLE_PT:
                    continue
                if in_chrome and pt >= 6:
                    continue
                level = "WARN" if pt >= 6 else "ERROR"
                report.add(
                    level,
                    "type floor",
                    index,
                    f"{pt:g}pt is below the {MIN_READABLE_PT}pt floor for a large room",
                )

        # -- geometry --------------------------------------------------------
        for shape in _iter_shapes(slide.shapes):
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
            except (AttributeError, TypeError):
                continue
            if None in (left, top, width, height):
                continue
            if width <= 0 or height <= 0:
                report.add(
                    "ERROR",
                    "geometry",
                    index,
                    f"shape {shape.shape_id} has non-positive size {width}x{height}",
                )
                continue
            # A deliberately bleeding element (the section numeral) may exceed the slide; only
            # flag a shape that is ENTIRELY outside, which is never intentional.
            if (
                left > grid.slide_w
                or top > grid.slide_h
                or left + width < 0
                or top + height < 0
            ):
                report.add(
                    "ERROR",
                    "geometry",
                    index,
                    f"shape {shape.shape_id} is entirely off-slide at ({left}, {top})",
                )

            # Only a TEXT BOX with no text is a defect. Every autoshape carries a text
            # frame it is not obliged to use, so flagging those makes the report 90% noise --
            # and a report nobody reads catches nothing.
            if (
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                and shape.has_text_frame
                and not shape.text_frame.text.strip()
            ):
                report.add(
                    "WARN",
                    "empty shapes",
                    index,
                    f"text box {shape.shape_id} is empty",
                )


def check_wrap_and_autofit(path: str, report: Report) -> None:
    """
    Assert no text frame has word wrap OFF or "resize shape to fit text" ON.

    This guards the single most common root cause of overflow in a generated deck, and it is a
    library default rather than anything an author does. `add_textbox()` writes:

        <a:bodyPr wrap="none"><a:spAutoFit/></a:bodyPr>

    -- word wrap OFF, and "shape resizes to fit text" ON, while python-pptx never actually
    recomputes the shape height. So every textbox created the naive way is a box whose stored
    height is a lie and whose text will not wrap: in PowerPoint the line simply runs off the
    slide.

    It is worth checking mechanically rather than trusting the builder, for a nasty reason:
    LibreOffice HIDES this defect. Given wrap="none" it silently wraps the text inside the box
    anyway, so the render loop shows a tidy block while PowerPoint would run the line several
    inches past the edge. This is the one defect class the render loop cannot be trusted on,
    which is exactly why it belongs in the mechanical gate instead.

    `text.write()` sets word_wrap=True and pins noAutofit, so a failure here means something
    bypassed it.
    """
    report.checked.append("wrap and autofit")
    with zipfile.ZipFile(path) as z:
        for name in sorted(z.namelist()):
            m = re.match(r"ppt/slides/slide(\d+)\.xml$", name)
            if not m:
                continue
            idx = int(m.group(1))
            for bodyPr in etree.fromstring(z.read(name)).iter(qn("a:bodyPr")):
                if bodyPr.get("wrap") == "none":
                    report.add(
                        "ERROR",
                        "wrap and autofit",
                        idx,
                        'a text frame has wrap="none" -- text will run off the slide '
                        "in PowerPoint, and LibreOffice will not show you",
                    )
                if bodyPr.find(qn("a:spAutoFit")) is not None:
                    report.add(
                        "ERROR",
                        "wrap and autofit",
                        idx,
                        "a text frame carries spAutoFit, whose stored height python-pptx "
                        "never recomputes",
                    )


def _filled_rects(slide):
    """
    Every filled shape on the slide, in z-order (document order = back to front).

    Returns [(left, top, right, bottom, "#RRGGBB"), ...]. Used to work out what a run of text
    is ACTUALLY sitting on.

    GRADIENT FILLS COUNT, AND SKIPPING THEM PRODUCED A FALSE FAILURE THAT LOOKED EXACTLY LIKE A
    REAL ONE. This used to match `a:solidFill` only, so the cover's full-bleed gradient field
    was invisible to it and every run on the cover was checked against the SLIDE background --
    white. The report then said "#ACEAFF on #FFFFFF is 1.32:1", naming a pairing that does not
    occur anywhere in the file. A gate that cries wolf on a correct page is worse than no gate:
    the next real failure is one line in a list of thirty phantoms.
    A gradient is resolved to its DARKEST stop, which is the worst case for light text and the
    best case for catching a genuine failure -- the check should err toward reporting.
    """
    out = []
    for shape in _iter_shapes(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            continue  # a text box is not a background for itself
        spPr = getattr(shape._element, "spPr", None)
        if spPr is None:
            continue
        srgb = spPr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
        if srgb is None or not srgb.get("val"):
            grad = spPr.find(qn("a:gradFill"))
            if grad is not None:
                stops = [
                    s.get("val") for s in grad.iter(qn("a:srgbClr")) if s.get("val")
                ]
                if stops:
                    from .tokens import relative_luminance

                    darkest = min(stops, key=lambda v: relative_luminance("#" + v))
                    try:
                        left, top, w, h = (
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height,
                        )
                    except (AttributeError, TypeError):
                        continue
                    if None not in (left, top, w, h) and w > 0 and h > 0:
                        out.append((left, top, left + w, top + h, "#" + darkest))
            continue
        # A TRANSLUCENT OVERLAY IS NOT A GROUND. A 13%-white streak over a dark blue field
        # leaves that field dark; treating the streak's own #FFFFFF as the ground turns every
        # run under it into a reported white-on-white failure. The alpha is right there in the
        # XML, so read it: anything under 85% opaque is a treatment, not a background.
        alpha_el = srgb.find(qn("a:alpha"))
        if alpha_el is not None:
            try:
                if int(alpha_el.get("val", "100000")) < 85000:
                    continue
            except (TypeError, ValueError):
                pass
        try:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
        except (AttributeError, TypeError):
            continue
        if None in (left, top, w, h) or w <= 0 or h <= 0:
            continue
        out.append((left, top, left + w, top + h, "#" + srgb.get("val")))
    return out


def _effective_bg(shape, rects, slide_bg: str) -> str:
    """
    The colour a text shape is actually sitting on.

    Checking every run against the SLIDE background is wrong the moment a deck uses colour
    panels: white type on a cobalt panel reads as white-on-white and the gate reports a
    contrast failure on a slide that is perfectly legible. This walks the filled shapes from
    front to back and takes the first one containing the point where the INK actually is.

    SAMPLING THE BOX'S CENTRE WAS THAT POINT, AND IT WAS WRONG FOR EVERY LEFT-ALIGNED RUN IN A
    WIDE BOX. The cover's eyebrow sits at x=72 in a box 560pt wide, so its centre is 352pt
    away from its own glyphs -- far enough to land inside a decorative overlay it never
    touches, which is exactly what it did: the gate reported "#ACEAFF on #FFFFFF is 1.32:1" for
    pale blue type sitting on a dark blue field. Sampling by ALIGNMENT costs one attribute read
    and removes a whole class of phantom failure.
    """
    try:
        left_e, top_e, w, h = shape.left, shape.top, shape.width, shape.height
    except (AttributeError, TypeError):
        return slide_bg
    if None in (left_e, top_e, w, h):
        return slide_bg

    align = None
    if shape.has_text_frame and shape.text_frame.paragraphs:
        align = shape.text_frame.paragraphs[0].alignment
    inset = min(w // 8, EMU_PER_POINT * 40)
    if align is not None and "RIGHT" in str(align):
        cx = left_e + w - inset
    elif align is not None and "CENTER" in str(align):
        cx = left_e + w // 2
    else:
        cx = left_e + inset
    cy = top_e + h // 2

    for left, top, right, bottom, color in reversed(rects):
        if left <= cx <= right and top <= cy <= bottom:
            return color
    return slide_bg


def check_contrast(path: str, brand: Brand, report: Report) -> None:
    """
    Check every run's colour against the ground it is actually sitting on.

    Resolves the ground per slide from <p:bg>, falling back to the brand canvas. This catches
    the pairing that is easiest to produce by accident with a blue-heavy palette -- a mid blue
    on the dark blue ground, which measures around 2.3:1 and is invisible in a lit room.
    """
    report.checked.append("contrast")
    prs = Presentation(path)

    for index, slide in enumerate(prs.slides, start=1):
        bg = brand.color("canvas")
        bg_el = slide._element.find(f"{qn('p:cSld')}/{qn('p:bg')}")
        if bg_el is not None:
            srgb = bg_el.find(f".//{qn('a:srgbClr')}")
            if srgb is not None and srgb.get("val"):
                bg = "#" + srgb.get("val")

        rects = _filled_rects(slide)

        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            bg_here = _effective_bg(shape, rects, bg)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rpr = run.font._rPr
                    srgb = (
                        rpr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
                        if rpr is not None
                        else None
                    )
                    if srgb is None or not srgb.get("val"):
                        continue
                    fg = "#" + srgb.get("val")
                    size_pt = run.font.size.pt if run.font.size else 18
                    # Alpha-ghosted runs are deliberately low-contrast; skip them.
                    if srgb.find(qn("a:alpha")) is not None:
                        continue
                    # A RUN WITH NO LETTERS OR DIGITS IS A GRAPHIC MARK, NOT TEXT. The
                    # oversized quotation glyph is the case in hand: measured Pacific #00B8F5
                    # on Cobalt, a fixed-size KPMG device recurring on five pages, and 2.95:1
                    # -- half a hundredth under the large-text threshold. WCAG's floors are
                    # about reading load, and a mark carrying no characters imposes none; the
                    # quotation is signalled by the words, not by the glyph.
                    # Deliberately narrow: one run, no alphanumerics, and only at display size.
                    # Anything a reader must actually read still gets the full check.
                    if size_pt >= 48 and not any(c.isalnum() for c in run.text):
                        continue
                    verdict, ratio = contrast_verdict(
                        fg, bg_here, size_pt, bool(run.font.bold)
                    )
                    if verdict == "FAIL":
                        report.add(
                            "ERROR",
                            "contrast",
                            index,
                            f"{fg} on {bg_here} is {ratio:.2f}:1 at {size_pt:g}pt "
                            f'-- "{run.text[:44]}"',
                        )


def check_pairs(path: str, brand: Brand, report: Report) -> None:
    """
    Check every two-colour page against KPMG's own pairing rule and against the AA rider.

    A "two-colour layout" here is detected by GEOMETRY, not by asking the builder what it meant:
    two or more full-height bands whose x-intervals together cover the slide. That is exactly
    what the brand book's Group 3 pages are -- a narrow vertical element against a wide field --
    and it is the only description of them that survives a caller who draws them by hand.

    Three defects, and each is a different kind of thing:

      SAME COLOUR TWICE      the book's own words, p.115: "under no circumstances can the same
                             color be used twice in one layout". ERROR.
      NOT A GROUP 3 COLOUR   a band painted in something outside the seven-colour set. WARN,
                             because check_deck already owns off-palette-for-the-brand and this
                             is the narrower claim that it is off-palette FOR A PAIR.
      WHITE TYPE TOO PALE    the field carrying body copy below 4.5:1, or display type below
                             3.0:1, against the type sitting on it. ERROR.

    THE THIRD ONE IS NOT A DUPLICATE OF check_contrast, and the difference is the whole point.
    check_contrast applies WCAG, which relaxes to 3.0:1 for anything at 18pt or above -- and
    nearly all slide type is. So white body copy on Pacific at 20pt passes WCAG and is still
    unreadable from the back of a 140-person room. The rider says a field carrying BODY COPY
    needs 4.5:1 whatever size it is set at, and that stricter floor is what this check adds.

    This check is KPMG-specific by construction: it enforces a rule out of KPMG's brand book. A
    different brand has its own page grammar or none, so it is skipped rather than guessed at.
    """
    if brand.name.upper() != "KPMG":
        return
    from . import palette

    report.checked.append("colour pairs")

    prs = Presentation(path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    if not slide_w or not slide_h:
        return

    set_hexes = {palette._HEX[r] for r in palette.PAIR_SET}
    body_hexes = {palette._HEX[r] for r in palette.BODY_FIELDS}
    display_hexes = {palette._HEX[r] for r in palette.DISPLAY_FIELDS}
    # White and the near-white canvas are NOT pair failures. "One colour against white" is a
    # Group 1/2 page -- a legitimate and very common KPMG layout that simply is not this system,
    # and flagging it would make this check fire on most of the deck.
    exempt = {
        "#FFFFFF",
        brand.color("canvas").upper(),
        brand.color("canvas_alt").upper(),
    }

    for index, slide in enumerate(prs.slides, start=1):
        rects = _filled_rects(slide)
        bands = [
            r
            for r in rects
            if (r[3] - r[1]) >= 0.95 * slide_h and (r[2] - r[0]) >= 0.01 * slide_w
        ]
        if len(bands) < 2:
            continue
        covered = sum(min(r[2], slide_w) - max(r[0], 0) for r in bands)
        if covered < 0.95 * slide_w:
            continue

        colours = [r[4].upper() for r in bands]
        if all(c in exempt for c in colours):
            continue

        # The field is the widest band, the element the narrowest. On a page built as one
        # full-bleed field with the element laid over it, the widest is still the field.
        by_width = sorted(bands, key=lambda r: r[2] - r[0])
        element_hex, field_hex = by_width[0][4].upper(), by_width[-1][4].upper()

        if element_hex == field_hex or len(set(colours) - exempt) < 2:
            report.add(
                "ERROR",
                "colour pairs",
                index,
                f"a two-colour layout painted {element_hex} on {element_hex}. The brand book: "
                f"'under no circumstances can the same color be used twice in one layout'",
            )
            continue

        for role, hexval in (("element", element_hex), ("field", field_hex)):
            if hexval not in set_hexes and hexval not in exempt:
                report.add(
                    "WARN",
                    "colour pairs",
                    index,
                    f"{hexval} is the {role} of a two-colour layout but is not one of the "
                    f"seven Group 3 colours ({sorted(set_hexes)})",
                )

        if (
            frozenset(
                {palette._BY_HEX.get(element_hex), palette._BY_HEX.get(field_hex)}
            )
            in palette.NEAR_IDENTICAL_PAIRS
        ):
            report.add(
                "WARN",
                "colour pairs",
                index,
                f"{element_hex} on {field_hex} are two different colours by the book's rule and "
                f"{contrast_ratio(element_hex, field_hex):.2f}:1 apart in the same hue family -- "
                f"at projector distance this reads as one page with a stripe, not as a pair",
            )

        # -- the AA rider on the type-bearing field ---------------------------
        if field_hex in exempt:
            continue
        bg_default = brand.color("canvas")
        carries_body = False
        worst: tuple[float, str, str] | None = None
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            if _effective_bg(shape, rects, bg_default).upper() != field_hex:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rpr = run.font._rPr
                    srgb = (
                        rpr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
                        if rpr is not None
                        else None
                    )
                    if srgb is None or not srgb.get("val"):
                        continue
                    if srgb.find(qn("a:alpha")) is not None:
                        continue
                    fg = "#" + srgb.get("val").upper()
                    size_pt = run.font.size.pt if run.font.size else 18
                    # BODY COPY IS A JOB, NOT A SIZE, and the two only mostly agree. A run set
                    # under 28pt is body copy at this package's scale; so is a long run at any
                    # size, because 200 characters is a paragraph however large it is set. The
                    # heuristic errs toward the STRICTER threshold on purpose -- calling display
                    # type body costs a legal pair, calling body copy display costs legibility.
                    if size_pt < 28 or len(text) > 200:
                        carries_body = True
                    ratio = contrast_ratio(fg, field_hex)
                    if worst is None or ratio < worst[0]:
                        worst = (ratio, fg, run.text[:44])

        if worst is None:
            continue
        need = palette.BODY_MIN_RATIO if carries_body else palette.DISPLAY_MIN_RATIO
        allowed = body_hexes if carries_body else display_hexes
        if worst[0] < need:
            kind = "body copy" if carries_body else "display type"
            report.add(
                "ERROR",
                "colour pairs",
                index,
                f"{worst[1]} {kind} on the {field_hex} field is {worst[0]:.2f}:1, under the "
                f'{need}:1 the rider requires -- "{worst[2]}". Legal fields for {kind} are '
                f"{sorted(allowed)}; use {field_hex} as the element instead",
            )

    # CHART-ONLY COLOURS USED AS PAGE FURNITURE.
    #
    # FY22 p.48 fences eight colours as "for use in infographics and charts only"; p.53 states it
    # as a prohibition. This is invisible to every other check in this file, and that is the point:
    # such a colour is on-palette, and it can be perfectly legible, and it is still wrong. Nothing
    # derivable from contrast or palette membership would catch it.
    #
    # It is a real defect and not a hypothetical: Dark Purple #510DBC cleared 9.84:1 against white,
    # so it entered the body-field set, so block pages shipped with it as a full-width fill.
    #
    # Threshold: 4% of the slide. Below that a mark is plausibly a chart element, a legend swatch
    # or a rule, all of which are legitimate uses -- the rule is about FURNITURE, so the check has
    # to be about area rather than about presence.
    chart_only = {h.upper() for h in brand.chart_only}
    if chart_only:
        area = slide_w * slide_h
        for index, slide in enumerate(prs.slides, start=1):
            for x0, y0, x1, y1, hexv in _filled_rects(slide):
                if hexv.upper() not in chart_only:
                    continue
                share = ((x1 - x0) * (y1 - y0)) / area if area else 0
                if share >= 0.04:
                    report.add(
                        "ERROR",
                        "colour pairs",
                        index,
                        f"{hexv} fills {share:.0%} of the slide, but the brand reserves it for "
                        f"charts and infographics only -- it is not page furniture. Use a "
                        f"primary-palette colour for the field",
                    )


def check_overflow(
    path: str, brand: Brand, report: Report, *, grid: Grid | None = None
) -> None:
    """
    RE-MEASURE EVERY TEXT FRAME AGAINST ITS OWN BOX, and flag the ones that do not fit.

    This is the gate that was missing, and its absence is why a rebuild that passed every check
    in this file still had text running through other text on ten pages and off the bottom of
    the page on three. The reason it is invisible to everything else here is worth stating: AN
    OVERFLOWING TEXT FRAME IS A VALID TEXT FRAME. The XML is well-formed, the shape is on the
    slide, the colours are in the palette, the size is above the floor. PowerPoint just draws
    the extra lines outside the box, over whatever is beneath them.

    The building code already measures -- `text.fit` is what every component uses to decide
    what fits. The failure mode is a component that measures and then a CALLER that allocates a
    fixed height anyway, or one that passes `strict=False` and moves on. So this re-runs the
    same measurement over the FINISHED file, where no caller can have skipped it.

    Reported as WARN rather than ERROR: deliberate overhang is a real technique (a bleeding
    numeral, a figure that breaks its frame), and text.fit's `safety` multiplier errs toward
    declaring an overflow. A page with one 4pt overhang is fine; a page with 40pt of it is the
    defect this exists to catch, so the report says by how much.
    """
    report.checked.append("text overflow")
    prs = Presentation(path)
    grid = grid or Grid(slide_w=prs.slide_width, slide_h=prs.slide_height)

    for index, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text.strip():
                continue
            try:
                w, h = shape.width, shape.height
            except (AttributeError, TypeError):
                continue
            if None in (w, h) or w <= 0 or h <= 0:
                continue

            runs = [
                r for p in shape.text_frame.paragraphs for r in p.runs if r.text.strip()
            ]
            if not runs:
                continue
            size_pt = max((r.font.size.pt for r in runs if r.font.size), default=18)
            bold = bool(runs[0].font.bold)
            # TRACKING IS PART OF THE MEASUREMENT. Display type carries negative tracking --
            # -1.08pt a character at 60pt -- and ignoring it over-measures a headline by 15pt
            # or more, which is the difference between one line and two. The first version of
            # this check reported five headline "overflows" that were all correctly set; a gate
            # that has to be argued with is a gate people stop reading.
            spc = 0.0
            rpr = runs[0].font._rPr
            if rpr is not None and rpr.get("spc"):
                try:
                    spc = int(rpr.get("spc")) / 100.0
                except (TypeError, ValueError):
                    spc = 0.0
            family = runs[0].font.name or brand.font_minor
            spacing = shape.text_frame.paragraphs[0].line_spacing
            if not isinstance(spacing, (int, float)):
                spacing = 1.2

            from .canvas import Box as _Box

            measured = _fit_text(
                text, _Box(0, 0, w, h), family, size_pt, bold, float(spacing), spc
            )
            if measured is None or measured.fits:
                continue
            over = measured.overflow_pt
            if over < 4:
                continue
            report.add(
                "WARN" if over < 24 else "ERROR",
                "text overflow",
                index,
                f"{over:.0f}pt of text past the bottom of its own box at {size_pt:g}pt "
                f'-- "{text.strip()[:44]}". It will be drawn over whatever is beneath it.',
            )


def _fit_text(text, box, family, size_pt, bold, spacing, tracking_pt=0.0):
    from . import text as T

    try:
        return T.fit(
            text,
            box,
            family,
            size_pt,
            bold=bold,
            line_spacing=spacing,
            tracking_pt=tracking_pt,
        )
    except (ValueError, OSError, T.FontNotFound):
        # A face this machine cannot load cannot be measured, and that is itself worth knowing
        # -- but it is `check_deck`'s business, not this check's. Skip rather than crash: a
        # gate that dies on the first unmeasurable run reports nothing about the other 400.
        return None


def check_siblings(path: str, report: Report, *, max_share: float = 0.20) -> None:
    """
    Assert that NO SINGLE LAYOUT ACCOUNTS FOR MORE THAN 20% OF SLIDES.

    THIS CHECK USED TO ASSERT THE OPPOSITE AND THAT WAS THE DEFECT. It clustered slides by
    headline position and warned on any position used exactly ONCE -- treating a one-off as the
    thing to catch. Measured against real material, the ratio runs the other way: KPMG has
    21-23 distinct layouts across 22-23 pages, with a largest repeat of one or two (4-9%). Our
    deck had NINE layouts across 37 slides with one repeated FOURTEEN times, 38% of the deck.

    So the old check would have passed a deck whose every page looked the same, and warned
    about the pages that did not. The rhythm strip of a deck like that is a square wave where
    theirs is a continuous spectrum, and it is legible at contact-sheet size before a single
    word is read.

    Uniformity within an ARCHETYPE is still correct and is enforced somewhere better: the
    archetypes themselves place their own headline, so two `statement()` pages cannot disagree.
    What this catches is one archetype being used for most of the deck, which is a content
    problem wearing a layout costume.
    """
    report.checked.append("layout diversity")
    prs = Presentation(path)

    positions: dict[tuple[int, int, float], list[int]] = {}
    for index, slide in enumerate(prs.slides, start=1):
        best = None
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            sizes = [
                r.font.size.pt
                for p in shape.text_frame.paragraphs
                for r in p.runs
                if r.font.size
            ]
            if not sizes:
                continue
            top_size = max(sizes)
            if best is None or top_size > best[0]:
                best = (top_size, shape.left, shape.top)
        if best:
            positions.setdefault((best[1], best[2], best[0]), []).append(index)

    total = len(prs.slides)
    if total < 5:
        return
    for (left, top, size), slides in sorted(
        positions.items(), key=lambda kv: -len(kv[1])
    ):
        share = len(slides) / total
        if share > max_share:
            report.add(
                "WARN",
                "layout diversity",
                None,
                f"one layout (headline {size:g}pt at {left / EMU_PER_POINT:.0f}, "
                f"{top / EMU_PER_POINT:.0f}pt) carries {len(slides)} of {total} slides "
                f"({share:.0%}), over the {max_share:.0%} ceiling. KPMG's largest repeat "
                f"across a comparable deck is 1-2 pages (4-9%).",
            )


def check_rhythm(
    render_dir: str,
    report: Report,
    *,
    min_ink: float = 0.08,
    genre: str = "deck",
) -> None:
    """
    THE DENSITY GATE, and the only check in this file that needs the deck RENDERED.

    Every other check reads the .pptx. This one reads PNGs, because ink coverage is a fact
    about the drawn page and cannot be computed from the XML -- a slide can carry forty shapes
    and still be empty, which is precisely the failure it exists to catch.

    Fails any page under 8% non-white pixels. The threshold is measured, not chosen: 28 of our
    37 pages sat under it while roughly 2% of 62 real KPMG pages do, and their one instance is
    a BACK COVER. A near-empty interior page is not a minimal one; it is a page whose
    composition was decided by what was removed.

    It also reports the page-type mix against a measured reference, as a WARNING rather than an
    error, because the mix is a judgment and the empty page is not. `genre` picks the reference:
    "deck" (the default, from KPMG's own nine slides) or "report" (from the CEO-Outlook PDFs).
    Passing the wrong one produces two confident, meaningless warnings -- see REFERENCE below.

        from kpmg_deck.verify import check_rhythm, Report
        r = Report(); check_rhythm("render", r); print(r.render())
    """
    import glob

    report.checked.append("page rhythm")
    try:
        from PIL import Image
    except ImportError:
        report.add(
            "WARN",
            "page rhythm",
            None,
            "Pillow is not installed, so ink coverage could not be measured. "
            "pip install pillow -- this check is the density gate and skipping it is how "
            "28 near-empty pages shipped last time.",
        )
        return

    pages = sorted(glob.glob(posixpath.join(render_dir, "slide-*.png")))
    if not pages:
        report.add(
            "ERROR",
            "page rhythm",
            None,
            f"no slide-*.png in {render_dir!r}. Render the deck first: "
            f"scripts/render.sh deck.pptx {render_dir}",
        )
        return

    bands = {"field": 0, "split": 0, "white": 0}
    for i, page in enumerate(pages, start=1):
        img = Image.open(page).convert("L").resize((240, 135))
        px = list(img.getdata())
        non_white = sum(1 for v in px if v < 246) / len(px)
        if i > 1 and non_white < min_ink:
            report.add(
                "ERROR",
                "page rhythm",
                i,
                f"{non_white:.1%} ink is below the {min_ink:.0%} floor -- this page is "
                f"near-empty. Add the second and third element, do not add colour: a "
                f"coloured void is still a void.",
            )
        if non_white > 0.80:
            bands["field"] += 1
        elif non_white > 0.30:
            bands["split"] += 1
        else:
            bands["white"] += 1

    n = len(pages)
    # THE REFERENCE MIX IS PER GENRE, AND THE DEFAULT USED TO BE A REPORT'S.
    #
    # 13/40/46 was measured off KPMG's CEO-Outlook PDFs -- an editorial REPORT family. It was
    # applied to a DECK, and it duly fired on both bands at once: field pages "too many" at 26%
    # against 13, white pages "too few" at 33% against 46. Both readings were arithmetically
    # right and neither meant anything, because a report and a deck are not the same object.
    #
    # This is the same defect as the colour error that started this whole rebuild, one layer up:
    # a measurement of ONE GENRE presented as the system. It is worth naming twice because the
    # second instance was inside the gate built to catch the first.
    #
    # The deck reference is measured off the nine real KPMG slides printed in the FY22 brand
    # guidelines p.113 -- the only KPMG-authored slides we hold: 4 of 9 white ground, 3 of 9
    # coloured ground, 2 of 9 photographic. Ink coverage there runs 8.7% to 77.1%. Mapped onto
    # this check's three bands, with the photographic slides counting as field pages because
    # that is what they measure as: field 0.33, split 0.22, white 0.44.
    #
    # ONE HONEST CAVEAT, kept because it bounds the claim: those nine are brand-book EXEMPLARS,
    # composed to demonstrate range, so the real spread of a 40-slide engagement deck is likely
    # narrower. Treat the deck numbers as an envelope rather than a target -- which is why this
    # stays a WARNING with a 10pp band and not an error.
    REFERENCE = {
        "report": {"field": 0.13, "split": 0.40, "white": 0.46},
        "deck": {"field": 0.33, "split": 0.22, "white": 0.44},
    }
    target = REFERENCE.get(genre, REFERENCE["deck"])
    for band, want in target.items():
        got = bands[band] / n
        if abs(got - want) > 0.10:
            report.add(
                "WARN",
                "page rhythm",
                None,
                f"{band} pages are {got:.0%} of the deck against a measured {want:.0%} "
                f"(+/-10pp). {bands[band]} of {n}.",
            )


# ---------------------------------------------------------------------------
# Monotony -- the gate that reads the PICTURE rather than the authorship
# ---------------------------------------------------------------------------

# GRID CONSTANTS, AND WHY THEY ARE THIS COARSE.
#
# 32x32 grayscale for the duplicate test and 12x7 cells for the shape test. Both are far below
# the resolution at which text is legible, deliberately: this is a contact-sheet judgement, not
# image forensics. At 32x32 a headline is four grey smudges, which is exactly the level at which
# "these two are the same page" is decided by an eye scanning a printed contact sheet -- and it
# keeps 903 pairwise comparisons on a 43-page deck under a tenth of a second in pure Python.
_DUP_GRID = 32
_SHAPE_COLS, _SHAPE_ROWS, _SHAPE_CELL = 12, 7, 8
_SHAPE_CELLS = _SHAPE_COLS * _SHAPE_ROWS


def _ink_signature(px: bytes, *, delta: int = 24) -> tuple[int, ...]:
    """
    Reduce one rendered page to 84 bits: which cells of a 12x7 grid carry marks.

    INK IS MEASURED AGAINST THE PAGE'S OWN GROUND, NOT AGAINST WHITE, and that one decision is
    what makes the signature mean "shape" rather than "colour". The first version counted
    non-white pixels the way `check_rhythm` does. On a full-bleed KPMG Blue page every cell is
    non-white, so every coloured page collapsed to the same all-ones signature and the cover,
    six section dividers and five dark statement pages came back as ONE shape carrying 26% of the
    deck -- a reading that is arithmetically true and says nothing, because those pages do not
    look alike.

    So the background is the page's own modal luminance (16-level buckets), and a cell is
    "inked" when at least half its samples sit more than `delta` away from it. A white page and
    its blue-ground twin then produce the SAME signature when the type sits in the same places,
    which is the correct answer: they are one composition in two colourways.
    """
    hist = Counter(v // 16 for v in px)
    bg = hist.most_common(1)[0][0] * 16 + 8
    width = _SHAPE_COLS * _SHAPE_CELL
    half = (_SHAPE_CELL * _SHAPE_CELL) // 2
    sig: list[int] = []
    for row in range(_SHAPE_ROWS):
        for col in range(_SHAPE_COLS):
            inked = 0
            for y in range(_SHAPE_CELL):
                start = (row * _SHAPE_CELL + y) * width + col * _SHAPE_CELL
                for v in px[start : start + _SHAPE_CELL]:
                    if abs(v - bg) > delta:
                        inked += 1
            sig.append(1 if inked >= half else 0)
    return tuple(sig)


def check_monotony(
    render_dir: str,
    report: Report,
    *,
    max_rms: float = 6.0,
    max_shape_share: float = 0.25,
    shape_tolerance: int = 4,
) -> None:
    """
    THE GATE THAT MEASURES WHAT THE READER SEES, NOT WHAT THE AUTHOR CALLED IT.

    `check_siblings` counts ARCHETYPE NAMES and is structurally blind to this: `statement`,
    `split` and `points_slide` are three different names that render as the same picture, so a
    deck can spread itself evenly across the archetype list and still be one page over and over.
    Measured on this package's own 43-slide reference deck, which `check_siblings` passed clean:
    slides 9, 16, 21, 26, 34 and 42 are the SAME PAGE six times, differing by one word ("Open the
    section N visualization and play it"). Nothing in this file could see it. SKILL.md's own
    line 194 says the failure mode this package fears is monotony rather than variety; until
    this check existed, nothing measured monotony.

    It reads the rendered PNGs, like `check_rhythm` and for the same reason -- what a page LOOKS
    like is a fact about the drawn page, and two slides built from different code can draw the
    same picture.

    (a) NEAR-DUPLICATES ARE AN ERROR. Two pages that render alike is a defect with no defensible
        reading, so unlike the rest of the render-reading family this one fails the build.

    (b) SHAPE CONCENTRATION IS A WARNING, with the offending slide numbers named. One layout
        carrying a quarter of the deck is a judgement call -- a workshop deck may legitimately
        run a worked-example page nine times -- so it reports and does not fail.
    """
    import glob

    report.checked.append("page monotony")
    try:
        from PIL import Image
    except ImportError:
        report.add(
            "WARN",
            "page monotony",
            None,
            "Pillow is not installed, so the rendered pages could not be compared. "
            "pip install pillow -- this is the check that catches six copies of one page, "
            "which every structural gate in this file passes clean.",
        )
        return

    pages = sorted(glob.glob(posixpath.join(render_dir, "slide-*.png")))
    if not pages:
        report.add(
            "ERROR",
            "page monotony",
            None,
            f"no slide-*.png in {render_dir!r}. Render the deck first: "
            f"scripts/render.sh deck.pptx {render_dir}",
        )
        return

    coarse: list[bytes] = []
    shapes: list[tuple[int, ...]] = []
    for page in pages:
        img = Image.open(page).convert("L")
        coarse.append(img.resize((_DUP_GRID, _DUP_GRID), Image.BILINEAR).tobytes())
        small = img.resize(
            (_SHAPE_COLS * _SHAPE_CELL, _SHAPE_ROWS * _SHAPE_CELL), Image.BILINEAR
        )
        shapes.append(_ink_signature(small.tobytes()))

    total = len(pages)

    # (a) NEAR-DUPLICATES.
    #
    # THE THRESHOLD IS MEASURED, NOT CHOSEN, and the separation is wide enough to state as a
    # fact. Root-mean-square grayscale difference at 32x32, on this deck's 903 pairs:
    #
    #     worst pair among the six known duplicates ... 4.39   (slides 26 and 42)
    #     best pair among the six known duplicates .... 2.13   (slides 16 and 34)
    #     closest genuinely DISTINCT pair ............. 9.01   (slides 5 and 29 -- same layout,
    #                                                           different headline and body)
    #
    # 6.0 sits in the middle of a gap with nothing in it, so the check does not depend on the
    # third decimal place of anything. RMS was chosen over a perceptual hash because a dHash
    # answers "is the gradient pattern alike" in a fixed number of bits and gives no distance to
    # report; RMS is in luminance units, so the finding can quote how close the pages actually
    # are and the reader can check it.
    #
    # The mean-difference pre-filter below is not an approximation: |mean(a) - mean(b)| is a
    # lower bound on RMS(a, b), so skipping on it can never miss a pair. It removes essentially
    # every white-vs-blue comparison before the inner loop runs.
    means = [sum(c) / len(c) for c in coarse]
    parent = list(range(total))

    def _root(i: int) -> int:
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    worst: dict[int, float] = {}
    for i in range(total):
        for j in range(i + 1, total):
            if abs(means[i] - means[j]) > max_rms:
                continue
            a, b = coarse[i], coarse[j]
            rms = (sum((x - y) ** 2 for x, y in zip(a, b)) / len(a)) ** 0.5
            if rms > max_rms:
                continue
            ri, rj = _root(i), _root(j)
            if ri != rj:
                parent[max(ri, rj)] = min(ri, rj)
            key = min(_root(i), _root(j))
            worst[key] = max(worst.get(key, 0.0), rms)

    clusters: dict[int, list[int]] = {}
    for i in range(total):
        clusters.setdefault(_root(i), []).append(i + 1)
    for key, members in sorted(clusters.items()):
        if len(members) < 2:
            continue
        listed = ", ".join(str(m) for m in members)
        report.add(
            "ERROR",
            "page monotony",
            None,
            f"slides {listed} render as the same page (grayscale RMS {worst[key]:.1f} at "
            f"{_DUP_GRID}x{_DUP_GRID}; the closest genuinely distinct pair in a comparable "
            f"deck measures 9.0). {len(members)} copies of one composition is not a rhythm, "
            f"it is a stutter -- give each its own page or collapse them into one.",
        )

    # (b) SHAPE CONCENTRATION.
    #
    # Greedy single-pass clustering against each cluster's FIRST member, never against a moving
    # centroid: chaining is what turns a tolerance into a bucket that swallows the deck. The
    # tolerance is 4 differing cells out of 84 -- about 5% of the grid, one small element moved
    # or dropped. Measured on the reference deck it produces 17 shapes across 43 pages, the
    # largest being 11 white pages that carry the identical skeleton (eyebrow, two-tone headline
    # top-left, two or three columns of 9pt body, blue closing line): slides 4, 7, 11, 12, 18,
    # 19, 23, 28, 31, 33, 36. That is 26%, and it is the shape the deck is actually made of.
    reps: list[tuple[int, ...]] = []
    groups: list[list[int]] = []
    for index, sig in enumerate(shapes, start=1):
        for k, rep in enumerate(reps):
            if sum(a != b for a, b in zip(rep, sig)) <= shape_tolerance:
                groups[k].append(index)
                break
        else:
            reps.append(sig)
            groups.append([index])

    if total < 5:
        return

    # A KNOWN LIMIT OF THIS ARM, RECORDED RATHER THAN TUNED AWAY. The signature is ink position,
    # and roughly a third of every page's ink is CHROME the deck holds constant on purpose --
    # nav strip, eyebrow, two-tone headline at the top left, margin note, footer. So ten pages
    # carrying ten genuinely different exhibits (a timeline, two tables, two versus, two metric
    # rows, a stacked bar) can still land in one bucket, because what they share is the house
    # skeleton and that skeleton is a feature.
    #
    # Observed on the proof deck at 10 of 35: inspected page by page, and the near-duplicate it
    # was really pointing at was ONE pair -- two pages built from the same giant "1x" figure with
    # three leaders. Fixing that pair was right; converting the other eight onto other devices to
    # get under the ceiling would have been variety for its own sake, which is the failure this
    # file warns about in the opposite direction.
    #
    # The threshold is NOT relaxed to accommodate that, because a ceiling moved to make the
    # current deck pass measures nothing. Read a shape-share warning as "go and look at these
    # pages", not as "this deck is wrong" -- the duplicate ERROR arm above is the one that
    # decides, and it is measured against a gap with 2.05 of clearance either side.
    for members in sorted(groups, key=len, reverse=True):
        share = len(members) / total
        if share <= max_shape_share:
            break
        listed = ", ".join(str(m) for m in members)
        report.add(
            "WARN",
            "page monotony",
            None,
            f"one visual shape carries {len(members)} of {total} pages ({share:.0%}), over "
            f"the {max_shape_share:.0%} ceiling: slides {listed}. These differ in wording, not "
            f"in composition -- at contact-sheet size they are one page. Re-cut some of them "
            f"onto a different device (the window, a full-bleed field, an exhibit).",
        )


def check_device_mix(path: str, report: Report) -> None:
    """
    Report which BRAND DEVICES the deck actually uses, against FY22's own nine slides.

    WARNING-ONLY, AND DELIBERATELY LOOSE, because the reference is an envelope and not a target.
    The nine PowerPoint slides printed on FY22 p.113 measure: white ground 4/9, coloured ground
    3/9, photographic 2/9, carries the window 5/9, carries a photograph anywhere 6/9. Those nine
    were composed to demonstrate the system's RANGE. A real 40-slide engagement deck is
    legitimately narrower, and BUILD-SPEC-FY22.md section 5 says so in as many words. A gate that
    reported a real deck as wrong for not matching a demo reel would be worse than no gate, so
    this one fires only at a 25pp gap -- wide enough that it means a device is MISSING rather
    than merely rarer, since the reference itself is n=9 and one slide there is worth 11pp.

    DETECTION IS GEOMETRIC, NOT BY API NAME, so it keeps working when an archetype is renamed:

      ground        the slide's own background fill, plus any filled non-picture shape covering
                    60%+ of the page; a picture at 60%+ makes it photographic.
      window        FY22 p.72 states the proportion -- "7:10 for vertical and 10:7 for
                    horizontal" -- so a filled, non-text rectangle within 10% of either ratio and
                    covering 3-60% of the page is the window -- whether it is FILLED (style 2,
                    holding type) or a PICTURE (style 3, holding an image). Both are windows and
                    counting only the first inverted the gate's advice; see the note at the
                    picture branch. It does NOT fire on the 0.65-ratio body columns, which carry no
                    fill.
      photograph    picture parts, EXCLUDING THE LOGO. The logo is excluded two ways: anything
                    under 2% of the page (ours measures 0.6%), and any image whose sha1 recurs on
                    more than a quarter of the slides, because a mark repeated on most pages is
                    furniture rather than photography. Untested against a deck that actually
                    carries photographs -- the reference deck has zero picture parts of any kind,
                    which is itself the loudest thing this check reports.
    """
    report.checked.append("device mix")
    prs = Presentation(path)
    total = len(prs.slides)
    if not total:
        return
    area = prs.slide_width * prs.slide_height

    def _filled(shape) -> bool:
        try:
            return shape.fill.type in (1, 3)  # SOLID, GRADIENT
        except (AttributeError, TypeError, ValueError, NotImplementedError):
            return False

    # First pass: how often each distinct image blob appears, so the logo can be told from a
    # photograph without knowing the logo's filename.
    recurrence: Counter[str] = Counter()
    for slide in prs.slides:
        seen: set[str] = set()
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    seen.add(shape.image.sha1)
                except (AttributeError, KeyError, ValueError):
                    continue
        recurrence.update(seen)

    grounds = {"white": 0, "coloured": 0, "photographic": 0}
    with_window = 0
    with_photo = 0

    for slide in prs.slides:
        ground = "white"
        try:
            bg = slide.background.fill
            if bg.type == 1 and str(bg.fore_color.rgb).upper() != "FFFFFF":
                ground = "coloured"
        except (AttributeError, TypeError, ValueError, NotImplementedError):
            pass

        window = False
        photo = False
        for shape in _iter_shapes(slide.shapes):
            w, h = shape.width, shape.height
            if not w or not h or w <= 0 or h <= 0:
                continue
            cover = (w * h) / area
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    sha = shape.image.sha1
                except (AttributeError, KeyError, ValueError):
                    sha = ""
                if cover >= 0.02 and recurrence[sha] <= total / 4:
                    photo = True
                    if cover >= 0.60:
                        ground = "photographic"
                    # A WINDOW HOLDING AN IMAGE IS STILL A WINDOW. This branch used to `continue`
                    # straight past the ratio test, so window STYLE 3 -- the image-holding shape,
                    # which p.74 makes the characteristic use and style 2 the least -- counted as
                    # a photograph and never as a window. The effect was backwards: a deck that
                    # moved its dividers from style 2 to style 3, which is the correction this
                    # whole gate exists to encourage, saw its window share FALL from 16% to 2%
                    # and got told to add windows it had just added. Measured on the course deck:
                    # seven style-3 windows at ratio 1.40 and 21% coverage, all invisible here.
                    elif 0.03 <= cover < 0.60:
                        ratio = w / h
                        if abs(ratio - 10 / 7) <= 0.143 or abs(ratio - 7 / 10) <= 0.07:
                            window = True
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or not _filled(shape):
                continue
            if cover >= 0.60 and ground != "photographic":
                ground = "coloured"
            if 0.03 <= cover < 0.60:
                ratio = w / h
                if abs(ratio - 10 / 7) <= 0.143 or abs(ratio - 7 / 10) <= 0.07:
                    window = True

        grounds[ground] += 1
        with_window += 1 if window else 0
        with_photo += 1 if photo else 0

    measured = {
        "white ground": grounds["white"] / total,
        "coloured ground": grounds["coloured"] / total,
        "photographic ground": grounds["photographic"] / total,
        "carries the window": with_window / total,
        "carries a photograph": with_photo / total,
    }
    counts = {
        "white ground": grounds["white"],
        "coloured ground": grounds["coloured"],
        "photographic ground": grounds["photographic"],
        "carries the window": with_window,
        "carries a photograph": with_photo,
    }
    REFERENCE = {
        "white ground": 4 / 9,
        "coloured ground": 3 / 9,
        "photographic ground": 2 / 9,
        "carries the window": 5 / 9,
        "carries a photograph": 6 / 9,
    }

    gaps = [k for k, v in measured.items() if abs(v - REFERENCE[k]) > 0.25]
    if not gaps:
        return

    got = "; ".join(f"{k} {measured[k]:.0%} ({counts[k]}/{total})" for k in REFERENCE)
    want = "; ".join(f"{k} {v:.0%}" for k, v in REFERENCE.items())
    report.add(
        "WARN",
        "device mix",
        None,
        f"{got}. Reference (FY22 p.113, the nine KPMG-authored PowerPoint slides): {want}. "
        f"Outside a 25pp band on: {', '.join(gaps)}. TREAT AS AN ENVELOPE, NOT A TARGET -- "
        f"those nine were composed to demonstrate the system's range, so a real engagement "
        f"deck is legitimately narrower and this is never an error. What it is worth reading "
        f"for is a device at zero: the deck is not using part of the system at all.",
    )


def check_notes(path: str, report: Report, *, skip: Iterable[int] = ()) -> None:
    """
    Warn on content slides with no speaker notes.

    Notes are free and slides are not. A content slide with none usually means its argument is
    still crammed onto the slide, because there was nowhere else to put it.
    """
    report.checked.append("speaker notes")
    prs = Presentation(path)
    skip = set(skip)
    for index, slide in enumerate(prs.slides, start=1):
        if index in skip:
            continue
        if (
            not slide.has_notes_slide
            or not slide.notes_slide.notes_text_frame.text.strip()
        ):
            report.add("WARN", "speaker notes", index, "no speaker notes")


def verify(
    path: str,
    brand: Brand,
    *,
    notes_skip: Iterable[int] = (),
    check_notes_too: bool = True,
    render_dir: str | None = None,
) -> Report:
    """
    Run every mechanical gate. Returns a Report; `.ok` is False if anything is an ERROR.

    `render_dir` additionally runs the two gates that READ THE PICTURE -- density
    (`check_rhythm`) and monotony (`check_monotony`). Pass it. A deck can pass every structural
    check in this file and still be 28 near-empty pages, which is exactly what happened -- the
    XML was valid on every one of them -- and it can equally be six copies of the same page,
    which is what happened next: `check_siblings` counts archetype NAMES and cannot see that
    three different names draw one picture.
    """
    report = Report()
    check_package(path, report)
    check_deck(path, brand, report)
    check_contrast(path, brand, report)
    check_pairs(path, brand, report)
    check_overflow(path, brand, report)
    check_wrap_and_autofit(path, report)
    check_siblings(path, report)
    check_device_mix(path, report)
    if render_dir:
        check_rhythm(render_dir, report)
        check_monotony(render_dir, report)
    if check_notes_too:
        check_notes(path, report, skip=notes_skip)
    return report


if __name__ == "__main__":
    import os
    import sys

    if len(sys.argv) < 2:
        print("usage: python -m kpmg_deck.verify <deck.pptx> [brand] [render-dir]")
        raise SystemExit(2)
    deck_path = sys.argv[1]
    brand_name = sys.argv[2] if len(sys.argv) > 2 else "kpmg"
    # Pick up a render beside the deck without being told, so the density gate is the default
    # rather than a thing to remember. It is the one check that catches the defect class this
    # package most recently shipped.
    default_render = os.path.join(os.path.dirname(os.path.abspath(deck_path)), "render")
    renders = (
        sys.argv[3]
        if len(sys.argv) > 3
        else (default_render if os.path.isdir(default_render) else None)
    )
    result = verify(deck_path, Brand.load(brand_name), render_dir=renders)
    print(result.render())
    if not renders:
        print(
            "\n  NOTE: no render directory found, so the DENSITY and MONOTONY GATES did not "
            "run. Render the deck and re-run:\n"
            "    scripts/render.sh <deck.pptx> render && python -m kpmg_deck.verify "
            "<deck.pptx> kpmg render"
        )
    raise SystemExit(0 if result.ok else 1)
